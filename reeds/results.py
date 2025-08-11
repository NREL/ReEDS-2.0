# Imports
import os
import sys
import io
import numpy as np
from glob import glob
import pandas as pd
import matplotlib.pyplot as plt
from itertools import product
import pptx
from pptx.util import Inches, Pt

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
import reeds

reeds_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if 'runs' in reeds_path.split(os.path.sep):
    reeds_path = reeds_path[: reeds_path.index(os.sep + 'runs' + os.sep)]

### Globals
DEFAULT_DOLLAR_YEAR = 2024
DEFAULT_PV_YEAR = 2024
### Source for default discount rate:
### https://www.whitehouse.gov/wp-content/uploads/2023/11/CircularA-4.pdf
DEFAULT_DISCOUNT_RATE = 0.02
DEFAULT_END_YEAR = 2050
SLIDE_HEIGHT = 6.88
SLIDE_WIDTH = 13.33

#%% ===========================================================================
### --- Powerpoint Tools ---
### ===========================================================================
def init_pptx(
    fpath_template=os.path.join(reeds_path, 'postprocessing', 'template.pptx'),
):
    prs = pptx.Presentation(fpath_template)
    return prs

def add_to_pptx(
        title=None,
        prs=None,
        file=None,
        left=0,
        top=0.62,
        width=13.33,
        height=None,
        verbose=1,
        slide=None,
        layout=3,
    ):
    """Add current matplotlib figure (or file if specified) to new powerpoint slide"""
    if not file:
        image = io.BytesIO()
        plt.savefig(image, format='png')
    else:
        image = file
        if not os.path.exists(image):
            raise FileNotFoundError(image)

    if prs is None:
        prs = init_pptx()

    if slide is None:
        slide = prs.slides.add_slide(prs.slide_layouts[layout])
        slide.shapes.title.text = title
    slide.shapes.add_picture(
        image,
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    if verbose:
        print(title)
    return slide

def add_textbox(
        text,
        slide,
        left=0,
        top=7.2,
        width=13.33,
        height=0.3,
        fontsize=14,
    ):
    """Add a textbox to the specified slide"""
    textbox = slide.shapes.add_textbox(
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    p = textbox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    return slide


#%% ===========================================================================
### --- Tools to get cases
### ===========================================================================
def parse_caselist(caselist, casenames, basecase_in, titleshorten=0):
    use_table_casenames = False
    use_table_colors = False
    use_table_bases = False
    _caselist = caselist.copy()
    _casenames = casenames
    if len(_caselist) == 1:
        ## If it's a .csv, read the cases to compare
        if _caselist[0].endswith('.csv'):
            dfcase = pd.read_csv(_caselist[0], header=None, comment='#', quoting=3)
            ## First check it's a simple csv with one case per row
            if dfcase.shape[1] == 1:
                caselist = dfcase[0].tolist()
            ## Then check if it's a csv with [casepath,casename] in the header
            elif (
                ('casepath' in dfcase.loc[0].tolist())
                and ('casename' in dfcase.loc[0].tolist())
            ):
                dfcase = dfcase.T.set_index(0).T
                ## Drop cases that haven't finished yet
                unfinished = dfcase.loc[
                    ~dfcase.casepath.map(
                        lambda x: os.path.isfile(os.path.join(x,'outputs','outputs.h5')))
                ].index
                if len(unfinished):
                    print('The following cases have not yet finished:')
                    print('\n'.join(dfcase.loc[unfinished].casepath.tolist()))
                dfcase = dfcase.drop(unfinished).copy()
                caselist = dfcase.casepath.tolist()
                use_table_casenames = True
                if 'color' in dfcase:
                    if not dfcase.color.isnull().any():
                        use_table_colors = True
                if 'base' in dfcase:
                    if not dfcase.base.isnull().any():
                        use_table_bases = True
            ## Otherwise assume it's a copy of a cases_{batchname}.csv file in a case folder
            ## This approach is less robust; the others are preferred.
            else:
                prefix_plus_tail = os.path.dirname(_caselist[0])
                tails = [i for i in dfcase.iloc[0] if i not in ['Default Value',np.nan]]
                prefix = prefix_plus_tail[:-len([i for i in tails if prefix_plus_tail.endswith(i)][0])]
                caselist = [prefix+i for i in tails]
        ## Otherwise look for all runs starting with the provided string
        else:
            caselist = sorted(glob(_caselist[0]+'*'))
            ## If no titleshorten is provided, use the provided prefix
            if not titleshorten:
                titleshorten = len(os.path.basename(_caselist))
    else:
        caselist = _caselist

    ## Remove cases that haven't finished yet
    caselist = [
        i for i in caselist
        if os.path.isfile(os.path.join(i,'outputs','outputs.h5'))
    ]

    ## Get the casenames
    if use_table_casenames:
        casenames = [c.replace('\\n','\n') for c in dfcase.casename.tolist()]
    else:
        casenames = (
            _casenames.split(',') if len(_casenames)
            else [os.path.basename(c)[titleshorten:] for c in caselist]
        )

    if len(caselist) != len(casenames):
        err = (
            f"len(caselist) = {len(caselist)} but len(casenames) = {len(casenames)}\n\n"
            'caselist:\n' + '\n'.join(caselist) + '\n\n'
            'casenames:\n' + '\n'.join(casenames) + '\n'
        )
        raise ValueError(err)

    cases = dict(zip(casenames, caselist))

    # check to ensure there are at least two cases
    if len(cases) <= 1: 
        err = f"There are less than two cases being compared: {', '.join(cases.values())}"
        raise ValueError(err)

    ### Get the base cases
    if not len(basecase_in):
        basecase = list(cases.keys())[0]
    else:
        basepath = [c for c in cases.values() if c.endswith(basecase_in)]
        if len(basepath) == 0:
            err = (
                f"Use a basecase that matches one case.\nbasecase={basecase_in} matches none of:\n"
                + '\n'.join(basepath)
            )
            raise ValueError(err)
        elif len(basepath) > 1:
            err = (
                f"Use a basecase that only matches one case.\nbasecase={basecase_in} matches:\n"
                + '\n'.join(basepath)
            )
            raise ValueError(err)
        else:
            basepath = basepath[0]
            ## basecase is the short name; basepath is the full path
            basecase = casenames[caselist.index(basepath)]
            ## Put it first in the list
            cases = {**{basecase:cases[basecase]}, **{k:v for k,v in cases.items() if k != basecase}}

    ## Make case->base dictionary
    if use_table_bases:
        basemap = dfcase.set_index('casename').base.to_dict()
    else:
        basemap = dict(zip(cases, [basecase]*len(cases)))

    ## Get the colors
    if use_table_colors:
        colors = dict(zip(dfcase.casename, dfcase.color))
        for k, v in colors.items():
            if v.startswith('plt.cm.') or v.startswith('cmocean.cm.'):
                colors[k] = eval(v)
    else:
        colors = reeds.plots.rainbowmapper(cases)

    ## Take a look
    print('Analyzing the following cases:')
    for case, path in cases.items():
        print(
            f'{path} -> {case}'
            + (' (base)' if ((not use_table_bases) and (case == basecase)) else '')
        )

    return cases, colors, basecase, basemap

    
#%% ===========================================================================
### --- Postprocessing tools --- Based on bokeh
### ===========================================================================
def inflate_series(dfin):
    df_deflator = pd.read_csv(
        os.path.join(reeds_path, 'inputs', 'financials', 'deflator.csv'),
        index_col=0,
    )
    return dfin * 1 / df_deflator.loc[DEFAULT_DOLLAR_YEAR, 'Deflator']


def gather_cost_types(df):
    # Gather lists of capital and operation labels
    cost_cats_df = df['cost_cat'].unique().tolist()
    cost_cat_map = pd.read_csv(
        os.path.join(
            reeds_path, 'postprocessing', 'bokehpivot', 'in', 'reeds2', 'cost_cat_map.csv'
        ),
    )
    if (~cost_cat_map.cost_type.isin(['Capital','Operation'])).any():
        error = (
            'Invalid values in cost_type column of cost_cat_map.csv. '
            'Only valid values are "Capital" and "Operation".'
        )
        raise KeyError(error)

    # Make sure all cost categories in df are in cost_cat_map and throw error if not
    if not set(cost_cats_df).issubset(cost_cat_map['raw'].values.tolist()):
        error = (
            'Not all cost categories have been mapped. Categories without a mapping are:\n'
            '\n'.join([c for c in cost_cats_df if c not in cost_cat_map['raw'].values.tolist()])
        )
        raise KeyError(error)

    cap_type_ls = [
        c for c in cost_cats_df
        if c in cost_cat_map[cost_cat_map['cost_type']=='Capital']['raw'].tolist()]
    op_type_ls = [
        c for c in cost_cats_df
        if c in cost_cat_map[cost_cat_map['cost_type']=='Operation']['raw'].tolist()]

    return cost_cat_map, cap_type_ls, op_type_ls

def system_generation(
    case: str,
    r_subset: str = "all",
) -> pd.DataFrame:
    """
    Calculate system costs from ReEDS output.
    This function is based on pre_systemcost from bokehpivot/reeds2.py

    Args:
        case (str): Path to the ReEDS run case or outputs.h5 file.
        r_subset (str): If performing an analysis of a specific region within a run, this can be set to 
        a specific r values for isolating the costs. If used, the revenues and costs of the transmission
        will be included. 
    """
    # Find the inputs_case directory
    if case.endswith('.h5'):
        # If we are pointing to outputs.h5
        inputs_case = os.path.abspath(os.path.join(case, '..', '..', 'inputs_case'))
    else:
        inputs_case = os.path.join(case, 'inputs_case')
    # Get case inputs
    systemgen = reeds.io.read_output(case, 'gen_ivrt')
    if r_subset != 'all':
        systemgen = systemgen.loc[systemgen.r.isin(r_subset)]
    return systemgen

def system_capacity(
        case: str,
        r_subset: str = "all",
    ) -> pd.DataFrame:
        """
        Calculate system costs from ReEDS output.
        This function is based on pre_systemcost from bokehpivot/reeds2.py

        Args:
            case (str): Path to the ReEDS run case or outputs.h5 file.
            r_subset (str): If performing an analysis of a specific region within a run, this can be set to 
            a specific r values for isolating the costs. If used, the revenues and costs of the transmission
            will be included. 
        """
        # Get case output
        if r_subset != 'all':
            systemcap = reeds.io.read_output(case, 'cap_ivrt',r_filter = r_subset)
        else:
            systemcap = reeds.io.read_output(case, 'cap_ivrt')
        return systemcap

def calc_systemcost(
    case: str,
    cost_type: str = "annualized",
    group_r: bool = False,
    rename_as_bokeh: bool = False,
    through_2050: bool = False,
    r_subset: str = None,
    RECS: bool = False,
    st_subset: str = None,
    # The following parameters are only used if cost_type == 'annualized'
    discount_rate: float = DEFAULT_DISCOUNT_RATE,
    present_value_year: int = DEFAULT_PV_YEAR,
    shift_capital: bool = True,
    remove_existing: bool = False,
    crf_from_user: bool = False,
) -> pd.DataFrame:
    """
    Calculate system costs from ReEDS output.
    This function is based on pre_systemcost from bokehpivot/reeds2.py

    Args:
        case (str): Path to the ReEDS run case or outputs.h5 file.
        cost_type (str): Type of cost to calculate. Options are:
            - 'objective': Calculate objective function system costs.
            - 'annualized': Calculate annualized system costs.
        group_r (bool): If True, group results across regions (dropping the region column).
        rename_as_bokeh (bool): If True, rename cost categories to match Bokeh output.
        through_2050 (bool): If True, only include costs from present_value_year to 2050.
        discount_rate (float): Discount rate for present value calculations.
            (Only used if cost_type == 'annualized'.)
        present_value_year (int): Base year for present value calculations.
            (Only used if cost_type == 'annualized'.)
        shift_capital (bool): 
            - If True: Start capital payments in the year of the investment,
              even though loan payments typically start in the following year, so that
              investments made in 2050 are still reflected in 2050 capital payments.
            - If False: Reflects typical loan payments that start in the year after the loan.
            (Only used if cost_type == 'annualized'.)
        remove_existing (bool): If False, include historical capacity costs from before the start year (typically 2010).
            (Only used if cost_type == 'annualized'.)
        crf_from_user (bool): Use a user-specified Capital Recovery Factor (CRF) instead of the model’s.
            (Only used if cost_type == 'annualized'.)
        r_subset (str): If performing an analysis of a specific region within a run, this can be set to 
        a specific r values for isolating the costs. If used, the revenues and costs of the transmission
        will be included. 
    """
    # Find the inputs_case directory
    if case.endswith('.h5'):
        # If we are pointing to outputs.h5
        inputs_case = os.path.abspath(os.path.join(case, '..', '..', 'inputs_case'))
    else:
        inputs_case = os.path.join(case, 'inputs_case')

    # Get case inputs
    systemcost = reeds.io.read_output(case, 'systemcost_ba')
    pvf_capital = reeds.io.read_output(case, 'pvf_capital', valname='pvfcap')
    pvf_onm = reeds.io.read_output(case, 'pvf_onm', valname='pvfonm')
    crf_in = pd.read_csv(os.path.join(inputs_case, 'crf.csv'))
    df_capex_init = pd.read_csv(os.path.join(inputs_case, 'df_capex_init.csv'))
    expenditure_flow = reeds.io.read_output(case, 'expenditure_flow')
    expenditure_flow_rps = reeds.io.read_output(case, 'expenditure_flow_rps')
    expenditure_flow_int = reeds.io.read_output(case, 'expenditure_flow_int')
    if st_subset is not None:
        RECS = True # Default assumption

    sw = pd.read_csv(
        os.path.join(inputs_case, 'switches.csv'),
        header=None, index_col=0,
    ).squeeze(1)

    scalars = pd.read_csv(
        os.path.join(inputs_case, 'scalars.csv'),
        header=None, usecols=[0,1], index_col=0,
    ).squeeze(1)

    # Valid regions
    val_r = pd.read_csv(
        os.path.join(inputs_case, 'val_r.csv'),
        header=None,
    ).squeeze(1).values

    # Rename columns to match the expected format
    systemcost.rename(columns={'t':'year', 'sys_costs':'cost_cat'}, inplace=True)
    pvf_capital.rename(columns={'t':'year'}, inplace=True)
    pvf_onm.rename(columns={'t':'year'}, inplace=True)
    crf_in.rename(columns={'*t':'year'}, inplace=True)
    df_capex_init.rename(columns={'t':'year','region':'r'}, inplace=True)

    # Redefine regions for specific subset if defined
    if r_subset is not None:
        systemcost = systemcost.loc[systemcost.r.isin(r_subset)]
        df_capex_init = df_capex_init.loc[df_capex_init.r.isin(r_subset)]
        
        expenditure_flow = expenditure_flow.loc[
            expenditure_flow.r.isin(r_subset) | expenditure_flow.rr.isin(r_subset)]
        expenditure_flow = expenditure_flow.loc[~(
            expenditure_flow.r.isin(r_subset) * expenditure_flow.rr.isin(r_subset))]
        expenditure_flow.loc[expenditure_flow.r.isin(r_subset),'Value'] *= -1 # aligning signs for import/export costs & revenues
        expenditure_flow = pd.concat((
            expenditure_flow.loc[expenditure_flow.r.isin(r_subset)],
            expenditure_flow.loc[expenditure_flow.rr.isin(r_subset)].rename(
                columns = {'r':'rr','rr':'r'}
            ))).rename(columns = {'*':'cost_cat','t':'year'})
        expenditure_flow = pd.pivot_table(data = expenditure_flow, index = ['cost_cat','r','year'],
            values = ['Value'], aggfunc = 'sum').reset_index(drop = False)
        systemcost = pd.concat((systemcost,expenditure_flow))

        expenditure_flow_int = expenditure_flow_int.loc[
            expenditure_flow_int.r.isin(r_subset)]
        expenditure_flow_int['Value'] *= -1 # aligning signs for import/export costs & revenues
        expenditure_flow_int.rename(columns = {'t':'year'}, inplace = True)
        expenditure_flow_int['cost_cat'] = 'expenditure_int'
        systemcost = pd.concat((systemcost,expenditure_flow_int))

        if RECS:
            expenditure_flow_rps = expenditure_flow_rps.loc[
                expenditure_flow_rps.st.isin(st_subset) | expenditure_flow_rps.ast.isin(st_subset)]
            expenditure_flow_rps = expenditure_flow_rps.loc[~(
                expenditure_flow_rps.st.isin(st_subset) * expenditure_flow_rps.ast.isin(st_subset))]
            expenditure_flow_rps.loc[expenditure_flow_rps.st.isin(st_subset),'Value'] *= -1 # aligning signs for import/export costs & revenues
            expenditure_flow_rps = pd.concat((
                expenditure_flow_rps.loc[expenditure_flow_rps.st.isin(st_subset)],
                expenditure_flow_rps.loc[expenditure_flow_rps.ast.isin(st_subset)].rename(
                    columns = {'st':'ast','ast':'st'}
                ))).rename(columns = {'t':'year'})
            expenditure_flow_rps = pd.pivot_table(data = expenditure_flow_rps, index = ['st','year'],
                values = ['Value'], aggfunc = 'sum').reset_index(drop = False)
            expenditure_flow_int['cost_cat'] = 'expenditure_rps'
            systemcost = pd.concat((systemcost,expenditure_flow_rps))


    # Convert to Billion dollars
    systemcost['Value'] = systemcost['Value'].astype("float64")* 1e-9
    df_capex_init['capex'] = df_capex_init['capex'].astype("float64")* 1e-9

    # Years Optimized
    sim_years = sorted(systemcost['year'].unique())
    sys_eval_years = int(sw['sys_eval_years'])
    trans_crp = int(scalars['trans_crp'])
    addyears = max(sys_eval_years, trans_crp)
    firstmodelyear, lastmodelyear = int(systemcost.year.min()), int(systemcost.year.max())

    # Begin processing
    df = systemcost.copy()

    # Apply inflation
    df['Value'] = inflate_series(df['Value'])
    df.rename(columns={'Value':'Cost (Bil $)'}, inplace=True)

    cost_cat_map, cap_type_ls, op_type_ls = gather_cost_types(df)
    # Get the list of transmission investment categories (subset of cap_type_ls)
    trans_cap_type_ls = [
        c for c in cap_type_ls if c in [
            'inv_converter_costs',
            'inv_transmission_line_investment',
        ]
    ]
    nontrans_cap_type_ls = [c for c in cap_type_ls if c not in trans_cap_type_ls]

    # Calculate objective function system costs
    if cost_type == 'objective':
        # Multiply all capital costs by pvf_capital and operation by pvf_onm
        df = pd.merge(left=df, right=pvf_capital, how='left', on=['year'], sort=False)
        df = pd.merge(left=df, right=pvf_onm, how='left', on=['year'], sort=False)
        cap_cond = df['cost_cat'].isin(cap_type_ls)
        onm_cond = df['cost_cat'].isin(op_type_ls)
        df.loc[cap_cond, 'Cost (Bil $)'] *= df.loc[cap_cond, 'pvfcap']
        df.loc[onm_cond, 'Cost (Bil $)'] *= df.loc[onm_cond, 'pvfonm']
        df.drop(['pvfcap','pvfonm'], axis='columns', inplace=True)

    # Annualize if specified
    elif cost_type == 'annualized':
        # Turn each cost category into a column
        df = df.pivot_table(index=['year', 'r'], columns='cost_cat', values='Cost (Bil $)').reset_index()

        # Add rows for all regions and years (including extra years after end year for financial recovery)
        full_yrs = list(range(firstmodelyear - sys_eval_years, lastmodelyear + addyears + 1))
        allyrs = pd.DataFrame(list(product(full_yrs, val_r)), columns=['year','r'])
        df = pd.merge(allyrs, df, on=['year', 'r'], how='left')
    
        if not remove_existing:
            # Add payments for pre-2010 capacity; 
            # Keep data up until the year before the first modeled year
            df_capex_hist = df_capex_init.loc[
                df_capex_init.r.isin(val_r)].groupby(['year','r']).capex.sum()

            # Convert to billion dollars and remove years after the first modeled year
            df_capex_hist=df_capex_hist.loc[:firstmodelyear-1,:]

            # Convert to output dollar year
            df_capex_hist = pd.DataFrame(inflate_series(df_capex_hist))

            # Insert into full cost table
            df = df.join(df_capex_hist, on=['year','r'])

            # Add historical capex only where valid (non-null)
            mask = (df['year'] < firstmodelyear) 
            df.loc[mask, 'inv_investment_capacity_costs'] = df.loc[mask, 'capex']
            df = df.drop('capex',axis=1)

        # For capital costs, multiply by CRF to annualize, and sum over previous 20 years.
        # This requires 20 years before 2010 to sum properly.
        if crf_from_user:
            crf = pd.DataFrame({
                'year': full_yrs,
                'crf': (
                    discount_rate * (1 + discount_rate)**sys_eval_years
                    / ((1 + discount_rate)**sys_eval_years - 1)
                )
            }).set_index('year')

        # Otherwise use the crf from model run
        else:
            crf = crf_in.copy()
            crf = crf.set_index('year').reindex(full_yrs)
            crf = crf.interpolate(method='linear')
            crf['crf'] = crf['crf'].fillna(method='bfill')
            
        if shift_capital:
            # This means we start capital payments in the year of the investment,
            # even though loan payments typically start in the following year, so that
            # investments made in 2050 are still reflected in 2050 capital payments.
            # This requires dividing the crf by discount rate to result in the
            # same present value calculation.
            crf = crf / (1 + discount_rate)
        else:
            # This method reflects typical loan payments that start in the year after the loan.
            df[cap_type_ls] = df.groupby('r')[cap_type_ls].shift()

        df = pd.merge(left=df.set_index('year'), right=crf, how='left',on=['year'], sort=False)
        df[cap_type_ls] = df[cap_type_ls].fillna(0)
        
        df[cap_type_ls] = df[cap_type_ls].multiply(df["crf"], axis="index")

        # Sort to ensure rolling sum is correct
        df = df.reset_index().set_index(['r','year']).sort_index()

        df[nontrans_cap_type_ls] = (
            df.groupby('r')[nontrans_cap_type_ls].rolling(sys_eval_years).sum()
        ).reset_index(level=0, drop=True)

        df[trans_cap_type_ls] = (
            df.groupby('r')[trans_cap_type_ls].rolling(trans_crp).sum()
        ).reset_index(level=0, drop=True)
        df = df.reset_index()

        # Remove years before first modeled year
        keep_yrs = list(range(firstmodelyear, lastmodelyear + addyears + 1))
        df = df.loc[df.year.isin(keep_yrs)]

        # For operation costs, simply fill missing years with model year values.
        # Assuming sys_eval_years = 20, operation payments last for 20 yrs starting in the
        # modeled year, so fill 19 empty years after the modeled year.

        # Set NaN values to 0 for missing costs in simulation years.
        # This approach prevents incorrect forward-filling from earlier years into years
        # where the cost is actually zero (since GAMS drops all zeros).
        df.loc[df['year'].isin(sim_years), op_type_ls] = (
            df.loc[df['year'].isin(sim_years), op_type_ls]
            .fillna(0)
        )

        df.loc[:,op_type_ls] = (
            df.groupby('r')[op_type_ls]
            .fillna(method='ffill', limit=sys_eval_years-1)
        )
        df = df.fillna(0)

        df = pd.melt(
            df.reset_index(),
            id_vars=['year', 'r'],
            value_vars=cap_type_ls + op_type_ls,
            var_name='cost_cat',
            value_name='Cost (Bil $)'
        )

        # Add Discounted Cost column 
        df['Discounted Cost (Bil $)'] = (
            df['Cost (Bil $)']
            / (1 + discount_rate)**(df['year'] - present_value_year)
        )

    else:
        raise ValueError(
            'Invalid cost_type. Valid options are "annualized" or "objective".'
        )

    if group_r:
        df = df.groupby(['year', 'cost_cat'], as_index=False).sum()
        df.drop(columns=['r'], inplace=True)

    if through_2050:
        df = df.loc[(df.year >= present_value_year)&(df.year <= 2050)]

    if rename_as_bokeh:
        cost_cat_map = cost_cat_map.set_index('raw')["display"].to_dict()
        df['cost_cat'] = df['cost_cat'].map(lambda x: cost_cat_map.get(x, x))

        df = df.groupby(['year', 'r', 'cost_cat'], as_index=False).sum()

    # Remove rows with zero values to reduce file size
    df = df.loc[df['Cost (Bil $)'] != 0].reset_index(drop=True)
    return df


def calc_tech_trans(case: str) -> pd.DataFrame:
    """
    Compute total reinforcement, and spur line transmission by region.
    (TW-mi)

    Args:
        case (str): Path to the ReEDS case directory or outputs.h5 file.

    Returns:
        pd.DataFrame: Transmission capacity DataFrame (TW-mi) for reinforcement, spur line.
    """
    # Identify the inputs_case directory based on provided path
    inputs_case = (
        os.path.abspath(os.path.join(case, '..', '..', 'inputs_case'))
        if case.endswith('.h5')
        else os.path.join(case, 'inputs_case')
    )

    # Load scalar parameters and spur-line distances
    scalars = pd.read_csv(
        os.path.join(inputs_case, 'scalars.csv'),
        header=None, usecols=[0, 1], index_col=0
    ).squeeze()

    # Valid regions
    val_r = pd.read_csv(
        os.path.join(inputs_case, 'val_r.csv'),
        header=None,
    ).squeeze(1).values

    spur_parameters = pd.read_csv(os.path.join(inputs_case, 'spur_parameters.csv'))

    # Load ReEDS outputs
    cap_new_bin_out = reeds.io.read_output(case, 'cap_new_bin_out')

    cap_new_bin_out['Value'] = cap_new_bin_out['Value'] * 1e-6  # Convert to TW
    cap_new_bin_out.rename(
        columns={'t': 'year', 'Value': 'New Cap (TW)'},
        inplace=True)

    # Sum new capacity by technology and resource bin
    cap_new_bin_out = cap_new_bin_out.groupby(
        ['i', 'r', 'rscbin', 'year'], as_index=False)['New Cap (TW)'].sum()

    # Convert UPV from DC to AC
    upv_mask = cap_new_bin_out['i'].str.contains('upv')
    cap_new_bin_out.loc[upv_mask, 'New Cap (TW)'] *= float(scalars['ilr_utility'])

    # Filter technologies of interest
    cap_new_filtered = cap_new_bin_out[cap_new_bin_out['i'].str.startswith(('wind', 'upv', 'csp'))]

    # Merge spur parameters
    cap_new_filtered = cap_new_filtered.merge(spur_parameters, on=['i', 'r', 'rscbin'], how='left')

    # Compute spur and reinforcement distances (TW-mi)
    tech_trans = pd.concat([
        cap_new_filtered[['year', 'r', 'New Cap (TW)', 'dist_spur_km']]
        .rename(columns={'dist_spur_km': 'dist'})
        .assign(trtype='spur'),

        cap_new_filtered[['year', 'r', 'New Cap (TW)', 'dist_reinforcement_km']]
        .rename(columns={'dist_reinforcement_km': 'dist'})
        .assign(trtype='reinforcement')
    ])

    tech_trans['Trans (TW-mi)'] = tech_trans['New Cap (TW)'] * tech_trans['dist'] / 1.60934  # km to mi

    # Compute cumulative TW-mi per year and transmission type
    years = sorted(cap_new_bin_out.year.unique())
    full_idx = pd.MultiIndex.from_product(
        [years, val_r, ['spur', 'reinforcement']], names=['year', 'r', 'trtype'])

    tech_trans_out = (
        tech_trans.groupby(['year', 'r' , 'trtype'])['Trans (TW-mi)'].sum()
        # cap_new_bin_out is new additions, so fill missing years with zero
        .reindex(full_idx).fillna(0)
        # Do a cumulative sum of the transmission capacity 
        # for each region and transmission type over the years
        .groupby(['trtype', 'r']).cumsum()
        .reset_index()
        .sort_values(by=['r', 'trtype', 'year'])
    )

    if True:
        """
        Due to some old ReEDS outputs runs we need to zero out 
        the reinforcement line for county level regions 
        (This can be removed in the future if countly level
        supply curves get reinforcement distance zeroed out)
        """
        # Get county level regions
        county_regions = pd.read_csv(
            os.path.join(inputs_case, 'county2zone.csv'),
            dtype={'FIPS':str},
        )

        county_regions['county'] = 'p' + county_regions.FIPS

        # Set reinforcement distance to zero for county level regions
        tech_trans_out = tech_trans_out.loc[
            ~tech_trans_out.r.isin(county_regions['county'])
        ]        

    return tech_trans_out

# %%
