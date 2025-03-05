#%% Imports
import numpy as np
import pandas as pd
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import patheffects as pe
import os
import sys
import io
import argparse
import site
import subprocess as sp
import platform
from glob import glob
from tqdm import tqdm
import traceback
import cmocean
import pptx
from pptx.util import Inches, Pt
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))
import reeds
from reeds import plots
from reeds import reedsplots
from bokehpivot.defaults import DEFAULT_DOLLAR_YEAR, DEFAULT_PV_YEAR, DEFAULT_DISCOUNT_RATE

reeds_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

### Format plots and load other convenience functions
site.addsitedir(os.path.join(reeds_path,'postprocessing'))

plots.plotparams()

#%% Argument inputs
parser = argparse.ArgumentParser(
    description='Compare multiple ReEDS cases',
    formatter_class=argparse.ArgumentDefaultsHelpFormatter,
)
parser.add_argument(
    'caselist', type=str, nargs='+',
    help=('space-delimited list of cases to plot, OR shared casename prefix, '
          'OR csv file of cases. The first case is treated as the base case '
          'unless a different one is provided via the --basecase/-b argument.'))
parser.add_argument(
    '--casenames', '-n', type=str, default='',
    help='comma-delimited list of shorter case names to use in plots')
parser.add_argument(
    '--titleshorten', '-s', type=str, default='',
    help='characters to cut from start of case name (only used if no casenames)')
parser.add_argument(
    '--startyear', '-t', type=int, default=2020,
    help='First year to show')
parser.add_argument(
    '--sharey', '-y', action='store_true',
    help='Use same y-axis scale for absolute and difference plots')
parser.add_argument(
    '--basecase', '-b', type=str, default='',
    help='Substring of case path to use as default (if empty, uses first case in list)')
parser.add_argument(
    '--skipbp', '-p', action='store_true',
    help='flag to prevent bokehpivot report from being generated')
parser.add_argument(
    '--bpreport', '-r', type=str, default='standard_report_reduced',
    help='which bokehpivot report to generate')
parser.add_argument(
    '--detailed', '-d', action='store_true',
    help='Include more detailed plots')
parser.add_argument(
    '--forcemulti', '-m', action='store_true',
    help='Always use multi-case plots (even for 2 cases)')
parser.add_argument(
    '--lesslabels', '-l', action='count',
    help='Add less value labels to plots')
parser.add_argument(
    '--nowrap', '-w', action='store_true',
    help="Don't wrap subplot titles")

args = parser.parse_args()
_caselist = args.caselist
_casenames = args.casenames
try:
    titleshorten = int(args.titleshorten)
except ValueError:
    titleshorten = len(args.titleshorten)
_basecase = args.basecase
startyear = args.startyear
sharey = True if args.sharey else 'row'
bpreport = args.bpreport
skipbp = args.skipbp
detailed = args.detailed
forcemulti = args.forcemulti
lesslabels = args.lesslabels
nowrap = args.nowrap
interactive = False

#%% Inputs for testing
# reeds_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# _caselist = [os.path.join(reeds_path,'postprocessing','example.csv')]
# _casenames = ''
# titleshorten = 0
# startyear = 2020
# sharey = 'row'
# _basecase = ''
# skipbp = True
# bpreport = 'standard_report_reduced'
# interactive = True
# forcemulti = False
# lesslabels = 0
# nowrap = False
# detailed = False


#%%### Fixed inputs
cmap = cmocean.cm.rain
cmap_diff = plt.cm.RdBu_r
## https://www.whitehouse.gov/wp-content/uploads/2023/11/CircularA-4.pdf
discountrate_social = DEFAULT_DISCOUNT_RATE
## https://www.epa.gov/environmental-economics/scghg
discountrate_scghg = 0.02
assert discountrate_scghg in [0.015, 0.02, 0.025]
central_health = {'cr':'ACS', 'model':'EASIUR'}
reeds_dollaryear = 2004
output_dollaryear = DEFAULT_DOLLAR_YEAR
startyear_notes = DEFAULT_PV_YEAR
SLIDE_HEIGHT = 6.88
SLIDE_WIDTH = 13.33

colors_social = {
    'CO2': plt.cm.tab20b(4),
    'CH4': plt.cm.tab20b(5),
    'health': plt.cm.tab20b(7),
}

techmap = {
    **{f'upv_{i}':'Utility PV' for i in range(20)},
    **{f'dupv_{i}':'Utility PV' for i in range(20)},
    **{f'wind-ons_{i}':'Land-based wind' for i in range(20)},
    **{f'wind-ofs_{i}':'Offshore wind' for i in range(20)},
    **dict(zip(['nuclear','nuclear-smr'], ['Nuclear']*20)),
    **dict(zip(
        ['gas-cc_re-cc','gas-ct_re-ct','re-cc','re-ct',
         'gas-cc_h2-ct','gas-ct_h2-ct','h2-cc','h2-ct',],
        ['H2 turbine']*20)),
    **{f'battery_{i}':'Battery/PSH' for i in range(20)}, **{'pumped-hydro':'Battery/PSH'},
    **dict(zip(
        ['coal-igcc', 'coaloldscr', 'coalolduns', 'gas-cc', 'gas-ct', 'coal-new', 'o-g-s',],
        ['Fossil']*20)),
    **dict(zip(
        ['gas-cc_gas-cc-ccs_mod','gas-cc_gas-cc-ccs_max','gas-cc-ccs_mod','gas-cc-ccs_max',
         'gas-cc_gas-cc-ccs_mod','coal-igcc_coal-ccs_mod','coal-new_coal-ccs_mod',
        'coaloldscr_coal-ccs_mod','coalolduns_coal-ccs_mod','cofirenew_coal-ccs_mod',
        'cofireold_coal-ccs_mod','gas-cc_gas-cc-ccs_max','coal-igcc_coal-ccs_max',
        'coal-new_coal-ccs_max','coaloldscr_coal-ccs_max','coalolduns_coal-ccs_max',
        'cofirenew_coal-ccs_max','cofireold_coal-ccs_max',],
        ['Fossil+CCS']*50)),
    **dict(zip(['dac','beccs_mod','beccs_max'],['CO2 removal']*20)),
}

maptechs = [
    'Utility PV',
    'Land-based wind',
    'Offshore wind',
    'Nuclear',
    'H2 turbine',
    'Battery/PSH',
    'Fossil+CCS',
]

plotdiffvals = [
    'Generation (TWh)',
    'Capacity (GW)',
    'New Annual Capacity (GW)',
    'Annual Retirements (GW)',
    'Firm Capacity (GW)',
    'Curtailment Rate',
    'Transmission (GW-mi)',
    'Transmission (PRM) (GW-mi)',
    'Bulk System Electricity Pric',
    'National Average Electricity',
    'Present Value of System Cost',
    'NEUE (ppm)',
    'Runtime (hours)',
    'Runtime by year (hours)',
]
onlytechs = None
i_plots = ['wind-ons','upv','battery','h2-ct','nuclear','gas-cc-ccs','coal-ccs',]
## mapdiff: 'cap' or 'gen_ann'
mapdiff = 'cap'



#%%### Functions
def add_to_pptx(
        title=None, file=None, left=0, top=0.62, width=SLIDE_WIDTH, height=None,
        verbose=1, slide=None,
    ):
    """Add current matplotlib figure (or file if specified) to new powerpoint slide"""
    if not file:
        image = io.BytesIO()
        plt.savefig(image, format='png')
    else:
        image = file
        if not os.path.exists(image):
            raise FileNotFoundError(image)

    if slide is None:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = title
    slide.shapes.add_picture(
        image,
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    if verbose:
        print(title)
    return slide


def add_textbox(
        text, slide,
        left=0, top=7.2, width=SLIDE_WIDTH, height=0.3,
        fontsize=14,
    ):
    """Add a textbox to the specified slide"""
    textbox = slide.shapes.add_textbox(
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    p = textbox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    return slide


def plot_bars_abs_stacked(
        dfplot, basecase, colors, ax, col=0,
        net=True, label=True, ypad=0.02, fontsize=9,
    ):
    """
    * ax must have at least 2 rows
    * dfplot must have cases as rows and stacked bar elements (matching colors) as cols
    """
    ## Absolute and difference
    if isinstance(basecase, str):
        dfdiff = dfplot - dfplot.loc[basecase]
    elif isinstance(basecase, list):
        dfdiff = dfplot - dfplot.loc[basecase].values
    elif isinstance(basecase, dict):
        dfdiff = dfplot - dfplot.loc[basecase.values()].values

    for (row, df) in enumerate([dfplot, dfdiff]):
        plots.stackbar(df=df, ax=ax[row,col], colors=colors, net=(net or row), width=0.8)
        ymin, ymax = ax[row,col].get_ylim()
        _ypad = (ymax - ymin) * ypad
        ## label net value
        if label:
            for x, case in enumerate(df.index):
                val = df.loc[case].sum()
                if np.around(val, 0) == 0:
                    continue
                ax[row,col].annotate(
                    f'{val:.0f}', (x, val - _ypad), ha='center', va='top',
                    color='k', size=fontsize,
                    path_effects=[pe.withStroke(linewidth=2.0, foreground='w', alpha=0.7)],
                )
    ## Legend info
    legend_handles = [
        mpl.patches.Patch(facecolor=colors[i], edgecolor='none', label=i)
        for i in (colors if isinstance(colors, dict) else colors.index) if i in dfplot
    ]
    return legend_handles


#%%### Procedure
#%% Parse arguments
use_table_casenames = False
use_table_colors = False
use_table_bases = False
if len(_caselist) == 1:
    ## If it's a .csv, read the cases to compare
    if _caselist[0].endswith('.csv'):
        dfcase = pd.read_csv(_caselist[0], header=None, comment='#', quoting=3)
        ## First check it's a simple csv with one case per row
        if dfcase.shape[1] == 1:
            caselist = dfcase[0].tolist()
        ## Then check if it's a csv with [casepath,casename] in the header
        elif (
            ('casepath' in dfcase.loc[0].tolist())
             and ('casename' in dfcase.loc[0].tolist())
        ):
            dfcase = dfcase.T.set_index(0).T
            ## Drop cases that haven't finished yet
            unfinished = dfcase.loc[
                ~dfcase.casepath.map(
                    lambda x: os.path.isfile(os.path.join(x,'outputs','reeds-report','report.xlsx')))
            ].index
            if len(unfinished):
                print('The following cases have not yet finished:')
                print('\n'.join(dfcase.loc[unfinished].casepath.tolist()))
            dfcase = dfcase.drop(unfinished).copy()
            caselist = dfcase.casepath.tolist()
            use_table_casenames = True
            if 'color' in dfcase:
                if not dfcase.color.isnull().any():
                    use_table_colors = True
            if 'base' in dfcase:
                if not dfcase.base.isnull().any():
                    use_table_bases = True
        ## Otherwise assume it's a copy of a cases_{batchname}.csv file in a case folder
        ## This approach is less robust; the others are preferred.
        else:
            prefix_plus_tail = os.path.dirname(_caselist[0])
            tails = [i for i in dfcase.iloc[0] if i not in ['Default Value',np.nan]]
            prefix = prefix_plus_tail[:-len([i for i in tails if prefix_plus_tail.endswith(i)][0])]
            caselist = [prefix+i for i in tails]
    ## Otherwise look for all runs starting with the provided string
    else:
        caselist = sorted(glob(_caselist[0]+'*'))
        ## If no titleshorten is provided, use the provided prefix
        if not titleshorten:
            titleshorten = len(os.path.basename(_caselist))
else:
    caselist = _caselist

## Remove cases that haven't finished yet
caselist = [
    i for i in caselist
    if os.path.isfile(os.path.join(i,'outputs','reeds-report','report.xlsx'))
]

## Get the casenames
if use_table_casenames:
    casenames = [c.replace('\\n','\n') for c in dfcase.casename.tolist()]
else:
    casenames = (
        _casenames.split(',') if len(_casenames)
        else [os.path.basename(c)[titleshorten:] for c in caselist]
    )

if len(caselist) != len(casenames):
    err = (
        f"len(caselist) = {len(caselist)} but len(casenames) = {len(casenames)}\n\n"
        'caselist:\n' + '\n'.join(caselist) + '\n\n'
        'casenames:\n' + '\n'.join(casenames) + '\n'
    )
    raise ValueError(err)

cases = dict(zip(casenames, caselist))
maxlength = max([len(c) for c in cases])

# check to ensure there are at least two cases
if len(cases) <= 1: 
    err = f"There are less than two cases being compared: {', '.join(cases.values())}"
    raise ValueError(err)

### Get the base cases
if not len(_basecase):
    basecase = list(cases.keys())[0]
else:
    basepath = [c for c in cases.values() if c.endswith(_basecase)]
    if len(basepath) == 0:
        err = (
            f"Use a basecase that matches one case.\nbasecase={_basecase} matches none of:\n"
            + '\n'.join(basepath)
        )
        raise ValueError(err)
    elif len(basepath) > 1:
        err = (
            f"Use a basecase that only matches one case.\nbasecase={_basecase} matches:\n"
            + '\n'.join(basepath)
        )
        raise ValueError(err)
    else:
        basepath = basepath[0]
        ## basecase is the short name; basepath is the full path
        basecase = casenames[caselist.index(basepath)]
        ## Put it first in the list
        cases = {**{basecase:cases[basecase]}, **{k:v for k,v in cases.items() if k != basecase}}

## Make case->base dictionary
if use_table_bases:
    basemap = dfcase.set_index('casename').base.to_dict()
else:
    basemap = dict(zip(cases, [basecase]*len(cases)))

## Get the colors
if use_table_colors:
    colors = dict(zip(dfcase.casename, dfcase.color))
    for k, v in colors.items():
        if v.startswith('plt.cm.') or v.startswith('cmocean.cm.'):
            colors[k] = eval(v)
else:
    colors = plots.rainbowmapper(cases)

## Arrange the maps
nrows, ncols, coords = plots.get_coordinates(cases, aspect=2)

## Take a look
print('Analyzing the following cases:')
for case, path in cases.items():
    print(
        f'{path} -> {case}'
        + (' (base)' if ((not use_table_bases) and (case == basecase)) else '')
    )

#%% Create output folder
firstcasepath = list(cases.values())[0]
outpath = os.path.join(firstcasepath, 'outputs', 'comparisons')
os.makedirs(outpath, exist_ok=True)
## Remove disallowed characters and clip filename to max length
max_filename_length = 250
savename = os.path.join(
    outpath,
    (f"results-{','.join(cases.keys())}"
     .replace(':','').replace('/','').replace(' ','').replace('\\n','').replace('\n','')
     [:max_filename_length-len('.pptx')]) + '.pptx'
)
print(f'Saving results to {savename}')

#%% Create bokehpivot report as subprocess
if not skipbp:
    start_str = 'start ' if platform.system() == 'Windows' else ''
    bp_path = f'{reeds_path}/postprocessing/bokehpivot'
    bp_py_file = f'{bp_path}/reports/interface_report_model.py'
    report_path = f'{bp_path}/reports/templates/reeds2/{bpreport}.py'
    bp_outpath = f'{outpath}/{bpreport}-diff-multicase'
    add_diff = 'Yes'
    auto_open = 'Yes'
    bp_colors = pd.read_csv(f'{bp_path}/reeds_scenarios.csv')['color'].tolist()
    bp_colors = bp_colors*10 #Up to 200 scenarios
    bp_colors = bp_colors[:len(casenames)]
    df_scenarios = pd.DataFrame({'name':casenames, 'color':bp_colors, 'path':caselist})
    scenarios_path = f'{outpath}/scenarios.csv'
    df_scenarios.to_csv(scenarios_path, index=False)
    call_str = (
        f'{start_str}python "{bp_py_file}" "ReEDS 2.0" "{scenarios_path}" all ' +
        f'{add_diff} "{basecase}" "{report_path}" "html,excel" one "{bp_outpath}" {auto_open}'
    )
    sp.Popen(call_str, shell=True)

#%%### Load data
#%% Shared
sw = reeds.io.get_switches(cases[case])
scalars = reeds.io.get_scalars(cases[case])
phaseout_trigger = float(scalars.co2_emissions_2022) * float(sw.GSw_TCPhaseout_trigger_f)

inflatable = reeds.io.get_inflatable(os.path.join(
    reeds_path,'inputs','financials','inflation_default.csv'))
inflator = inflatable[reeds_dollaryear, output_dollaryear]

scghg = pd.read_csv(
    os.path.join(reeds_path, 'postprocessing', 'plots', 'scghg_annual.csv'),
    comment='#', thousands=','
).rename(columns={
    'gas':'e',
    'emission.year':'t',
    '2.5% Ramsey':'2020_2.5%',
    '2.0% Ramsey':'2020_2.0%',
    '1.5% Ramsey':'2020_1.5%',
}).set_index(['e','t'])
scghg_central = (
    scghg[f'2020_{discountrate_scghg*100:.1f}%'].unstack('e')
    * inflatable[2020, output_dollaryear]
)

#%% Colors
bokehcostcolors = pd.read_csv(
    os.path.join(
        reeds_path,'postprocessing','bokehpivot','in','reeds2','cost_cat_style.csv'),
    index_col='order').squeeze(1)
bokehcostcolors = bokehcostcolors.loc[~bokehcostcolors.index.duplicated()]

colors_time = pd.read_csv(
    os.path.join(
        reeds_path,'postprocessing','bokehpivot','in','reeds2','process_style.csv'),
    index_col='order',
).squeeze(1)

bokehcolors = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_style.csv'),
    index_col='order').squeeze(1)

tech_map = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_map.csv'),
    index_col='raw').squeeze(1)
    
bokehcolors = pd.concat([
    bokehcolors.loc['smr':'electrolyzer'],
    pd.Series('#D55E00', index=['dac'], name='color'),
    bokehcolors.loc[:'Canada'],
])

bokehcolors['canada'] = bokehcolors['Canada']

techcolors = {
    'gas-cc_gas-cc-ccs':bokehcolors['gas-cc-ccs_mod'],
    'cofire':bokehcolors['biopower'],
    'gas-cc':'#5E1688',
    'gas-cc-ccs':'#9467BD',
}
for i in bokehcolors.index:
    if i in techcolors:
        pass
    elif i in bokehcolors.index:
        techcolors[i] = bokehcolors[i]
    else:
        raise Exception(i)

techcolors = {i: techcolors[i] for i in bokehcolors.index}

trtype_map = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','trtype_map.csv'),
    index_col='raw')['display']
colors_trans = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','trtype_style.csv'),
    index_col='order')['color']

#%% Parse excel report sheet names
val2sheet = reeds.io.get_report_sheetmap(cases[basecase])

#%% Read input files
renametechs = {
    'h2-cc_upgrade':'h2-cc',
    'h2-ct_upgrade':'h2-ct',
    'gas-cc-ccs_mod_upgrade':'gas-cc-ccs_mod',
    'coal-ccs_mod_upgrade':'coal-ccs_mod',
}
dictin_sw = {case: reeds.io.get_switches(cases[case]) for case in cases}

hierarchy = {case: reeds.io.get_hierarchy(cases[case]) for case in cases}

dictin_error = {}
for case in tqdm(cases, desc='system cost error'):
    dictin_error[case] = reeds.io.read_output(cases[case], 'error_check').set_index('*').squeeze(1)

dictin_cap = {}
for case in tqdm(cases, desc='national capacity'):
    dictin_cap[case] = reeds.io.read_report(cases[case], 'Capacity (GW)', val2sheet)
    ### Simplify techs
    dictin_cap[case].tech = dictin_cap[case].tech.map(lambda x: renametechs.get(x,x))
    dictin_cap[case] = (
        dictin_cap[case].groupby(['tech','year'], as_index=False)
        ['Capacity (GW)'].sum())
    dictin_cap[case] = dictin_cap[case].loc[
        ~dictin_cap[case].tech.isin(['electrolyzer','smr','smr-ccs','dac'])].copy()

dictin_gen = {}
for case in tqdm(cases, desc='national generation'):
    dictin_gen[case] = reeds.io.read_report(cases[case], 'Generation (TWh)', val2sheet)
    ### Simplify techs
    dictin_gen[case].tech = dictin_gen[case].tech.map(lambda x: renametechs.get(x,x))
    dictin_gen[case] = (
        dictin_gen[case].groupby(['tech','year'], as_index=False)
        ['Generation (TWh)'].sum())

costcat_rename = {
    'CO2 Spurline':'CO2 T&S Capex',
    'CO2 Pipeline':'CO2 T&S Capex',
    'CO2 Storage':'CO2 T&S Capex',
    'CO2 Spurline FOM':'CO2 T&S O&M',
    'CO2 Pipeline FOM':'CO2 T&S O&M',
    'CO2 Incentive Payments':'CCS Incentives',
    'Capital': 'Gen & Stor Capex',
    'O&M': 'Gen & Stor O&M',
    'CO2 Network':'CO2 T&S Capex',
    'CO2 Incentives':'CCS Incentives',
    'CO2 FOM':'CO2 T&S O&M',
    'CO2 Capture':'CO2 T&S Capex',
    'H2 Fuel':'Fuel',
    'H2 VOM':'H2 Prod O&M',
}
dictin_npv = {}
for case in tqdm(cases, desc='NPV of system cost'):
    dictin_npv[case] = (
        reeds.io.read_report(cases[case], 'Present Value of System Cost', val2sheet)
        .set_index('cost_cat')['Discounted Cost (Bil $)']
    )
    dictin_npv[case].index = pd.Series(dictin_npv[case].index).replace(costcat_rename)
    dictin_npv[case] = dictin_npv[case].groupby(level=0, sort=False).sum()

dictin_scoe = {}
for case in tqdm(cases, desc='SCOE'):
    dictin_scoe[case] = reeds.io.read_report(cases[case], 'National Average Electricity', val2sheet)
    dictin_scoe[case].cost_cat = dictin_scoe[case].cost_cat.replace(
        {**costcat_rename,**{'CO2 Incentives':'CCS Incentives'}})
    dictin_scoe[case] = (
        dictin_scoe[case].groupby(['cost_cat','year'], sort=False, as_index=False)
        ['Average cost ($/MWh)'].sum())

dictin_syscost = {}
for case in tqdm(cases, desc='annual system cost'):
    dictin_syscost[case] = reeds.io.read_report(cases[case], 'Undiscounted Annualized Syst', val2sheet)
    dictin_syscost[case].cost_cat = dictin_syscost[case].cost_cat.replace(
        {**costcat_rename,**{'CO2 Incentives':'CCS Incentives'}})
    dictin_syscost[case] = (
        dictin_syscost[case].groupby(['cost_cat','year'], sort=False)
        ['Cost (Bil $)'].sum().unstack('cost_cat'))

dictin_emissions = {}
for case in tqdm(cases, desc='national emissions'):
    dictin_emissions[case] = (
        reeds.io.read_output(cases[case], 'emit_nat', valname='ton')
        .set_index(['e','t']).squeeze(1).unstack('e')
    )

dictin_trans = {}
for case in tqdm(cases, desc='national transmission'):
    dictin_trans[case] = reeds.io.read_report(cases[case], 'Transmission (GW-mi)')

dictin_trans_r = {}
for case in tqdm(cases, desc='regional transmission'):
    dictin_trans_r[case] = reeds.io.read_output(cases[case], 'tran_out', valname='MW')
    for _level in ['interconnect','transreg','transgrp','st']:
        dictin_trans_r[case][f'inter_{_level}'] = (
            dictin_trans_r[case].r.map(hierarchy[case][_level])
            != dictin_trans_r[case].rr.map(hierarchy[case][_level])
        ).astype(int)

dictin_cap_r = {}
for case in tqdm(cases, desc='regional capacity'):
    dictin_cap_r[case] = reeds.io.read_output(cases[case], 'cap', valname='MW')
    ### Simplify techs
    dictin_cap_r[case].i = dictin_cap_r[case].i.map(lambda x: renametechs.get(x,x))
    dictin_cap_r[case].i = dictin_cap_r[case].i.str.lower().map(lambda x: techmap.get(x,x))
    dictin_cap_r[case] = dictin_cap_r[case].groupby(['i','r','t'], as_index=False).MW.sum()

dictin_cap_firm = {}
for case in tqdm(cases, desc='firm capacity'):
    dictin_cap_firm[case] = reeds.io.read_output(cases[case], 'cap_firm', valname='MW')
    ### Simplify techs
    dictin_cap_firm[case].i = reedsplots.simplify_techs(dictin_cap_firm[case].i)
    dictin_cap_firm[case] = dictin_cap_firm[case].groupby(['i','r','ccseason','t'], as_index=False).MW.sum()

dictin_runtime = {}
for case in tqdm(cases, desc='runtime'):
    dictin_runtime[case] = (
        reeds.io.read_report(cases[case], 'Runtime by year (hours)', val2sheet)
        .drop(columns='Net Level processtime')
    )

dictin_neue = {}
for case in tqdm(cases, desc='NEUE'):
    infiles = sorted(glob(os.path.join(cases[case],'outputs','neue_*.csv')))
    if not len(infiles):
        continue
    df = {}
    for f in infiles:
        y, i = [int(s) for s in os.path.basename(f).strip('neue_.csv').split('i')]
        df[y,i] = pd.read_csv(f)
        df[y,i] = df[y,i].loc[
            (df[y,i].level=='country') & (df[y,i].metric=='sum'),
            'NEUE_ppm'
        ].iloc[0]
    df = (
        pd.Series(df, name='NEUE [ppm]').rename_axis(['t','iteration'])
        .sort_index().reset_index()
        .drop_duplicates(subset='t', keep='last')
        .set_index('t')['NEUE [ppm]']
    )
    dictin_neue[case] = df

### Model years and discount rates
years = {}
yearstep = {}
for case in cases:
    years[case] = sorted(dictin_cap[case].year.astype(int).unique())
    years[case] = [y for y in years[case] if y >= startyear]
    yearstep[case] = years[case][-1] - years[case][-2]
lastyear = max(years[case])
## Years for which to add data notes
startyear_sums = 2023
allyears = range(startyear_sums,lastyear+1)
noteyears = [2035, 2050]
if all([lastyear < y for y in noteyears]):
    noteyears = [lastyear]
startyear_growth = 2035

discounts = pd.Series(
    index=range(startyear_notes,lastyear+1),
    data=[1/(1+discountrate_social)**(y-startyear_notes)
          for y in range(startyear_notes,lastyear+1)]
).rename_axis('t')

### Health impacts
dictin_health = {}
dictin_health_central = {}
dictin_health_central_mort = {}
for case in tqdm(cases, desc='health'):
    try:
        dictin_health[case] = (
            reeds.io.read_output(cases[case], 'health_damages_caused_r.csv')
            .groupby(['year','pollutant','model','cr']).sum()
        )
        dictin_health_central[case] = (
            dictin_health[case]
            .xs(central_health['cr'], level='cr')
            .xs(central_health['model'], level='model')
            .groupby('year').sum()
            ['damage_$']
            ### Inflate from reeds_dollaryear (2004) to bokeh output_dollaryear (2021)
            * inflator
            ### Convert to $B
            / 1e9
        )
        dictin_health_central_mort[case] = (
            dictin_health[case]
            .xs(central_health['cr'], level='cr')
            .xs(central_health['model'], level='model')
            .groupby('year').sum()
            ['mortality']
        )
    except (FileNotFoundError, KeyError) as err:
        print(f'Health impacts error for {case}: {err}')
        dictin_health_central[case] = (
            pd.Series(np.nan, index=years[case], name='damage_$')
            .rename_axis('year')
        )
        dictin_health_central_mort[case] = (
            pd.Series(np.nan, index=years[case], name='mortality')
            .rename_axis('year')
        )


#%% Detailed inputs
if detailed:
    ### Timeslice generation by region
    dictin_gen_h = {}
    for case in tqdm(cases, desc='gen_h'):
        dictin_gen_h[case] = reeds.io.read_output(cases[case], 'gen_h', valname='GW')
        dictin_gen_h[case].GW /= 1e3
        dictin_gen_h[case].i = reedsplots.simplify_techs(dictin_gen_h[case].i)
        dictin_gen_h[case] = dictin_gen_h[case].groupby(['i','r','h','t'], as_index=False).GW.sum()
        ## Separate charge and discharge
        dictin_gen_h[case].loc[
            (dictin_gen_h[case].i.str.startswith('battery')
            | dictin_gen_h[case].i.str.startswith('pumped-hydro'))
            & (dictin_gen_h[case].GW < 0),
            'i'
        ] += '|charge'
        dictin_gen_h[case].loc[
            (dictin_gen_h[case].i.str.startswith('battery')
            | dictin_gen_h[case].i.str.startswith('pumped-hydro'))
            & (~dictin_gen_h[case].i.str.endswith('|charge')),
            'i'
        ] += '|discharge'

    ### Aggregated generation by region
    dictin_gen_h_twh = {}
    for case in tqdm(dictin_gen_h):
        numhours = pd.read_csv(
            os.path.join(cases[case],'inputs_case','numhours.csv'),
        ).rename(columns={'*h':'h'}).set_index('h').squeeze(1)

        dictin_gen_h_twh[case] = dictin_gen_h[case].copy()
        dictin_gen_h_twh[case]['TWh'] = (
            dictin_gen_h_twh[case]['GW'] * dictin_gen_h_twh[case]['h'].map(numhours)
            / 1e3
        ).round(3)

        dictin_gen_h_twh[case] = dictin_gen_h_twh[case].groupby(['t','i']).TWh.sum().unstack('i')

    ### Stress period dispatch
    dictin_gen_h_stress = {}
    for case in tqdm(cases, desc='gen_h_stress'):
        dictin_gen_h_stress[case] = reeds.io.read_output(cases[case], 'gen_h_stress', valname='GW')
        dictin_gen_h_stress[case].GW /= 1e3
        dictin_gen_h_stress[case].i = reedsplots.simplify_techs(dictin_gen_h_stress[case].i)
        ## Separate charge and discharge
        dictin_gen_h_stress[case].loc[dictin_gen_h_stress[case].GW < 0,'i'] += '|charge'
        dictin_gen_h_stress[case].loc[dictin_gen_h_stress[case].i.isin(
            ['battery_4','battery_8','pumped-hydro']),'i'] += '|discharge'

    ### Stress period flows
    dictin_tran_flow_stress = {}
    for case in tqdm(cases, desc='tran_flow_stress'):
        dictin_tran_flow_stress[case] = reeds.io.read_output(
            cases[case], 'tran_flow_stress', valname='GW')
        dictin_tran_flow_stress[case].GW /= 1e3

    ### Stress period load
    dictin_load_stress = {}
    for case in tqdm(cases, desc='load_stress'):
        dictin_load_stress[case] = reeds.io.read_output(cases[case], 'load_stress', valname='GW')
        dictin_load_stress[case].GW /= 1e3

    ### Peak load (for capacity credit)
    distloss = 0.05
    dictin_peak_ccseason = {}
    for case in tqdm(cases, desc='peak_ccseason'):
        dictin_peak_ccseason[case] = pd.read_csv(
            os.path.join(cases[case],'inputs_case','peak_ccseason.csv'),
        ).rename(columns={'*r':'r', 'MW':'GW'})
        dictin_peak_ccseason[case].GW /= (1e3 * (1 - distloss))

    ### Capacity credit PRMTRADE
    dictin_prmtrade = {}
    for case in tqdm(cases, desc='prmtrade'):
        dictin_prmtrade[case] = reeds.io.read_output(cases[case], 'captrade', valname='GW')
        dictin_prmtrade[case].GW /= 1e3


#%%### Plots ######
### Set up powerpoint file
prs = pptx.Presentation(os.path.join(reeds_path,'postprocessing','template.pptx'))
blank_slide_layout = prs.slide_layouts[3]


#%%### System cost error
dfplot = pd.concat(dictin_error, axis=1).replace(0,np.nan).dropna(how='all').fillna(0)

ncols_err = len(dfplot)
data = {
    'z': {'title': 'System cost [fraction]', 'scale':1},
    'gen': {'title': 'Non-valgen generation [GWh]', 'scale':1e-3},
    'cap': {'title': 'Non-valcap capacity [GW]', 'scale':1e-3},
    'RPS': {'title': 'Non-RecMap RECS [GWh]', 'scale':1e-3},
    'OpRes': {'title': 'Non-valgen opres [MWh]', 'scale':1},
    'm_rsc_dat': {'title': 'Supply curve tweaks [GW]', 'scale':1e-3},
    'dropped': {'title': 'Dropped load [GWh]', 'scale':1e-3},
}
data = {k:v for k,v in data.items() if k in dfplot.index}

plt.close()
f,ax = plt.subplots(
    1, ncols_err,
    figsize=(min(ncols_err*3.5, SLIDE_WIDTH), max(3.75, 0.25*len(cases))),
)
for col, (datum, settings) in enumerate(data.items()):
    if datum not in dfplot.index:
        continue
    vals = dfplot.loc[datum] * settings['scale']
    _ax = ax if ncols_err == 1 else ax[col]
    _ax.bar(
        range(len(cases)),
        vals.values, color=[colors[c] for c in cases],
    )
    ## Formatting
    _ax.set_title(settings['title'], weight='bold', fontsize=14)
    _ax.set_xticks(range(len(cases)))
    _ax.set_xticklabels(cases.keys(), rotation=45, rotation_mode='anchor', ha='right')
    if _ax.get_ylim()[0] < 0:
        _ax.axhline(0, c='k', ls='--', lw=0.75)
    ## Notes
    for x, val in enumerate(vals):
        text = f"{val:.1e}" if datum == 'z' else f"{val:.0f}"
        _ax.annotate(text, (x, val), ha='center',
        xytext=(0, 2), textcoords='offset points',
    )
plots.despine(ax)

### Save it
title = 'Error check'
slide = add_to_pptx(title, width=None, height=SLIDE_HEIGHT)
if interactive:
    plt.show()


#%%### Generation capacity lines
aggtechsplot = {
    'Interregional\ntransmission': 'inter_transreg',
    'Land-based\nwind': ['wind-ons'],
    'Offshore\nwind': ['wind-ofs'],
    # 'Wind': ['wind-ons', 'wind-ofs'],
    'Solar': ['upv', 'dupv', 'distpv', 'csp', 'pvb'],
    'Battery': ['battery_{}'.format(i) for i in [2,4,6,8,10]],
    'Pumped\nstorage\nhydro': ['pumped-hydro'],
    # 'Storage': ['battery_{}'.format(i) for i in [2,4,6,8,10]] + ['pumped-hydro'],
    'Hydro, geo, bio': [
        'hydro','geothermal',
        'biopower','lfill-gas','cofire','beccs_mod','beccs'
    ],
    'Nuclear': ['nuclear', 'nuclear-smr'],
    'Hydrogen\nturbine': ['h2-cc', 'h2-cc-upgrade', 'h2-ct', 'h2-ct-upgrade'],
    'Gas CCS': ['gas-cc-ccs_mod'],
    'Coal CCS': ['coal-ccs_mod'],
    # 'Fossil\n(with CCS)': ['gas-cc-ccs_mod','coal-ccs_mod'],
    'Fossil\n(w/o CCS)': ['gas-cc', 'gas-ct', 'o-g-s', 'coal', 'cofire'],
#     'CDR': ['dac', 'beccs'],
#     'H2 production': ['smr', 'smr-ccs', 'electrolyzer'],
}
checktechs = [i for sublist in aggtechsplot.values() for i in sublist]
alltechs = pd.concat(dictin_cap).tech.unique()
printstring = (
    'The following techs are not plotted: '
    + ', '.join([c for c in alltechs if c not in checktechs])
)

offsetstart = {
    'Solar': (15,0),
    'Wind': (15,0), 'Land-based\nwind': (15,0),
}

val = '4_Capacity (GW)'
ycol = 'Capacity (GW)'
techrows, techcols = 2, len(aggtechsplot)//2+len(aggtechsplot)%2
techcoords = dict(zip(
    list(aggtechsplot.keys()),
    [(row,col) for row in range(techrows) for col in range(techcols)]
))

offset = dict()

plt.close()
f,ax = plt.subplots(
    techrows, techcols, sharex=True, sharey=True,
    figsize=(SLIDE_WIDTH, SLIDE_HEIGHT),
    gridspec_kw={'wspace':0.3, 'hspace':0.15},
)
df = {}
for tech in aggtechsplot:
    for case in cases:
        ### Central cases
        if 'transmission' in tech.lower():
            df[tech,case] = dictin_trans_r[case].loc[
                dictin_trans_r[case][aggtechsplot[tech]]==1
            ].groupby('t').MW.sum() / 1e3
        else:
            df[tech,case] = dictin_cap[case].loc[
                dictin_cap[case].tech.isin(aggtechsplot[tech])
            ].groupby('year')[ycol].sum().reindex(years[case]).fillna(0)
        ax[techcoords[tech]].plot(
            df[tech,case].index, df[tech,case].values,
            label=case, color=colors[case], ls='-',
        )
        ### Annotate the last value (with overlaps)
        fincap = df[tech,case].reindex([lastyear]).fillna(0).squeeze()
        ax[techcoords[tech]].annotate(
            ' {:.0f}'.format(fincap),
            (lastyear, fincap+offset.get((tech,case),0)),
            ha='left', va='center',
            color=colors[case], fontsize='small',
            annotation_clip=False,
        )
    ### Formatting
    ax[techcoords[tech]].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5 if lastyear>2040 else 1))
    ax[techcoords[tech]].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10 if lastyear>2040 else 5))
    ax[techcoords[tech]].annotate(
        tech,
        (0.05,1.0), va='top', ha='left',
        xycoords='axes fraction', 
        fontsize='x-large', weight='bold',)
    ### Annotate the 2020 value
    if len(df[tech,basecase]):
        plots.annotate(
            ax[techcoords[tech]], basecase,
            startyear, offsetstart.get(tech,(10,10)), color='C7',
            arrowprops={'arrowstyle':'-|>', 'color':'C7'})
if len(aggtechsplot) % 2:
    ax[-1,-1].axis('off')
## Legend
handles, labels = ax[-1,0].get_legend_handles_labels()
leg = ax[0,-1].legend(
    handles, labels,
    fontsize='large', frameon=False, 
    loc='center left', bbox_to_anchor=(1.1,-0.075),
    handletextpad=0.3, handlelength=0.7,
    ncol=1,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')
ax[techcoords[list(aggtechsplot.keys())[0]]].set_xlim(startyear,lastyear)
ax[techcoords[list(aggtechsplot.keys())[0]]].set_ylim(0)
ax[techcoords[list(aggtechsplot.keys())[0]]].set_ylabel('Capacity [GW]', y=-0.075)
# for row in range(techrows):
#     ax[row,0].set_ylabel('Capacity [GW]')
ax[techcoords[list(aggtechsplot.keys())[0]]].set_ylim(0)
# ## Annotate the last value (without overlaps)
# ymax = ax[techcoords[list(aggtechsplot.keys())[0]]].get_ylim()[1]
# df = pd.concat(df, axis=1)
# for tech in aggtechsplot:
#     plots.label_last(
#         df[tech], ax[techcoords[tech]], colors=colors, extend='both',
#         mindistance=ymax*0.05, fontsize='small',
#     )

plots.despine(ax)
plt.draw()
plots.shorten_years(ax[1,0])
### Save it
slide = add_to_pptx('Capacity')
add_textbox(printstring, slide)
if interactive:
    print(printstring)
    plt.show()


#%%### Capacity and generation bars
if (len(cases) == 2) and (not forcemulti):
    casebase, casecomp = list(cases.values())
    casebase_name, casecomp_name = list(cases.keys())
    for val in plotdiffvals:
        try:
            plt.close()
            f, ax, leg, dfdiff, printstring = reedsplots.plotdiff(
                val, casebase=casebase, casecomp=casecomp,
                casebase_name=casebase_name, casecomp_name=casecomp_name,
                onlytechs=onlytechs, titleshorten=titleshorten,
                yearmin=(2025 if 'NEUE' in val else startyear), yearmax=lastyear,
                # plot_kwds={'figsize':(4,4), 'gridspec_kw':{'wspace':0.7}},
            )
            slide = add_to_pptx(val, verbose=0)
            textbox = slide.shapes.add_textbox(
                left=Inches(0), top=Inches(7),
                width=Inches(SLIDE_WIDTH), height=Inches(0.5))
            textbox.text_frame.text = printstring
            if interactive:
                plt.show()
        except Exception as err:
            print(err)
else:
    toplot = {
        'Capacity': {
            'data': dictin_cap,
            'colors':techcolors,
            'columns':'tech',
            'values':'Capacity (GW)',
            'label':'Capacity [GW]'
        },
        'Generation': {
            'data': dictin_gen,
            'colors':techcolors,
            'columns':'tech',
            'values':'Generation (TWh)',
            'label':'Generation [TWh]'
        },
        'Runtime': {
            'data': dictin_runtime,
            'colors':colors_time.to_dict(),
            'columns':'process',
            'values':'processtime',
            'label':'Runtime [hours]'
        },
    }
    plotwidth = 2.0
    figwidth = plotwidth * len(cases)
    dfbase = {}
    for slidetitle, data in toplot.items():
        plt.close()
        f,ax = plt.subplots(
            2, len(cases), figsize=(figwidth, 6.8),
            sharex=True, sharey=sharey, dpi=None,
        )
        ax[0,0].set_ylabel(data['label'], y=-0.075)
        ax[0,0].set_xlim(2017.5, lastyear+2.5)
        ax[1,0].annotate(
            f'Diff\nfrom\n{basecase}', (0.03,0.03), xycoords='axes fraction',
            fontsize='x-large', weight='bold')
        ###### Absolute
        alltechs = set()
        for col, case in enumerate(cases):
            if case not in data['data']:
                continue
            dfplot = data['data'][case].pivot(index='year', columns=data['columns'], values=data['values'])
            dfplot = (
                dfplot[[c for c in data['colors'] if c in dfplot]]
                .round(3).replace(0,np.nan)
                .dropna(axis=1, how='all')
            )
            if case == basecase:
                dfbase[slidetitle] = dfplot.copy()
            alltechs.update(dfplot.columns)
            plots.stackbar(df=dfplot, ax=ax[0,col], colors=data['colors'], width=yearstep[case], net=False)
            ax[0,col].set_title(
                (case if nowrap else plots.wraptext(case, width=plotwidth*0.9, fontsize=14)),
                fontsize=14, weight='bold', x=0, ha='left', pad=8,)
            ax[0,col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
            ax[0,col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))


        ### Legend
        handles = [
            mpl.patches.Patch(
                facecolor=data['colors'][i], edgecolor='none',
                label=i.replace('Canada','imports').split('/')[-1]
            )
            for i in data['colors'] if i in alltechs
        ]
        leg = ax[0,-1].legend(
            handles=handles[::-1], loc='upper left', bbox_to_anchor=(1.0,1.0), 
            fontsize='medium', ncol=1,  frameon=False,
            handletextpad=0.3, handlelength=0.7, columnspacing=0.5, 
        )

        ###### Difference
        for col, case in enumerate(cases):
            ax[1,col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
            ax[1,col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))
            ax[1,col].axhline(0,c='k',ls='--',lw=0.75)

            if (case not in data['data']) or (case == basecase):
                continue
            dfplot = data['data'][case].pivot(index='year', columns=data['columns'], values=data['values'])
            dfplot = (
                dfplot
                .round(3).replace(0,np.nan)
                .dropna(axis=1, how='all')
            )
            dfplot = dfplot.subtract(dfbase[slidetitle], fill_value=0)
            dfplot = dfplot[[c for c in data['colors'] if c in dfplot]].copy()
            alltechs.update(dfplot.columns)
            plots.stackbar(df=dfplot, ax=ax[1,col], colors=data['colors'], width=yearstep[case], net=True)

        plots.despine(ax)
        plt.draw()
        plots.shorten_years(ax[1,0])
        ### Save it
        slide = add_to_pptx(slidetitle+' stack', width=min(figwidth, SLIDE_WIDTH))
        if interactive:
            plt.show()


#%% Alternate view: Stacks with bars labeled
barwidth = 0.35
labelpad = 0.08
width = 1.6*len(cases) + 0.5
aggstack = {
    **{f'battery_{i}':'Storage' for i in [2,4,6,8,10]},
    **{f'battery_{i}|charge':'Storage|charge' for i in [2,4,6,8,10]},
    **{f'battery_{i}|discharge':'Storage|discharge' for i in [2,4,6,8,10]},
    **{
        'pumped-hydro':'Storage',
        'pumped-hydro|charge':'Storage|charge', 'pumped-hydro|discharge':'Storage|discharge',

        'h2-cc':'H2 turbine', 'h2-ct':'H2 turbine',

        'beccs_mod':'Bio/BECCS',
        'biopower':'Bio/BECCS', 'lfill-gas':'Bio/BECCS', 'cofire':'Bio/BECCS',

        'gas-cc-ccs_mod':'Gas+CCS',
        'coal-ccs_mod':'Coal+CCS',
        'gas-cc':'Gas', 'gas-ct':'Gas', 'o-g-s':'Gas',
        'coal':'Coal',

        'Canada':'Canadian imports', 'canada':'Canadian imports',

        'hydro':'Hydro',
        'geothermal':'Geothermal',

        'csp':'CSP',
        'upv':'PV', 'dupv':'PV', 'distpv':'PV',
        'pvb':'PVB',

        'wind-ofs':'Offshore wind',
        'wind-ons':'Land-based wind',

        'nuclear':'Nuclear', 'nuclear-smr':'Nuclear',
    }
}
aggcolors = {
    'Nuclear':'C3',

    'Coal':plt.cm.binary(1.0),
    'Gas':plt.cm.tab20(8),
    'Coal+CCS':'C7',
    'Gas+CCS':plt.cm.tab20(9),

    'Hydro': techcolors['hydro'],
    'Geothermal': techcolors['geothermal'],
    'Canadian imports': techcolors['dr'],

    # 'Bio/BECCS':plt.cm.tab20(4),
    # 'H2 turbine':plt.cm.tab20(5),
    'Bio/BECCS':techcolors['biopower'],
    'H2 turbine':techcolors['h2-ct'],

    'Land-based wind':techcolors['wind-ons'],
    'Offshore wind':techcolors['wind-ofs'],

    'CSP':techcolors['csp'],
    'PV':techcolors['upv'],
    'PVB':techcolors['pvb'],

    # 'Storage':plt.cm.tab20(12),
    # 'Storage|charge':plt.cm.tab20(12),
    # 'Storage|discharge':plt.cm.tab20(12),
    'Storage':techcolors['battery_8'],
    'Storage|charge':techcolors['battery_8'],
    'Storage|discharge':techcolors['battery_8'],
}
aggtechs_disagg = aggstack.copy()
for k,v in aggstack.items():
    if v == 'Storage':
        aggtechs_disagg[k+'|charge'] = 'Storage|charge'
        aggtechs_disagg[k+'|discharge'] = 'Storage|discharge'


if len(cases) <= 4:
    plt.close()
    f,ax = plt.subplots(figsize=(width, 5))

    ### Final capacity and generation
    datum = 'Capacity'
    data = {
        'data': dictin_cap,
        'values':'Capacity (GW)',
        'label':f' {lastyear} Capacity [GW]',
    }
    ax.set_ylabel(data['label'])
    dfplot = pd.concat(
        {case:
            data['data'][case].loc[data['data'][case].year==lastyear]
            .set_index('tech')[data['values']]
            for case in cases},
        axis=1,
    ).T
    dfplot = dfplot.rename(columns=aggstack).groupby(axis=1, level='tech').sum()
    unmapped = [c for c in dfplot if c not in aggcolors]
    if len(unmapped):
        raise Exception(f"Unmapped techs: {unmapped}")
    dfplot = (
        dfplot[[c for c in aggcolors if c in dfplot]]
        .round(3).replace(0,np.nan).dropna(axis=1, how='all').fillna(0)
    )
    mindistance = dfplot.sum(axis=1).max() / 20
    dfcumsum = dfplot.cumsum(axis=1)
    dfdiff = dfplot - dfplot.loc[dfplot.index.map(basemap)].values

    ## Absolute and difference
    plots.stackbar(df=dfplot, ax=ax, colors=aggcolors, width=barwidth, net=False)

    ### Labels
    for x, case in enumerate(cases):
        labels = (dfcumsum.loc[case] - dfplot.loc[case]/2).rename('middle').to_frame()
        try:
            labels['ylabel'] = plots.optimize_label_positions(
                ydata=labels.middle.values, mindistance=mindistance, ypad=0,
            )
        except ValueError:
            labels['ylabel'] = labels['middle'].values
        labels['yval'] = labels.index.map(dfplot.loc[case])
        for i, row in labels.iterrows():
            ## Draw the line
            ax.annotate(
                '',
                xy=(x+barwidth/2, row.middle),
                xytext=(x+barwidth/2+labelpad, row.ylabel),
                arrowprops={
                    'arrowstyle':'-', 'shrinkA':0, 'shrinkB':0,
                    'color':aggcolors[i], 'lw':0.5},
                annotation_clip=False,
            )
            ## Write the label
            diff = np.around(dfdiff.loc[case,i], 0)
            ax.annotate(
                # f"{row.yval:.0f} {i}" + (f" ({diff:+.0f})" if diff else ''),
                (
                    f"{row.yval:.0f}"
                    + (f" {i}" if case == list(cases.keys())[-1] else '')
                    + (f" ({diff:+.0f})" if diff else '')
                ),
                (x+barwidth/2+labelpad+0.01, row.ylabel),
                va='center', ha='left', fontsize=9, color=aggcolors[i],
                weight=('bold' if abs(diff) >= 100 else 'normal'),
                annotation_clip=False,
            )

    ### Formatting
    ax.set_xticks(range(len(cases)))
    ax.set_xticklabels(cases.keys(), rotation=45, rotation_mode='anchor', ha='right')
    ax.yaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(5))
    plt.tight_layout()
    plots.despine(ax)
    plt.draw()
    ### Save it
    slide = add_to_pptx('Capacity stacks', width=width)
    if interactive:
        plt.show()


#%%### Hodgepodge: Final capacity, final generation, final transmission, runtime
width = max(11, len(cases)*1.3)
plt.close()
f,ax = plt.subplots(
    2, 4, figsize=(width, SLIDE_HEIGHT), sharex=True,
    sharey=('col' if (sharey is True) else False),
)
handles = {}

### Final capacity and generation
toplot = {
    'Capacity': {
        'data': dictin_cap,
        'values':'Capacity (GW)',
        'label':f'{lastyear} Capacity [GW]'},
    'Generation': {
        'data': dictin_gen,
        'values':'Generation (TWh)',
        'label':f'{lastyear} Generation [TWh]'},
}
ax[0,1].axhline(0, c='k', ls='--', lw=0.75)
for col, (datum, data) in enumerate(toplot.items()):
    ax[0,col].set_ylabel(data['label'], y=-0.075)
    dfplot = pd.concat(
        {case:
         data['data'][case].loc[data['data'][case].year==lastyear]
         .set_index('tech')[data['values']]
         for case in cases},
        axis=1,
    ).T
    dfplot = (
        dfplot[[c for c in bokehcolors.index if c in dfplot]]
        .round(3).replace(0,np.nan).dropna(axis=1, how='all')
    )

    handles[datum] = plot_bars_abs_stacked(
        dfplot=dfplot, basecase=basemap,
        colors=techcolors, fontsize=8,
        ax=ax, col=col, net=(True if datum == 'Generation' else False),
        label=(False if lesslabels else True),
    )

### Total transmission
col = 2
ax[0,col].set_ylabel(f'{lastyear} Transmission capacity [TW-mi]', y=-0.075)
dftrans = pd.concat({
    case:
    dictin_trans[case].groupby(['year','trtype'])['Amount (GW-mi)'].sum()
    .unstack('trtype')
    .reindex(allyears).interpolate('linear')
    / 1e3
    for case in cases
}, axis=1).loc[lastyear].unstack('trtype')

handles['Transmission'] = plot_bars_abs_stacked(
    dfplot=dftrans, basecase=basemap,
    colors=colors_trans,
    ax=ax, col=col, net=False,
    label=(False if lesslabels else True),
)

### Runtime
col = 3
ax[0,col].set_ylabel('Runtime [hours]', y=-0.075)
dfplot = pd.concat(
    {case: dictin_runtime[case].groupby('process').processtime.sum() for case in cases},
    axis=1).T.fillna(0)
dfplot = dfplot[[c for c in colors_time.index if c in dfplot]].copy()

handles['Runtime'] = plot_bars_abs_stacked(
    dfplot=dfplot, basecase=basemap,
    colors=colors_time,
    ax=ax, col=col, net=False,
    label=(False if lesslabels else True),
)

### Formatting
for col in range(4):
    ax[1,col].set_xticks(range(len(cases)))
    ax[1,col].set_xticklabels(cases.keys(), rotation=90)
    ax[1,col].annotate('Diff', (0.03,0.03), xycoords='axes fraction', fontsize='large')
    ax[1,col].axhline(0, c='k', ls='--', lw=0.75)
plt.tight_layout()
plots.despine(ax)
plt.draw()
### Save it
slide = add_to_pptx('Capacity, generation, transmission, runtime', width=width)
if interactive:
    plt.show()

### Add legends as separate figure below the slide
plt.close()
f,ax = plt.subplots(1, 4, figsize=(11, 0.1))
for col, datum in enumerate(handles):
    leg = ax[col].legend(
        handles=handles[datum][::-1], loc='upper center', bbox_to_anchor=(0.5,1.0), 
        fontsize='medium', ncol=1, frameon=False,
        handletextpad=0.3, handlelength=0.7, columnspacing=0.5, 
    )
    ax[col].axis('off')
add_to_pptx(slide=slide, width=width, top=7.5)


#%% Costs: NPV of system cost, NPV of climate + health costs
simple_npv = False
width = max(11, len(cases)*1.3)
plt.close()
f,ax = plt.subplots(
    2, 3, figsize=(width, 6), sharex=True,
    sharey=('col' if (sharey is True) else False),
)
handles = {}

### NPV of system cost
col = 0
ax[0,col].set_ylabel('NPV of system cost [$B]', y=-0.075)
ax[0,col].axhline(0, c='k', ls='--', lw=0.75)
dfcost_npv = pd.concat(dictin_npv, axis=1).fillna(0).T
dfcost_npv = dfcost_npv[[c for c in bokehcostcolors.index if c in dfcost_npv]].copy()
if simple_npv:
    dfcost_npv = dfcost_npv.sum(axis=1)
    dfcost_npv = pd.concat([pd.Series({case:dfcost_npv[case]}, name=case).to_frame() for case in cases])

handles['NPV'] = plot_bars_abs_stacked(
    dfplot=dfcost_npv, basecase=basemap,
    colors=bokehcostcolors,
    ax=ax, col=col, net=(not simple_npv),
    label=(False if lesslabels else True),
)

### NPV of climate and health costs
col = 1
ax[0,col].set_ylabel('NPV of climate + health cost [$B]', y=-0.075)

dfsocial = {}
for case in cases:
    dfsocial[case] = (
        dictin_emissions[case].reindex(allyears).interpolate('linear')
        * scghg_central
    )[['CO2','CH4']].dropna() / 1e9
    dfsocial[case]['health'] = dictin_health_central[case].reindex(allyears).interpolate('linear')
dfsocial = pd.concat(dfsocial, axis=1)

dfsocial_npv = dfsocial.multiply(discounts, axis=0).dropna().sum().unstack('e')

handles['social'] = plot_bars_abs_stacked(
    dfplot=dfsocial_npv, basecase=basemap,
    colors=colors_social,
    ax=ax, col=col, net=True,
    label=(False if lesslabels else True),
)

### Combined
col = 2
ax[0,col].set_ylabel('NPV of system\n+ climate + health cost [$B]', y=-0.075)
dfcombo_npv = pd.concat([dfcost_npv, dfsocial_npv], axis=1)

handles['combo'] = plot_bars_abs_stacked(
    dfplot=dfcombo_npv, basecase=basemap,
    colors={**bokehcostcolors.to_dict(), **colors_social},
    ax=ax, col=col, net=True,
    label=(False if lesslabels else True),
)

### Formatting
for col in range(3):
    ax[1,col].set_xticks(range(len(cases)))
    ax[1,col].set_xticklabels(cases.keys(), rotation=90)
    ax[1,col].annotate('Diff', (0.03,0.03), xycoords='axes fraction', fontsize='large')
    ax[1,col].axhline(0, c='k', ls='--', lw=0.75)
    ## Add commas to y axis labels
    if max([abs(i) for i in ax[0,col].get_ylim()]) >= 10000:
        ax[0,col].yaxis.set_major_formatter(mpl.ticker.StrMethodFormatter('{x:,.0f}'))
plt.tight_layout()
plots.despine(ax)
plt.draw()
### Save it
slide = add_to_pptx('NPV of system, climate, health costs', width=width)
if interactive:
    plt.show()

### Add legends as separate figure below the slide
plt.close()
f,ax = plt.subplots(1, 4, figsize=(11, 0.1))
for col, datum in enumerate(handles):
    leg = ax[col].legend(
        handles=handles[datum][::-1], loc='upper center', bbox_to_anchor=(0.5,1.0), 
        fontsize='medium', ncol=1, frameon=False,
        handletextpad=0.3, handlelength=0.7, columnspacing=0.5, 
    )
for col in range(4):
    ax[col].axis('off')
add_to_pptx(slide=slide, width=width, top=7.5)


#%% Simplifed NPV
width = len(cases)*1.3 + 2
plt.close()
f,ax = plt.subplots(
    2, 3, figsize=(width, 6), sharex=True,
    sharey=('col' if (sharey is True) else False),
)
handles = {}

### NPV of system cost
col = 0
ax[0,col].set_ylabel('NPV of system cost [$B]', y=-0.075)
ax[0,col].axhline(0, c='k', ls='--', lw=0.75)
dfcost_npv = pd.concat(dictin_npv, axis=1).fillna(0).T.sum(axis=1)

def twobars(dfplot, basecase, colors, ax, col=0, ypad=0.02):
    if isinstance(basecase, str):
        dfdiff = dfplot - dfplot.loc[basecase]
    elif isinstance(basecase, list):
        dfdiff = dfplot - dfplot.loc[basecase].values
    elif isinstance(basecase, dict):
        dfdiff = dfplot - dfplot.loc[basecase.values()].values

    for (row, df) in enumerate([dfplot, dfdiff]):
        ax[row,col].bar(
            range(len(df)), df.values,
            color=[colors[c] for c in dfplot.index],
            width=0.8,
        )
        ymin, ymax = ax[row,col].get_ylim()
        _ypad = (ymax - ymin) * ypad
        if ymin < 0:
            ax[row,col].set_ylim(ymin * (1+ypad))
        ## label net value
        if not lesslabels:
            for x, case in enumerate(df.index):
                val = df.loc[case].sum()
                ax[row,col].annotate(
                    f'{val:.0f}', (x, val - _ypad),
                    ha='center', va='top', color='k', size=9,
                    path_effects=[pe.withStroke(linewidth=2.0, foreground='w', alpha=0.7)],
                )

twobars(dfplot=dfcost_npv, basecase=basemap, colors=colors, ax=ax, col=col)

### NPV of climate and health costs
col = 1
ax[0,col].set_ylabel('NPV of climate + health cost [$B]', y=-0.075)

dfsocial = {}
for case in cases:
    dfsocial[case] = (
        dictin_emissions[case].reindex(allyears).interpolate('linear')
        * scghg_central
    )[['CO2','CH4']].dropna() / 1e9
    dfsocial[case]['health'] = dictin_health_central[case].reindex(allyears).interpolate('linear')
dfsocial = pd.concat(dfsocial, axis=1)

dfsocial_npv = dfsocial.multiply(discounts, axis=0).dropna().sum().unstack('e').sum(axis=1)

twobars(dfplot=dfsocial_npv, basecase=basemap, colors=colors, ax=ax, col=col)

### Combined
col = 2
ax[0,col].set_ylabel('NPV of system\n+ climate + health cost [$B]', y=-0.075)
dfcombo_npv = pd.concat([dfcost_npv, dfsocial_npv], axis=1).sum(axis=1)

twobars(dfplot=dfcombo_npv, basecase=basemap, colors=colors, ax=ax, col=col)

### Formatting
for col in range(3):
    ax[1,col].set_xticks(range(len(cases)))
    ax[1,col].set_xticklabels(cases.keys(), rotation=90)
    ax[1,col].annotate('Diff', (0.03,0.03), xycoords='axes fraction', fontsize='large')
    ax[1,col].axhline(0, c='k', ls='--', lw=0.75)
    ## Add commas to y axis labels
    if max([abs(i) for i in ax[0,col].get_ylim()]) >= 10000:
        ax[0,col].yaxis.set_major_formatter(mpl.ticker.StrMethodFormatter('{x:,.0f}'))
plt.tight_layout()
plots.despine(ax)

### Save it
slide = add_to_pptx('NPV of system, climate, health costs', width=min(width, SLIDE_WIDTH))
if interactive:
    plt.show()


#%%### SCOE, NEUE
width = 9 + len(cases)*0.5
plt.close()
f,ax = plt.subplots(
    1, 4, figsize=(width, 4.5),
    gridspec_kw={'wspace':0.7, 'width_ratios':[1,1,1,len(cases)*0.25]},
)

### SCOE
col = 0
dfscoe = {}
for case in cases:
    dfscoe[case] = dictin_scoe[case].groupby('year')['Average cost ($/MWh)'].sum().loc[years[case]]
    ax[col].plot(dfscoe[case].index, dfscoe[case].values, label=case, color=colors[case])
ax[col].set_ylim(0)
ax[col].set_ylabel('System cost of electricity [$/MWh]')
dfscoe = pd.concat(dfscoe, axis=1)
## annotate the last value
plots.label_last(dfscoe, ax[col], colors=colors, extend='below')

### Undiscounted annualized system cost
col = 1
dfsyscost = {}
for case in cases:
    dfsyscost[case] = dictin_syscost[case].sum(axis=1).loc[startyear:lastyear]
    ax[col].plot(dfsyscost[case].index, dfsyscost[case].values, label=case, color=colors[case])
ax[col].set_ylim(0)
ax[col].set_ylabel('Annualized system cost [$B/year]')
dfsyscost = pd.concat(dfsyscost, axis=1)
## annotate the last value
plots.label_last(dfsyscost, ax[col], colors=colors, extend='below')

### NEUE
col = 2
dfneue = {}
for case in cases:
    if case in dictin_neue:
        dfneue[case] = dictin_neue[case].reindex([y for y in years[case] if y >= 2025])
        ax[col].plot(dfneue[case].index, dfneue[case].values, label=case, color=colors[case])
ax[col].set_ylim(0)
if ax[col].get_ylim()[1] >= 10:
    ax[col].axhline(10, c='C7', ls='--', lw=0.75)
ax[col].set_ylabel('NEUE [ppm]')
## annotate the last value
if len(dfneue):
    dfneue = pd.concat(dfneue, axis=1)
    plots.label_last(dfneue, ax[col], colors=colors, extend='below')

### Spares
col = 3
ypad = 0.02
ax[col].set_ylabel('NPV of system cost [$billion]')
# ax[3].axis('off')
ax[col].bar(
    range(len(dfcost_npv)), dfcost_npv.values,
    color=[colors[c] for c in dfplot.index],
    width=0.8,
)
ymin, ymax = ax[col].get_ylim()
_ypad = (ymax - ymin) * ypad
if ymin < 0:
    ax[col].set_ylim(ymin * (1+ypad))
## label net value
for x, case in enumerate(dfcost_npv.index):
    val = dfcost_npv.loc[case].sum()
    if len(cases) <= 10:
        ax[col].annotate(
            f'{np.around(val,-1):.0f}',
            (x, val - _ypad), ha='center', va='top', color='k', size=9,
            path_effects=[pe.withStroke(linewidth=2.0, foreground='w', alpha=0.8)],
        )
    else:
        ax[col].annotate(
            f'{np.around(val,-1):.0f}',
            (x, val + _ypad), ha='center', va='bottom', color='k', size=9,
            rotation=90,
        )

ax[col].set_xticks(range(len(cases)))
ax[col].set_xticklabels(cases, rotation=45, rotation_mode='anchor', ha='right')

### Legend
leg = ax[0].legend(
    loc='upper left', bbox_to_anchor=(-0.3,-0.05), frameon=False, fontsize='large',
    handletextpad=0.3, handlelength=0.7,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')

### Formatting
# plt.tight_layout()
plots.despine(ax)
plt.draw()
for col in [0,1] + ([2] if len(dictin_neue) else []):
    ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
    ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    plots.shorten_years(ax[col])
### Save it
slide = add_to_pptx('Cost, reliability', width=width)
if interactive:
    plt.show()


#%% Emissions, health
width = 12
plt.close()
f,ax = plt.subplots(1, 4, figsize=(width, 4.5), gridspec_kw={'wspace':0.6})

### CO2 emissions
for col, pollutant in enumerate(['CO2','CO2e']):
    note = []
    df = {}
    for case in cases:
        df[case] = dictin_emissions[case].reindex(years[case]).fillna(0) / 1e6 # metric ton->MMT
        emissions_allyears = df[case].reindex(allyears).interpolate('linear').loc[startyear_notes:]
        ax[col].plot(df[case].index, df[case][pollutant].values, label=case, color=colors[case])
        ## collect more notes
        lives = dictin_health_central_mort[case].reindex(allyears).interpolate('linear').sum()
        note.append(
            f"{case:<{maxlength}} | {startyear_notes}–{lastyear}: "
            + f"{emissions_allyears['CO2'].sum()/1e3:.2f} GT CO2"
            + f"; {emissions_allyears['CO2e'].sum()/1e3:.2f} GT CO2e"
            + f"; {lives:,.0f} lives"
        )
    ax[col].set_ylim(0)
    ax[col].set_ylabel(f"{pollutant.replace('CO2e','CO2(e)')} emissions [MMT/year]")
    ## annotate the last value
    df = pd.concat(df, axis=1).xs(pollutant, 1, 'e')
    plots.label_last(df, ax[col], colors=colors, extend='both')

## Notes
ax[0].axhline(phaseout_trigger, c='C7', ls='--', lw=0.75)
ax[1].annotate(
    '\n'.join(note), (-0.2, -0.1), xycoords='axes fraction',
    annotation_clip=False, fontsize=9, fontfamily='monospace', va='top', ha='left',
)
print('\n'.join(note))

### Health impacts - mortality
col = 2
dfmort = {}
for case in cases:
    if case in dictin_health_central_mort:
        dfmort[case] = dictin_health_central_mort[case].loc[years[case]]
        ax[col].plot(dfmort[case].index, dfmort[case].values, label=case, color=colors[case])
ax[col].set_ylim(0)
ax[col].set_ylabel('Mortality [lives/year]')
## annotate the last value
dfmort = pd.concat(dfmort, axis=1)
plots.label_last(dfmort, ax[col], colors=colors, extend='both')

### Health impacts - dollars
col = 3
dfhealth = {}
for case in cases:
    if case in dictin_health_central:
        dfhealth[case] = dictin_health_central[case].loc[years[case]]
        ax[col].plot(dfhealth[case].index, dfhealth[case].values, label=case, color=colors[case])
ax[col].set_ylim(0)
ax[col].set_ylabel('Health costs [$B/year]')
## annotate the last value
dfhealth = pd.concat(dfhealth, axis=1)
plots.label_last(dfhealth, ax[col], colors=colors, extend='both')

### Legend
leg = ax[0].legend(
    loc='upper left', bbox_to_anchor=(-0.3,-0.05), frameon=False, fontsize='large',
    handletextpad=0.3, handlelength=0.7,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')

### Formatting
# plt.tight_layout()
plots.despine(ax)
plt.draw()
for col in range(4):
    ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
    ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    plots.shorten_years(ax[col])
### Save it
slide = add_to_pptx('Emissions', width=width)
if interactive:
    plt.show()


#%%### Generation fraction
ycol = 'Generation (TWh)'
stortechs = [f'battery_{i}' for i in [2,4,6,8,10]] + ['pumped-hydro']
vretechs = ['upv','wind-ons','wind-ofs','distpv','csp']
retechs = vretechs + ['hydro','geothermal','biopower']
zctechs = vretechs + ['hydro','geothermal','nuclear','nuclear-smr']
fossiltechs = ['coal','coal-ccs_mod','gas-cc','gas-cc-ccs_mod','gas-ct','o-g-s','cofire']
reccsnuctechs = retechs + ['coal-ccs_mod','gas-cc-ccs_mod','nuclear','nuclear-smr']

dftotal = pd.concat({
    case:
    dictin_gen[case].loc[~dictin_gen[case].tech.isin(stortechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dfvre = pd.concat({
    case:
    dictin_gen[case].loc[dictin_gen[case].tech.isin(vretechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dfre = pd.concat({
    case:
    dictin_gen[case].loc[dictin_gen[case].tech.isin(retechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dfzc = pd.concat({
    case:
    dictin_gen[case].loc[dictin_gen[case].tech.isin(zctechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dffossil = pd.concat({
    case:
    dictin_gen[case].loc[dictin_gen[case].tech.isin(fossiltechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dfreccsnuc = pd.concat({
    case:
    dictin_gen[case].loc[dictin_gen[case].tech.isin(reccsnuctechs)].groupby('year')[ycol].sum()
    for case in cases
}, axis=1)

dfplot = {
    'VRE share [%]': (dfvre / dftotal * 100).loc[startyear:],
    'RE share [%]': (dfre / dftotal * 100).loc[startyear:],
    'Zero carbon share [%]': (dfzc / dftotal * 100).loc[startyear:],
    'RE + Nuclear + CCS share [%]': (dfreccsnuc / dftotal * 100).loc[startyear:],
    'Fossil share [%]': (dffossil / dftotal * 100).loc[startyear:],
}

### Plot them
plt.close()
f,ax = plt.subplots(1, 5, figsize=(SLIDE_WIDTH, 4.5), gridspec_kw={'wspace':0.7})
for col, (ylabel, df) in enumerate(dfplot.items()):
    ax[col].set_ylabel(ylabel, labelpad=-4)
    ax[col].set_ylim(0,100)
    ax[col].yaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    for case in cases:
        ax[col].plot(df.index, df[case].values, label=case, color=colors[case])
    ## annotate the last value
    plots.label_last(df, ax[col], mindistance=3.5, colors=colors, extend='both', tail='%')

### Legend
leg = ax[0].legend(
    loc='upper left', bbox_to_anchor=(-0.3,-0.05), frameon=False, fontsize='large',
    handletextpad=0.3, handlelength=0.7,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')

### Formatting
plots.despine(ax)
plt.draw()
for col in range(len(dfplot)):
    ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
    ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    plots.shorten_years(ax[col])
### Save it
slide = add_to_pptx('Generation share')
if interactive:
    plt.show()


#%%### Firm capacity and capacity credit
capcreditcases = [c for c in cases if int(dictin_sw[c].GSw_PRM_CapCredit)]
if len(capcreditcases):
    capcredittechs = ['wind-ons','upv','storage']
    ccseasons = ['hot','cold']
    cccols = len(ccseasons) + len(capcredittechs)
    handles = {}

    plt.close()
    f,ax = plt.subplots(
        2, cccols, figsize=(11, SLIDE_HEIGHT), sharex='col',
        sharey=('col' if (sharey is True) else False),
    )
    ### Firm capacity stack
    dfplot = pd.concat(
        {case:
            dictin_cap_firm[case].loc[
                dictin_cap_firm[case].t==lastyear
            ].groupby(['ccseason','i']).MW.sum().rename('GW') / 1e3
            for case in capcreditcases},
        axis=1,
    ).T
    for col, ccseason in enumerate(ccseasons):
        ax[0,col].set_ylabel(f'Firm capacity, {ccseason} [GW]', y=-0.075)

        df = dfplot[ccseason][[c for c in bokehcolors.index if c in dfplot[ccseason]]].copy()

        handles[ccseason] = plot_bars_abs_stacked(
            dfplot=df, basecase=basecase,
            colors=techcolors,
            ax=ax, col=col, net=False, label=False,
        )

    ### Average capacity credit by technology
    for _col, tech in enumerate(capcredittechs):
        col = _col + len(ccseasons)
        if tech == 'storage':
            techs = [f'battery_{i}' for i in [2,4,6,8,10]] + ['pumped-hydro']
        else:
            techs = [tech]
        for case in capcreditcases:
            cap_firm = (
                dictin_cap_firm[case]
                .loc[dictin_cap_firm[case].i.isin(techs)]
                .groupby(['ccseason','t']).MW.sum().unstack('ccseason') / 1e3
            )
            cap_total = (
                dictin_cap[case]
                .loc[dictin_cap[case].tech.isin(techs)]
                .groupby('year')['Capacity (GW)'].sum().rename('GW').loc[2025:]
            )
            capcredit = cap_firm.divide(cap_total, axis=0).dropna()
            # for ccseason in lss:
            for row, ccseason in enumerate(['hot','cold']):
                    ax[row,col].plot(
                        capcredit.index, capcredit[ccseason].values,
                        color=colors[case], label=case,
                    )

    ### Legend
    leg = ax[1,len(ccseasons)].legend(
        loc='upper left', bbox_to_anchor=(-0.3,-0.07), frameon=False, fontsize='large',
        handletextpad=0.3, handlelength=0.7,
    )
    for legobj in leg.legend_handles:
        legobj.set_linewidth(8)
        legobj.set_solid_capstyle('butt')

    ### Formatting
    for col in range(2):
        ax[1,col].set_xticks(range(len(capcreditcases)))
        ax[1,col].set_xticklabels(capcreditcases, rotation=90)
        ax[1,col].annotate('Diff', (0.03,0.03), xycoords='axes fraction', fontsize='large')
        ax[1,col].axhline(0, c='k', ls='--', lw=0.75)
    for _col, tech in enumerate(capcredittechs):
        col = _col + len(ccseasons)
        ax[0,col].set_ylabel(f'Capacity credit, {tech} [fraction]', y=-0.075)
        for row, ccseason in enumerate(['hot','cold']):
            ax[row,col].set_ylim(0,1)
            ax[row,col].annotate(
                ccseason.title(), (0.5, 1.0), xycoords='axes fraction',
                fontsize='large', weight='bold', va='top', ha='center'
            )

    plt.tight_layout()
    plots.despine(ax)
    plt.draw()
    for col in range(len(ccseasons),cccols):
        ax[1,col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
        ax[1,col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
        plots.shorten_years(ax[1,col])

    ### Save it
    slide = add_to_pptx('Firm capacity, capacity credit')
    if interactive:
        plt.show()


#%%### Transmission
startyear_transgrowth = min(int(scalars.firstyear_trans_longterm), max(years[basecase]))
for interzonal_only in [False, True]:
    if interzonal_only:
        labelline = 'Interzonal transmission [TW-mi]'
        labelbar = f'{lastyear} Interzonal transmission [TW-mi]'
        labelgrowth = f'Interzonal transmission growth,\n{startyear_transgrowth}–{lastyear} [TWmi/year]'
    else:
        labelline = 'Transmission capacity [TW-mi]'
        labelbar = f'{lastyear} Transmission capacity [TW-mi]'
        labelgrowth = f'Transmission growth,\n{startyear_transgrowth}–{lastyear} [TWmi/year]'

    plt.close()
    f,ax = plt.subplots(1, 4, figsize=(11, 4.5), gridspec_kw={'wspace':0.6})

    ### Transmission TW-miles over time
    for case in cases:
        if interzonal_only:
            df = (
                dictin_trans[case]
                .loc[~dictin_trans[case].trtype.str.lower().isin(['spur','reinforcement'])]
                .groupby('year')['Amount (GW-mi)'].sum().reindex(years[case]) / 1e3
            )
        else:
            df = dictin_trans[case].groupby('year').sum()['Amount (GW-mi)'].reindex(years[case]) / 1e3
        ax[0].plot(df.index, df.values, label=case, color=colors[case])
        ## annotate the last value
        val = np.around(df.loc[max(years[case])], 0) + 0
        ax[0].annotate(
            f' {val:.0f}',
            (max(years[case]), df.loc[max(years[case])]), ha='left', va='center',
            color=colors[case], fontsize='medium',
        )
    ax[0].set_ylim(0)
    ax[0].set_ylabel(labelline)

    ### Disaggregated transmission (for next two plots)
    dftrans = pd.concat({
        case:
        dictin_trans[case].groupby(['year','trtype'])['Amount (GW-mi)'].sum()
        .unstack('trtype')
        .reindex(allyears).interpolate('linear')
        / 1e3
        for case in cases
    }, axis=1)
    if interzonal_only:
        dftrans = (
            dftrans[[c for c in dftrans if c[1].lower() not in ['spur','reinforcement']]]
        ).copy()

    ### Disaggregated final year transmission capacity
    df = dftrans.loc[lastyear].unstack('trtype')
    plots.stackbar(df=df, ax=ax[1], colors=colors_trans, width=0.8, net=False)
    ax[1].set_ylabel(labelbar)

    ### Transmission growth
    dftransgrowth = (
        (dftrans.loc[lastyear] - dftrans.loc[startyear_transgrowth])
        / (lastyear - startyear_transgrowth)
    ).unstack('trtype')
    plots.stackbar(df=dftransgrowth, ax=ax[2], colors=colors_trans, width=0.8, net=False)
    ax[2].set_ylabel(labelgrowth)
    ax[2].set_ylim(0, max(ax[2].get_ylim()[1], 3.8))
    ## Scales
    ymax = ax[2].get_ylim()[1]
    scales = {
        ## https://cigreindia.org/CIGRE%20Lib/CIGRE%20Session%202010%20paper/B4_306_2010.pdf
        1476 * 6.3 / 1e3: '1× Rio Madeira per year',
    }
    ## Values here from DOE 2024 Land-based Wind Market Report (page 64)
    ## and represent U.S-wide transmission capacity editions
    ## https://emp.lbl.gov/sites/default/files/2024-08/Land-Based%20Wind%20Market%20Report_2024%20Edition.pdf
    if interzonal_only:
        scales[0.73] = 'Mean since 2014 (345+ kV)'
        scales[1.46] = 'Max since 2014 (345+ kV)'
        scales[3.42] = 'Max since 2009 (345+ kV)'
    if not interzonal_only:
        scales[0.96] = 'Mean since 2014 (all kV)'
        scales[1.83] = 'Max since 2014 (all kV)'
        scales[3.64] = 'Max since 2009 (all kV)'

    ## Only add labels if doing a national-scale run
    if sw.GSw_Region.lower() == 'usa':
        for y, label in scales.items():
            if y > ymax:
                continue
            ax[2].annotate(
                label, xy=(len(cases), y), xytext=(len(cases)+1, y), annotation_clip=False,
                arrowprops={'arrowstyle':'-|>', 'color':'k'},
                ha='left', va='center', color='k', fontsize=11,
            )
            ax[2].axhline(
                y, c='k', lw=0.5, ls='--',
                path_effects=[pe.withStroke(linewidth=1.5, foreground='w', alpha=0.5)])

    ### Spare
    ax[3].axis('off')

    ### Legends
    ## Traces
    _h, _l = ax[0].get_legend_handles_labels()
    leg = ax[0].legend(
        # _h[::-1], _l[::-1],
        loc='upper left', bbox_to_anchor=(-0.4,-0.05), frameon=False, fontsize='large',
        handletextpad=0.3, handlelength=0.7,
    )
    for legobj in leg.legend_handles:
        legobj.set_linewidth(8)
        legobj.set_solid_capstyle('butt')
    ## Transmission types
    handles = [
        mpl.patches.Patch(facecolor=colors_trans[i], edgecolor='none', label=i)
        for i in colors_trans.index if i in dftrans.columns.get_level_values('trtype')
    ]
    leg = ax[2].legend(
        handles=handles[::-1],
        loc='upper left', bbox_to_anchor=(1,-0.05), frameon=False, fontsize='large',
        handletextpad=0.3, handlelength=0.7,
    )

    ### Formatting
    for col in [1,2]:
        ax[col].set_xticks(range(len(cases)))
        ax[col].set_xticklabels(cases.keys(), rotation=90)
        ax[col].yaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(5))

    plots.despine(ax)
    plt.draw()
    for col in [0]:
        ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
        ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
        plots.shorten_years(ax[col])
    ### Save it
    slide = add_to_pptx(
        'Interzonal transmission' if interzonal_only else 'Transmission (all types)')
    if interactive:
        plt.show()


#%%### More transmission plot styles
ylabels = {
    'transgrp': 'Transmission capacity between\nplanning regions [GW]',
    'transreg': 'Interregional transmission\ncapacity [GW]',
    'interconnect': 'Interconnection-seam-crossing\ntransmission capacity [GW]',
}

plt.close()
f,ax = plt.subplots(1, 4, figsize=(11, 4.5), gridspec_kw={'wspace':0.9})

### All transmission [TW-mi]
### Interzonal ("long-distance") [TW-mi]
for col, interzonal_only in enumerate([False,True]):
    if interzonal_only:
        labelline = 'Interzonal transmission capacity\n[TW-mi]'
    else:
        labelline = 'Total transmission capacity\n[TW-mi]'

    df = {}
    for case in cases:
        if interzonal_only:
            df[case] = (
                dictin_trans[case]
                .loc[~dictin_trans[case].trtype.str.lower().isin(['spur','reinforcement'])]
                .groupby('year')['Amount (GW-mi)'].sum().reindex(years[case]) / 1e3
            )
        else:
            df[case] = dictin_trans[case].groupby('year').sum()['Amount (GW-mi)'].reindex(years[case]) / 1e3
        ax[col].plot(df[case].index, df[case].values, label=case, color=colors[case])
    ax[col].set_ylim(0)
    ax[col].set_ylabel(labelline)
    ## annotate the last value
    df = pd.concat(df, axis=1)
    plots.label_last(df, ax[col], colors=colors, extend='both')
    ## Annotate the first value
    plots.annotate(
        ax[col], list(cases.keys())[0], 2020, (10,-10),
        color='C7', arrowprops={'arrowstyle':'-|>','color':'C7'})
ax[1].set_ylim(0, ax[0].get_ylim()[1])

### Interregional (FERC regions)
### Interconnection-seam-crossing
for _col, level in enumerate(['transreg','interconnect']):
    col = _col + 2
    df = {}
    for case in cases:
        df[case] = dictin_trans_r[case].loc[
            dictin_trans_r[case][f'inter_{level}'] == 1
        ].groupby('t').MW.sum().reindex(years[case]).fillna(0) / 1e3
        ax[col].plot(df[case].index, df[case].values, label=case, color=colors[case])
    ax[col].set_ylim(0)
    ax[col].set_ylabel(ylabels[level])
    ## annotate the last value
    df = pd.concat(df, axis=1)
    plots.label_last(df, ax[col], colors=colors, extend='both')
    ## Annotate the first value
    plots.annotate(
        ax[col], list(cases.keys())[0], 2020, (10,(-10 if col == 2 else 10)),
        color='C7', arrowprops={'arrowstyle':'-|>','color':'C7'},
        decimals=(0 if df[case][2020] >= 10 else 1),
    )
ax[3].set_ylim(0, ax[2].get_ylim()[1])

### Legends
## Traces
_h, _l = ax[0].get_legend_handles_labels()
leg = ax[0].legend(
    _h[::-1], _l[::-1],
    # loc='lower left', bbox_to_anchor=(-0.02,-0.02),
    loc='upper left', bbox_to_anchor=(-0.4,-0.05),
    frameon=False, fontsize='large',
    handletextpad=0.3, handlelength=0.7,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')

### Formatting
plots.despine(ax)
plt.draw()
for col in [0,1,2,3]:
    ax[col].yaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
    ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    plots.shorten_years(ax[col])
### Save it
slide = add_to_pptx('Transmission at different resolutions')
if interactive:
    plt.show()


#%%### Interregional transfer capability to peak demand ratio
try:
    f, ax, dfplot = reedsplots.plot_interreg_transfer_cap_ratio(
        case=list(cases.values()),
        colors={v: colors[k] for k,v in cases.items()},
        casenames={v:k for k,v in cases.items()},
        level='transreg', tstart=startyear,
        ymax=None,
    )
    ### Save it
    slide = add_to_pptx(
        'Interregional transmission / peak demand',
        height=(SLIDE_HEIGHT if ax.shape[1] <= 8 else None),
        width=(SLIDE_WIDTH if ax.shape[1] > 8 else None),
    )
    if interactive:
        plt.show()

except Exception:
    print(traceback.format_exc())


#%%### Max firm imports
try:
    f, ax, dfplot = reedsplots.plot_max_imports(
        case=list(cases.values()),
        colors={v: colors[k] for k,v in cases.items()},
        casenames={v:k for k,v in cases.items()},
        level='nercr', tstart=startyear,
    )
    ### Save it
    slide = add_to_pptx(
        'Max net stress imports / peak demand',
        height=(SLIDE_HEIGHT if ax.shape[1] <= 8 else None),
        width=(SLIDE_WIDTH if ax.shape[1] > 8 else None),
    )
    if interactive:
        plt.show()

except Exception:
    print(traceback.format_exc())


#%%### Transmission maps
if (len(cases) == 2) and (not forcemulti):
    plt.close()
    f,ax = reedsplots.plot_trans_diff(
        casebase=casebase,
        casecomp=casecomp,
        pcalabel=False,
        wscale=0.0004,
        subtract_baseyear=2020,
        yearlabel=True,
        year=lastyear,
        alpha=1, dpi=150,
        titleshorten=titleshorten,
    )
    add_to_pptx(f'Transmission ({lastyear})')
    if interactive:
        plt.show()
else:
    ### Absolute
    wscale = 0.0003
    alpha = 0.8
    for subtract_baseyear in [None, 2020]:
        plt.close()
        f,ax = plt.subplots(
            nrows, ncols, figsize=(SLIDE_WIDTH, SLIDE_HEIGHT),
            gridspec_kw={'wspace':0.0,'hspace':-0.1},
        )
        for case in cases:
            ### Plot it
            reedsplots.plot_trans_onecase(
                case=cases[case], pcalabel=False, wscale=wscale,
                yearlabel=False, year=lastyear, simpletypes=None,
                alpha=alpha, scalesize=8,
                f=f, ax=ax[coords[case]], title=False,
                subtract_baseyear=subtract_baseyear,
                thickborders='transreg', drawstates=False, drawzones=False, 
                label_line_capacity=10,
                scale=(True if case == basecase else False),
            )
            ax[coords[case]].set_title(case)
        ### Formatting
        title = (
            f'New interzonal transmission since {subtract_baseyear}' if subtract_baseyear
            else 'All interzonal transmission')
        for row in range(nrows):
            for col in range(ncols):
                if nrows == 1:
                    ax[col].axis('off')
                elif ncols == 1:
                    ax[row].axis('off')
                else:
                    ax[row,col].axis('off')
        ### Save it
        slide = add_to_pptx(title)
        if interactive:
            plt.show()


    ### Difference
    plt.close()
    f,ax = plt.subplots(
        nrows, ncols, figsize=(SLIDE_WIDTH, SLIDE_HEIGHT),
        gridspec_kw={'wspace':0.0,'hspace':-0.1},
    )
    for case in cases:
        ax[coords[case]].set_title(case)
        if case == basecase:
            ### Plot absolute
            reedsplots.plot_trans_onecase(
                case=cases[case], pcalabel=False, wscale=wscale,
                yearlabel=False, year=lastyear, simpletypes=None,
                alpha=alpha, scalesize=8,
                f=f, ax=ax[coords[case]], title=False,
                subtract_baseyear=subtract_baseyear,
                thickborders='transreg', drawstates=False, drawzones=False, 
                label_line_capacity=10,
                scale=(True if case == basecase else False),
            )
        else:
            ### Plot the difference
            reedsplots.plot_trans_diff(
                casebase=cases[basecase], casecomp=cases[case],
                pcalabel=False, wscale=wscale,
                yearlabel=False, year=lastyear, simpletypes=None,
                alpha=alpha,
                f=f, ax=ax[coords[case]],
                subtract_baseyear=subtract_baseyear,
                thickborders='transreg', drawstates=False, drawzones=False, 
                label_line_capacity=10,
                scale=False,
            )
    ### Formatting
    title = 'Interzonal transmission difference'
    for row in range(nrows):
        for col in range(ncols):
            if nrows == 1:
                ax[col].axis('off')
            elif ncols == 1:
                ax[row].axis('off')
            else:
                ax[row,col].axis('off')
    ### Save it
    slide = add_to_pptx(title)
    if interactive:
        plt.show()


#%%### RA sharing
if detailed:
    ralevel = 'transreg'
    scale = 10
    wscale = 7e3
    dfmap = reeds.io.get_dfmap(cases[basecase])
    regions = dfmap[ralevel].bounds.minx.sort_values().index

    ### Calculate aggregated load
    dictin_load_stress_agg = {}
    for case in cases:
        if int(dictin_sw[case].GSw_PRM_CapCredit):
            hcol = 'ccseason'
            df = dictin_peak_ccseason[case].copy()
        else:
            hcol = 'h'
            df = dictin_load_stress[case].copy()
        df = (
            df.assign(region=df.r.map(hierarchy[case][ralevel]))
            .groupby(['t','region',hcol]).GW.sum()
            .loc[lastyear].unstack('region')
        )
        dictin_load_stress_agg[case] = df.sum()
        # df['szn'] = df.index.map(lambda x: x.split('h')[0])
        # df = df.reset_index().set_index(['szn','h'])
        # dictin_load_stress_agg[case] = df.copy()

    ### Calculate aggregated stress period flows
    tran_flow_stress_agg = {}
    for case in cases:
        if int(dictin_sw[case].GSw_PRM_CapCredit):
            df = dictin_prmtrade[case].copy()
            hcol = 'ccseason'
        else:
            df = dictin_tran_flow_stress[case].copy()
            hcol = 'h'
        df['aggreg'] = df.r.map(hierarchy[case][ralevel])
        df['aggregg'] = df.rr.map(hierarchy[case][ralevel])
        df['interface'] = df.aggreg + '|' + df.aggregg

        df = (
            df
            .loc[df.aggreg != df.aggregg]
            .groupby(['t',hcol,'interface']).GW.sum().unstack('interface').fillna(0)
        )
        if df.empty:
            continue
        else:
            df = df.loc[lastyear].copy()
        ## Order interfaces alphabetically
        rename = {}
        for interface in df:
            r, rr = interface.split('|')
            if r > rr:
                rename[interface] = f'{rr}|{r}'
                df[interface] *= -1
        df = df.rename(columns=rename).groupby(axis=1, level=0).sum()
        ## Now reorder interfaces by flow
        rename = {}
        for interface in df:
            r, rr = interface.split('|')
            if df[interface].clip(lower=0).sum() < df[interface].clip(upper=0).abs().sum():
                rename[interface] = f'{rr}|{r}'
                df[interface] *= -1
        tran_flow_stress_agg[case] = df.rename(columns=rename).copy()

    ### Calculate regional imports/exports
    dfimportexport = {}
    for case in cases:
        df = {}
        for region in regions:
            df[region] = reedsplots.get_import_export(
                region=region, df=tran_flow_stress_agg[case]
            )
        dfimportexport[case] = pd.concat(df).sum(axis=1).unstack(level=0)
        dfimportexport[case].columns = dfimportexport[case].columns.rename('region')

    ### Plot it
    whiteout = dict(zip(
        [f'C{i}' for i in range(10)],
        [plt.cm.tab20(i*2+1) for i in range(10)]
    ))
    if any([v not in whiteout for v in list(colors.values())]):
        whiteout = {v: (v[0], v[1], v[2], v[3]*0.7) for v in list(colors.values())}

    for label in ['max','average']:
        plt.close()
        f,ax = plt.subplots(
            nrows, ncols, figsize=(SLIDE_WIDTH, SLIDE_HEIGHT),
            gridspec_kw={'wspace':0.0,'hspace':-0.1},
        )

        for case in cases:
            ### Formatting
            dfmap[ralevel].plot(ax=ax[coords[case]], facecolor='none', edgecolor='C7', lw=0.5)
            dfmap['interconnect'].plot(ax=ax[coords[case]], facecolor='none', edgecolor='k', lw=1)
            ax[coords[case]].set_title(
                case, y=0.95, weight='bold', color=colors[case], fontsize=14)
            ax[coords[case]].axis('off')

            ### RA flows
            if case not in tran_flow_stress_agg:
                continue

            if label == 'max':
                ## Max flow
                forwardwidth = tran_flow_stress_agg[case].clip(lower=0).max()
                reversewidth = abs(tran_flow_stress_agg[case].clip(upper=0).min())
            # ## GWh per day
            # wscale = 2.5e3
            # forwardwidth = gwh_forward / numdays
            # reversewidth = gwh_reverse / numdays
            else:
                ## Average when it's flowing
                forwardwidth = (
                    tran_flow_stress_agg[case].clip(lower=0).sum()
                    / tran_flow_stress_agg[case].clip(lower=0).astype(bool).sum()
                )
                reversewidth = (
                    tran_flow_stress_agg[case].clip(upper=0).abs().sum()
                    / tran_flow_stress_agg[case].clip(upper=0).abs().astype(bool).sum()
                )

            interfaces = tran_flow_stress_agg[case].columns
            numdays = (
                len(tran_flow_stress_agg[case])
                * int(dictin_sw[case].GSw_HourlyChunkLengthStress)
                // 24
            )

            ### Head/tail length:
            gwh_forward = tran_flow_stress_agg[case].clip(lower=0).sum()
            gwh_reverse = abs(tran_flow_stress_agg[case].clip(upper=0).sum())

            reversefrac = gwh_reverse / (gwh_reverse + gwh_forward)
            forwardfrac = gwh_forward / (gwh_reverse + gwh_forward)

            ### Plot it
            for interface in interfaces:
                r, rr = interface.split('|')
                startx, starty = dfmap[ralevel].loc[r, ['x', 'y']]
                endx, endy = dfmap[ralevel].loc[rr, ['x', 'y']]

                plots.plot_segmented_arrow(
                    ax[coords[case]],
                    reversefrac=reversefrac[interface],
                    forwardfrac=forwardfrac[interface],
                    reversewidth=reversewidth[interface]*wscale,
                    forwardwidth=forwardwidth[interface]*wscale,
                    startx=startx, starty=starty, endx=endx, endy=endy,
                    forwardcolor=colors[case], reversecolor=whiteout[colors[case]],
                    alpha=0.8, headwidthfrac=1.5,
                )
            ### Scale
            if scale:
                (startx, starty, endx, endy) = (-2.0e6, -1.2e6, -1.5e6, -1.2e6)
                yspan = ax[coords[case]].get_ylim()
                yspan = yspan[1] - yspan[0]
                plots.plot_segmented_arrow(
                    ax[coords[case]],
                    reversefrac=0, forwardfrac=1,
                    reversewidth=0, forwardwidth=scale*wscale,
                    startx=startx, starty=starty, endx=endx, endy=endy,
                    forwardcolor=colors[case], reversecolor=whiteout[colors[case]],
                    alpha=0.8, headwidthfrac=1.5,
                )
                ax[coords[case]].annotate(
                    f"{scale} GW\n{label}", ((startx+endx)/2, starty-(scale/2*wscale)-yspan*0.02),
                    ha='center', va='top', fontsize=14,
                )

        ### Save it
        title = f'{ralevel} {label} RA flows'
        slide = add_to_pptx(title)
        if interactive:
            plt.show()


#%%### Copy some premade single-case plots
level = dictin_sw[basecase]['GSw_PRM_StressThreshold'].split('_')[0]
wide = 1 if len(hierarchy[basecase]['transreg'].unique()) > 6 else 0
for figname, width, height in [
    (f'map_gencap_transcap-{lastyear}', None, SLIDE_HEIGHT),
    (f'plot_stressperiod_evolution-sum-{level}', SLIDE_WIDTH, None),
    (f'plot_dispatch-yearbymonth-1-{lastyear}', SLIDE_WIDTH, None),
    (
        f'plot_cap_rep_stress_mix-GW-transreg-{lastyear}',
        (SLIDE_WIDTH if wide else None),
        (None if wide else SLIDE_HEIGHT)
    ),
    (
        f'plot_cap_rep_stress_mix-percent-transreg-{lastyear}',
        (SLIDE_WIDTH if wide else None),
        (None if wide else SLIDE_HEIGHT)
    ),
]:
    for case in cases:
        try:
            slide = add_to_pptx(
                case,
                file=os.path.join(cases[case], 'outputs', 'maps', f'{figname}.png'),
                width=width, height=height,
            )
        except FileNotFoundError:
            print(f'No outputs/maps/{figname}.png for {os.path.basename(cases[case])}')


#%%### Generation capacity maps
### Shared data
base = cases[list(cases.keys())[0]]
val_r = dictin_cap_r[basecase].r.unique()
dfmap = reeds.io.get_dfmap(base)
dfba = dfmap['r']
dfstates = dfmap['st']
if (len(cases) == 2) and (not forcemulti):
    for i_plot in i_plots:
        plt.close()
        f,ax=plt.subplots(
            1, 3, sharex=True, sharey=True, figsize=(14,8),
            gridspec_kw={'wspace':-0.05, 'hspace':0.05},
            dpi=150,
        )

        _,_,dfplot = reedsplots.plotdiffmaps(
            val=mapdiff, i_plot=i_plot, year=lastyear, casebase=casebase, casecomp=casecomp,
            reeds_path=reeds_path, plot='base', f=f, ax=ax[0],
            cmap=cmocean.cm.rain,
        )
        ax[0].annotate(
            casebase_name,
            (0.1,1), xycoords='axes fraction', fontsize=10)

        _,_,dfplot = reedsplots.plotdiffmaps(
            val=mapdiff, i_plot=i_plot, year=lastyear, casebase=casebase, casecomp=casecomp,
            reeds_path=reeds_path, plot='comp', f=f, ax=ax[1],
            cmap=cmocean.cm.rain,
        )
        ax[1].annotate(
            casecomp_name,
            (0.1,1), xycoords='axes fraction', fontsize=10)

        _,_,dfplot = reedsplots.plotdiffmaps(
            val=mapdiff, i_plot=i_plot, year=lastyear, casebase=casebase, casecomp=casecomp,
            reeds_path=reeds_path, plot='absdiff', f=f, ax=ax[2],
            cmap=plt.cm.RdBu_r,
        )
        # print(dfplot.CAP_diff.min(), dfplot.CAP_diff.max())
        ax[2].annotate(
            '{}\n– {}'.format(
                casecomp_name,
                casebase_name),
            (0.1,1), xycoords='axes fraction', fontsize=10)

        add_to_pptx(f'{i_plot} capacity {lastyear} [GW]')
        if interactive:
            plt.show()
else:
    figwidth = SLIDE_WIDTH
    #### Absolute maps
    if (nrows == 1) or (ncols == 1):
        legendcoords = max(nrows, ncols) - 1
    elif (nrows-1, ncols-1) in coords.values():
        legendcoords = (nrows-1, ncols-1)
    else:
        legendcoords = (nrows-2, ncols-1)

    ### Set up plot
    for tech in maptechs:
        ### Get limits
        vmin = 0.
        vmax = float(pd.concat({
            case: dictin_cap_r[case].loc[
                (dictin_cap_r[case].i==tech)
                & (dictin_cap_r[case].t.astype(int)==lastyear)
            ].groupby('r').MW.sum()
            for case in cases
        }).max()) / 1e3
        if np.isnan(vmax):
            vmax = 0.
        if not vmax:
            print(f'{tech} has zero capacity in {lastyear}, so skipping maps')
            continue
        ### Set up plot
        plt.close()
        f,ax = plt.subplots(
            nrows, ncols, figsize=(figwidth, SLIDE_HEIGHT),
            gridspec_kw={'wspace':0.0,'hspace':-0.1},
        )
        ### Plot it
        for case in cases:
            dfval = dictin_cap_r[case].loc[
                (dictin_cap_r[case].i==tech)
                & (dictin_cap_r[case].t.astype(int)==lastyear)
            ].groupby('r').MW.sum()
            dfplot = dfba.copy()
            dfplot['GW'] = (dfval / 1e3).fillna(0)

            ax[coords[case]].set_title(
                case if nowrap else plots.wraptext(case, width=figwidth/ncols*0.9, fontsize=14)
            )
            dfba.plot(
                ax=ax[coords[case]],
                facecolor='none', edgecolor='k', lw=0.1, zorder=10000)
            dfstates.plot(
                ax=ax[coords[case]],
                facecolor='none', edgecolor='k', lw=0.2, zorder=10001)
            dfplot.plot(
                ax=ax[coords[case]], column='GW', cmap=cmap, vmin=vmin, vmax=vmax,
                legend=False,
            )
            ## Legend
            if coords[case] == legendcoords:
                plots.addcolorbarhist(
                    f=f, ax0=ax[coords[case]], data=dfplot.GW.values,
                    title=f'{tech} {lastyear}\ncapacity [GW]', cmap=cmap, vmin=vmin, vmax=vmax,
                    orientation='horizontal', labelpad=2.25, histratio=0.,
                    cbarwidth=0.05, cbarheight=0.85,
                    cbarbottom=-0.05, cbarhoffset=0.,
                )

        for row in range(nrows):
            for col in range(ncols):
                if nrows == 1:
                    ax[col].axis('off')
                elif ncols == 1:
                    ax[row].axis('off')
                else:
                    ax[row,col].axis('off')
        ### Save it
        slide = add_to_pptx(f'{tech} capacity {lastyear} [GW]')
        if interactive:
            plt.show()

    #### Difference maps
    ### Set up plot
    for tech in maptechs:
        ### Get limits
        dfval = pd.concat({
            case: dictin_cap_r[case].loc[
                (dictin_cap_r[case].i==tech)
                & (dictin_cap_r[case].t.astype(int)==lastyear)
            ].groupby('r').MW.sum()
            for case in cases
        }, axis=1).fillna(0) / 1e3
        dfdiff = dfval.subtract(dfval[basecase], axis=0)
        ### Get colorbar limits
        absmax = dfval.stack().max()
        diffmax = dfdiff.unstack().abs().max()

        if np.isnan(absmax):
            absmax = 0.
        if not absmax:
            print(f'{tech} has zero capacity in {lastyear}, so skipping maps')
            continue
        ### Set up plot
        plt.close()
        f,ax = plt.subplots(
            nrows, ncols, figsize=(figwidth, SLIDE_HEIGHT),
            gridspec_kw={'wspace':0.0,'hspace':-0.1},
        )
        ### Plot it
        for case in cases:
            dfplot = dfba.copy()
            dfplot['GW'] = dfval[case] if case == basecase else dfdiff[case]

            ax[coords[case]].set_title(
                case if nowrap else plots.wraptext(case, width=figwidth/ncols*0.9, fontsize=14)
            )
            dfba.plot(
                ax=ax[coords[case]],
                facecolor='none', edgecolor='k', lw=0.1, zorder=10000)
            dfstates.plot(
                ax=ax[coords[case]],
                facecolor='none', edgecolor='k', lw=0.2, zorder=10001)
            dfplot.plot(
                ax=ax[coords[case]], column='GW',
                cmap=(cmap if case == basecase else cmap_diff),
                vmin=(0 if case == basecase else -diffmax),
                vmax=(absmax if case == basecase else diffmax),
                legend=False,
            )
            ## Difference legend
            if coords[case] == legendcoords:
                plots.addcolorbarhist(
                    f=f, ax0=ax[coords[case]], data=dfplot.GW.values,
                    title=f'{tech} {lastyear}\ncapacity, difference\nfrom {basecase} [GW]',
                    cmap=(cmap if case == basecase else cmap_diff),
                    vmin=(0 if case == basecase else -diffmax),
                    vmax=(absmax if case == basecase else diffmax),
                    orientation='horizontal', labelpad=2.25, histratio=0.,
                    cbarwidth=0.05, cbarheight=0.85,
                    cbarbottom=-0.05, cbarhoffset=0.,
                )
        ## Absolute legend
        plots.addcolorbarhist(
            f=f, ax0=ax[coords[basecase]], data=dfval[basecase].values,
            title=f'{tech} {lastyear}\ncapacity [GW]',
            cmap=cmap, vmin=0, vmax=absmax,
            orientation='horizontal', labelpad=2.25, histratio=0.,
            cbarwidth=0.05, cbarheight=0.85,
            cbarbottom=-0.05, cbarhoffset=0.,
        )

        for row in range(nrows):
            for col in range(ncols):
                if nrows == 1:
                    ax[col].axis('off')
                elif ncols == 1:
                    ax[row].axis('off')
                else:
                    ax[row,col].axis('off')
        ### Save it
        slide = add_to_pptx(f'Difference: {tech} capacity {lastyear} [GW]')
        if interactive:
            plt.show()


#%% Save the powerpoint file
prs.save(savename)
print(f'\ncompare_casegroup.py results saved to:\n{savename}')

### Open it
if sys.platform == 'darwin':
    sp.run(f"open '{savename}'", shell=True)
elif platform.system() == 'Windows':
    sp.run(f'"{savename}"', shell=True)
