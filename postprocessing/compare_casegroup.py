#%% Imports
import numpy as np
import pandas as pd
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import patheffects as pe
import os
import sys
import io
import argparse
import site
import subprocess as sp
import platform
from glob import glob
from tqdm import tqdm
import traceback
import cmocean
import pptx
from pptx.util import Inches, Pt

reeds_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
remotepath = '/Volumes/ReEDS/' if sys.platform == 'darwin' else r'//nrelnas01/ReEDS/'

### Format plots and load other convenience functions
site.addsitedir(os.path.join(reeds_path,'postprocessing'))
import plots
import reedsplots
plots.plotparams()

#%% Argument inputs
parser = argparse.ArgumentParser(
    description='Compare multiple ReEDS cases',
    formatter_class=argparse.ArgumentDefaultsHelpFormatter,
)
parser.add_argument(
    'caselist', type=str,
    help=('comma-delimited list of cases to plot, OR shared casename prefix, '
          'OR csv file of cases. The first case is treated as the base case '
          'unless a different one is provided via the --basecase/-b argument.'))
parser.add_argument(
    '--casenames', '-n', type=str, default='',
    help='comma-delimited list of shorter case names to use in plots')
parser.add_argument(
    '--titleshorten', '-s', type=str, default='',
    help='characters to cut from start of case name (only used if no casenames)')
parser.add_argument(
    '--startyear', '-t', type=int, default=2020,
    help='First year to show')
parser.add_argument(
    '--sharey', '-y', action='store_true',
    help='Use same y-axis scale for absolute and difference plots')
parser.add_argument(
    '--basecase', '-b', type=str, default='',
    help='Substring of case path to use as default (if empty, uses first case in list)')
parser.add_argument(
    '--level', '-l', type=str, default='transreg',
    choices=['interconnect','nercr','transreg','transgrp','st'],
    help='hierarchy level at which to plot regional results')
parser.add_argument(
    '--skipbp', '-p', action='store_true',
    help='flag to prevent bokehpivot report from being generated')
parser.add_argument(
    '--bpreport', '-r', type=str, default='standard_report_reduced',
    help='which bokehpivot report to generate')

args = parser.parse_args()
_caselist = args.caselist.split(',')
try:
    titleshorten = int(args.titleshorten)
except ValueError:
    titleshorten = len(args.titleshorten)
_basecase = args.basecase
startyear = args.startyear
sharey = True if args.sharey else 'row'
level = args.level
bpreport = args.bpreport
skipbp = args.skipbp
interactive = False

# #%% Inputs for testing
# startyear = 2020
# reeds_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# casepath = os.path.join(reeds_path,'runs')
# _caselist = [
#     os.path.join(casepath,'v20231215_casegroupM0_Pacific_lim'),
#     os.path.join(casepath,'v20231215_casegroupM0_Pacific_ref'),
#     os.path.join(casepath,'v20231215_casegroupM0_Pacific_open'),
# ]
# _basecase = _caselist[1]
# titleshorten = len('v20231215_casegroupM0_Pacific_')
# sharey = 'row'
# interactive = True


#%%### Fixed inputs
cmap = cmocean.cm.rain
cmap_diff = cmocean.cm.balance

techmap = {
    **{f'upv_{i}':'Utility PV' for i in range(20)},
    **{f'dupv_{i}':'Utility PV' for i in range(20)},
    **{f'wind-ons_{i}':'Land-based wind' for i in range(20)},
    **{f'wind-ofs_{i}':'Offshore wind' for i in range(20)},
    **dict(zip(['nuclear','nuclear-smr'], ['Nuclear']*20)),
    **dict(zip(
        ['gas-cc_re-cc','gas-ct_re-ct','re-cc','re-ct',
         'gas-cc_h2-ct','gas-ct_h2-ct','h2-cc','h2-ct',],
        ['H2 turbine']*20)),
    **{f'battery_{i}':'Battery/PSH' for i in range(20)}, **{'pumped-hydro':'Battery/PSH'},
    **dict(zip(
        ['coal-igcc', 'coaloldscr', 'coalolduns', 'gas-cc', 'gas-ct', 'coal-new', 'o-g-s',],
        ['Fossil']*20)),
    **dict(zip(
        ['gas-cc_gas-cc-ccs_mod','gas-cc_gas-cc-ccs_max','gas-cc-ccs_mod','gas-cc-ccs_max',
         'gas-cc_gas-cc-ccs_mod','coal-igcc_coal-ccs_mod','coal-new_coal-ccs_mod',
        'coaloldscr_coal-ccs_mod','coalolduns_coal-ccs_mod','cofirenew_coal-ccs_mod',
        'cofireold_coal-ccs_mod','gas-cc_gas-cc-ccs_max','coal-igcc_coal-ccs_max',
        'coal-new_coal-ccs_max','coaloldscr_coal-ccs_max','coalolduns_coal-ccs_max',
        'cofirenew_coal-ccs_max','cofireold_coal-ccs_max',],
        ['Fossil+CCS']*50)),
    **dict(zip(['dac','beccs_mod','beccs_max'],['CO2 removal']*20)),
}

maptechs = [
    'Utility PV',
    'Land-based wind',
    'Offshore wind',
    'Nuclear',
    'H2 turbine',
    'Battery/PSH',
    'Fossil+CCS',
]


#%%### Functions
def add_to_pptx(
        title=None, file=None, left=0, top=0.62, width=13.33, height=None,
        verbose=1, slide=None,
    ):
    """Add current matplotlib figure (or file if specified) to new powerpoint slide"""
    if not file:
        image = io.BytesIO()
        plt.savefig(image, format='png')
    else:
        image = file
        if not os.path.exists(image):
            raise FileNotFoundError(image)

    if slide is None:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.title.text = title
    slide.shapes.add_picture(
        image,
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    if verbose:
        print(title)
    return slide


def add_textbox(
        text, slide,
        left=0, top=7.2, width=13.33, height=0.3,
        fontsize=14,
    ):
    """Add a textbox to the specified slide"""
    textbox = slide.shapes.add_textbox(
        left=(None if left is None else Inches(left)),
        top=(None if top is None else Inches(top)),
        width=(None if width is None else Inches(width)),
        height=(None if height is None else Inches(height)),
    )
    p = textbox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    return slide


#%%### Procedure
#%% Parse arguments
use_table_casenames = False
if len(_caselist) == 1:
    ## If it's a .csv, read the cases to compare
    if _caselist[0].endswith('.csv'):
        dfcase = pd.read_csv(_caselist[0], header=None)
        ## First check it's a simple csv with one case per row
        if dfcase.shape[1] == 1:
            caselist = dfcase[0].tolist()
        ## Then check if it's a 2-column csv with [casepath,casename] header
        elif (
            (dfcase.shape[1] == 2)
            and (dfcase.loc[0,[0,1]].tolist() == ['casepath', 'casename'])
        ):
            dfcase.columns = ['casepath','casename']
            dfcase.drop(0, inplace=True)
            caselist = dfcase.casepath.tolist()
            use_table_casenames = True
        ## Otherwise assume it's a copy of a cases_{batchname}.csv file in a case folder
        ## This approach is less robust; the others are preferred.
        else:
            prefix_plus_tail = os.path.dirname(_caselist[0])
            tails = [i for i in dfcase.iloc[0] if i not in ['Default Value',np.nan]]
            prefix = prefix_plus_tail[:-len([i for i in tails if prefix_plus_tail.endswith(i)][0])]
            caselist = [prefix+i for i in tails]
    ## Otherwise look for all runs starting with the provided string
    else:
        caselist = sorted(glob(_caselist[0]+'*'))
        ## If no titleshorten is provided, use the provided prefix
        if not titleshorten:
            titleshorten = len(os.path.basename(args.caselist))
else:
    caselist = _caselist

## Remove cases that haven't finished yet
caselist = [
    i for i in caselist
    if os.path.isfile(os.path.join(i,'outputs','reeds-report','report.xlsx'))
]

if use_table_casenames:
    casenames = dfcase.casename.tolist()
else:
    casenames = (
        args.casenames.split(',') if len(args.casenames)
        else [os.path.basename(c)[titleshorten:] for c in caselist]
    )

if len(caselist) != len(casenames):
    err = (
        f"len(caselist) = {len(caselist)} but len(casenames) = {len(casenames)}\n\n"
        'caselist:\n' + '\n'.join(caselist) + '\n\n'
        'casenames:\n' + '\n'.join(casenames) + '\n'
    )
    raise ValueError(err)

cases = dict(zip(casenames, caselist))

if not len(_basecase):
    basecase = list(cases.keys())[0]
else:
    basepath = [c for c in cases.values() if c.endswith(_basecase)]
    if len(basepath) == 0:
        err = (
            f"Use a basecase that matches one case.\nbasecase={_basecase} matches none of:\n"
            + '\n'.join(basepath)
        )
        raise ValueError(err)
    elif len(basepath) > 1:
        err = (
            f"Use a basecase that only matches one case.\nbasecase={_basecase} matches:\n"
            + '\n'.join(basepath)
        )
        raise ValueError(err)
    else:
        basepath = basepath[0]
        ## basecase is the short name; basepath is the full path
        basecase = casenames[caselist.index(basepath)]
        ## Put it first in the list
        cases = {**{basecase:cases[basecase]}, **{k:v for k,v in cases.items() if k != basecase}}

for case, path in cases.items():
    print(f'{path} -> {case}' + (' (base)' if case == basecase else ''))

colors = plots.rainbowmapper(cases)

#%% Create output folder
outpath = os.path.join(cases[basecase], 'outputs', 'comparisons')
os.makedirs(outpath, exist_ok=True)
## clip name to max length and removing disallowed characters
maxlength = os.pathconf(os.sep,'PC_NAME_MAX')
savename = os.path.join(
    outpath,
    (f"results-{','.join(cases.keys())}"
     .replace('/','').replace(' ','')
     [:maxlength-len('.pptx')]) + '.pptx'
)
print(f'Saving results to {savename}')

#%% Create bokehpivot report as subprocess
if not skipbp:
    start_str = 'start ' if platform.system() == 'Windows' else ''
    bp_path = f'{reeds_path}/postprocessing/bokehpivot'
    bp_py_file = f'{bp_path}/reports/interface_report_model.py'
    report_path = f'{bp_path}/reports/templates/reeds2/{bpreport}.py'
    bp_outpath = f'{outpath}/{bpreport}-diff-multicase'
    add_diff = 'Yes'
    auto_open = 'Yes'
    bp_colors = pd.read_csv(f'{bp_path}/in/example_reeds_scenarios.csv')['color'].tolist()
    bp_colors = bp_colors*10 #Up to 200 scenarios
    bp_colors = bp_colors[:len(casenames)]
    df_scenarios = pd.DataFrame({'name':casenames, 'color':bp_colors, 'path':caselist})
    scenarios_path = f'{outpath}/scenarios.csv'
    df_scenarios.to_csv(scenarios_path, index=False)
    call_str = (
        f'{start_str}python "{bp_py_file}" "ReEDS 2.0" "{scenarios_path}" all ' +
        f'{add_diff} "{basecase}" "{report_path}" "html,excel" one "{bp_outpath}" {auto_open}'
    )
    sp.Popen(call_str, shell=True)

#%%### Load data
#%% Shared
## Determine if we're on a branch before or after county-level capability was merged
countyreeds = (
    True if os.path.isfile(os.path.join(reeds_path,'inputs','transmission','r_rr_adj_county.csv'))
    else False
)
if countyreeds:
    hierarchy = pd.read_csv(
        os.path.join(reeds_path,'inputs','hierarchy.csv')
    ).drop(['*county','county_name'], axis=1).drop_duplicates().rename(columns={'ba':'r'}).set_index('r')
else:
    hierarchy = pd.read_csv(
        os.path.join(reeds_path,'inputs','hierarchy.csv')
    ).rename(columns={'*r':'r'}).set_index('r')    
hierarchy = hierarchy.loc[hierarchy.country.str.lower()=='usa'].copy()

sw = pd.read_csv(
    os.path.join(cases[case],'inputs_case','switches.csv'),
    header=None, index_col=0).squeeze(1)

scalars = pd.read_csv(
    os.path.join(cases[case], 'inputs_case', 'scalars.csv'),
    header=None, usecols=[0,1], index_col=0).squeeze(1)
phaseout_trigger = float(scalars.co2_emissions_2022) * float(sw.GSw_TCPhaseout_trigger_f)

#%% Colors
bokehcostcolors = pd.read_csv(
    os.path.join(
        reeds_path,'postprocessing','bokehpivot','in','reeds2','cost_cat_style.csv'),
    index_col='order').squeeze(1)
bokehcostcolors = bokehcostcolors.loc[~bokehcostcolors.index.duplicated()]

colors_time = pd.read_csv(
    os.path.join(
        reeds_path,'postprocessing','bokehpivot','in','reeds2','process_style.csv'),
    index_col='order',
).squeeze(1)

bokehcolors = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_style.csv'),
    index_col='order').squeeze(1)

tech_map = pd.read_csv(
    os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_map.csv'),
    index_col='raw').squeeze(1)
    
bokehcolors = pd.concat([
    bokehcolors.loc['smr':'electrolyzer'],
    pd.Series('#D55E00', index=['dac'], name='color'),
    bokehcolors.loc[:'Canada'],
])

bokehcolors['canada'] = bokehcolors['Canada']

techcolors = {
    'gas-cc_gas-cc-ccs':bokehcolors['gas-cc-ccs_mod'],
    'cofire':bokehcolors['biopower'],
    'gas-cc':'#5E1688',
    'gas-cc-ccs':'#9467BD',
}
for i in bokehcolors.index:
    if i in techcolors:
        pass
    elif i in bokehcolors.index:
        techcolors[i] = bokehcolors[i]
    else:
        raise Exception(i)

techcolors = {i: techcolors[i] for i in bokehcolors.index}

#%% Parse excel report sheet names
val2sheet = reedsplots.get_report_sheetmap(cases[basecase])

#%% National capacity and generation
renametechs = {
    'h2-cc_upgrade':'h2-cc',
    'h2-ct_upgrade':'h2-ct',
    'gas-cc-ccs_mod_upgrade':'gas-cc-ccs_mod',
    'coal-ccs_mod_upgrade':'coal-ccs_mod',
}
dictin_cap = {}
for case in tqdm(cases, desc='national capacity'):
    dictin_cap[case] = pd.read_excel(
        os.path.join(cases[case],'outputs','reeds-report','report.xlsx'),
        sheet_name=val2sheet['Capacity (GW)'],
    ).drop('scenario',axis=1)
    ### Simplify techs
    dictin_cap[case].tech = dictin_cap[case].tech.map(lambda x: renametechs.get(x,x))
    dictin_cap[case] = dictin_cap[case].groupby(['tech','year'], as_index=False)['Capacity (GW)'].sum()
    dictin_cap[case] = dictin_cap[case].loc[
        ~dictin_cap[case].tech.isin(['electrolyzer','smr','smr-ccs'])].copy()

dictin_gen = {}
for case in tqdm(cases, desc='national generation'):
    dictin_gen[case] = pd.read_excel(
        os.path.join(cases[case],'outputs','reeds-report','report.xlsx'),
        sheet_name=val2sheet['Generation (TWh)'],
    ).drop('scenario',axis=1)
    ### Simplify techs
    dictin_gen[case].tech = dictin_gen[case].tech.map(lambda x: renametechs.get(x,x))
    dictin_gen[case] = dictin_gen[case].groupby(['tech','year'], as_index=False)['Generation (TWh)'].sum()

costcat_rename = {
    'CO2 Spurline':'CO2 T&S Capex',
    'CO2 Pipeline':'CO2 T&S Capex',
    'CO2 Storage':'CO2 T&S Capex',
    'CO2 Spurline FOM':'CO2 T&S O&M',
    'CO2 Pipeline FOM':'CO2 T&S O&M',
    'CO2 Incentive Payments':'CCS Incentives',
    'Capital': 'Gen & Stor Capex',
    'O&M': 'Gen & Stor O&M',
    'CO2 Network':'CO2 T&S Capex',
    'CO2 Incentives':'CCS Incentives',
    'CO2 FOM':'CO2 T&S O&M',
    'CO2 Capture':'CO2 T&S Capex',
    'H2 Fuel':'Fuel',
    'H2 VOM':'H2 Prod O&M',
}
dictin_npv = {}
for case in tqdm(cases, desc='NPV of system cost'):
    try:
        dictin_npv[case] = pd.read_excel(
            os.path.join(cases[case],'outputs','reeds-report/report.xlsx'),
            sheet_name=val2sheet['Present Value of System Cost'], engine='openpyxl',
        ).drop('scenario',axis=1).set_index('cost_cat')['Discounted Cost (Bil $)']
        dictin_npv[case].index = pd.Series(dictin_npv[case].index).replace(costcat_rename)
        dictin_npv[case] = dictin_npv[case].groupby(level=0, sort=False).sum()
    except FileNotFoundError:
        print(case)

dictin_scoe = {}
for case in tqdm(cases, desc='SCOE'):
    try:
        dictin_scoe[case] = pd.read_excel(
            os.path.join(cases[case],'outputs','reeds-report/report.xlsx'),
            sheet_name=val2sheet['National Average Electricity'], engine='openpyxl',
        ).drop('scenario',axis=1)
    except FileNotFoundError:
        print(case)
    dictin_scoe[case].cost_cat = dictin_scoe[case].cost_cat.replace(
        {**costcat_rename,**{'CO2 Incentives':'CCS Incentives'}})
    dictin_scoe[case] = (
        dictin_scoe[case].groupby(['cost_cat','year'], sort=False, as_index=False)
        ['Average cost ($/MWh)'].sum())

pollutant = 'CO2'
dictin_emissions = {}
for case in tqdm(cases, desc='national emissions'):
    try:
        dictin_emissions[case] = pd.read_csv(
            os.path.join(cases[case], 'outputs', 'emit_nat.csv'),
            header=0, names=['e','t','tonne'], index_col=['e','t'],
        ### Convert to MMT
        ).squeeze(1).loc[pollutant] / 1e6
    except FileNotFoundError:
        print(case)

dictin_trans = {}
for case in tqdm(cases, desc='national transmission'):
    try:
        dictin_trans[case] = pd.read_excel(
            os.path.join(cases[case],'outputs','reeds-report','report.xlsx'),
            sheet_name=val2sheet['Transmission (GW-mi)'], engine='openpyxl',
        ).drop('scenario',axis=1)
    except FileNotFoundError:
        print(case)

dictin_trans_r = {}
for case in tqdm(cases, desc='regional transmission'):
    dictin_trans_r[case] = pd.read_csv(
        os.path.join(cases[case],'outputs','tran_out.csv')
    ).rename(columns={'Value':'MW'})
    for _level in ['interconnect','transreg','transgrp','st']:
        dictin_trans_r[case][f'inter_{_level}'] = (
            dictin_trans_r[case].r.map(hierarchy[_level])
            != dictin_trans_r[case].rr.map(hierarchy[_level])
        ).astype(int)

dictin_cap_r = {}
for case in tqdm(cases, desc='regional capacity'):
    dictin_cap_r[case] = pd.read_csv(
        os.path.join(cases[case],'outputs','cap.csv'),
        names=['i','r','t','MW'], header=0,
    )
    ### Simplify techs
    dictin_cap_r[case].i = dictin_cap_r[case].i.map(lambda x: renametechs.get(x,x))
    dictin_cap_r[case].i = dictin_cap_r[case].i.str.lower().map(lambda x: techmap.get(x,x))
    dictin_cap_r[case] = dictin_cap_r[case].groupby(['i','r','t'], as_index=False).MW.sum()

dictin_runtime = {}
for case in tqdm(cases, desc='runtime'):
    try:
        dictin_runtime[case] = pd.read_excel(
            os.path.join(cases[case],'outputs','reeds-report','report.xlsx'),
            sheet_name=val2sheet['Runtime (hours)'], engine='openpyxl',
        ).set_index('process').processtime
    except FileNotFoundError:
        print(case)

### Model years
years = sorted(dictin_cap[case].year.astype(int).unique())
years = [y for y in years if y >= startyear]
yearstep = years[-1] - years[-2]
lastyear = max(years)



#%%### Plots ######
### Set up powerpoint file
prs = pptx.Presentation(os.path.join(reeds_path,'postprocessing','template.pptx'))
blank_slide_layout = prs.slide_layouts[3]


#%%### Generation capacity lines
aggtechsplot = {
    'Interregional\ntransmission': 'inter_transreg',
    'Land-based\nwind': ['wind-ons'],
    'Offshore\nwind': ['wind-ofs'],
    # 'Wind': ['wind-ons', 'wind-ofs'],
    'Solar': ['upv', 'dupv', 'distpv', 'csp', 'pvb'],
    'Battery': ['battery_{}'.format(i) for i in [2,4,6,8,10]],
    'Pumped\nstorage\nhydro': ['pumped-hydro'],
    # 'Storage': ['battery_{}'.format(i) for i in [2,4,6,8,10]] + ['pumped-hydro'],
    'Hydro, geo, bio': [
        'hydro','geothermal',
        'biopower','lfill-gas','cofire','beccs_mod','beccs'
    ],
    'Nuclear': ['nuclear', 'nuclear-smr'],
    'Hydrogen\nturbine': ['h2-cc', 'h2-cc-upgrade', 'h2-ct', 'h2-ct-upgrade'],
    'Gas CCS': ['gas-cc-ccs_mod'],
    'Coal CCS': ['coal-ccs_mod'],
    # 'Fossil\n(with CCS)': ['gas-cc-ccs_mod','coal-ccs_mod'],
    'Fossil\n(w/o CCS)': ['gas-cc', 'gas-ct', 'o-g-s', 'coal', 'cofire'],
#     'CDR': ['dac', 'beccs'],
#     'H2 production': ['smr', 'smr-ccs', 'electrolyzer'],
}
checktechs = [i for sublist in aggtechsplot.values() for i in sublist]
alltechs = pd.concat(dictin_cap).tech.unique()
printstring = (
    'The following techs are not plotted: '
    + ', '.join([c for c in alltechs if c not in checktechs])
)

offsetstart = {
    'Solar': (15,0),
    'Wind': (15,0), 'Land-based\nwind': (15,0),
}

val = '4_Capacity (GW)'
ycol = 'Capacity (GW)'
nrows, ncols = 2, len(aggtechsplot)//2+len(aggtechsplot)%2
coords = dict(zip(
    list(aggtechsplot.keys()),
    [(row,col) for row in range(nrows) for col in range(ncols)]
))

offset = dict()

plt.close()
f,ax = plt.subplots(
    nrows, ncols, sharex=True, sharey=True,
    figsize=(13.33, 6.88),
    gridspec_kw={'wspace':0.3, 'hspace':0.15},
)
for tech in aggtechsplot:
    for case in cases:
        ### Central cases
        if 'transmission' in tech.lower():
            df = dictin_trans_r[case].loc[
                dictin_trans_r[case][aggtechsplot[tech]]==1
            ].groupby('t').MW.sum() / 1e3
        else:
            df = dictin_cap[case].loc[
                dictin_cap[case].tech.isin(aggtechsplot[tech])
            ].groupby('year')[ycol].sum().reindex(years).fillna(0)
        ax[coords[tech]].plot(
            df.index, df.values,
            label=case, color=colors[case], ls='-',
        )
        ### Annotate the last value
        fincap = df.reindex([lastyear]).fillna(0).squeeze()
        ax[coords[tech]].annotate(
            ' {:.0f}'.format(fincap),
            (lastyear, fincap+offset.get((tech,case),0)),
            ha='left', va='center',
            color=colors[case], fontsize='small',
            annotation_clip=False,
        )

    ### Formatting
    ax[coords[tech]].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5 if lastyear>2040 else 1))
    ax[coords[tech]].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10 if lastyear>2040 else 5))
    ax[coords[tech]].annotate(
        tech,
        (0.05,1.0), va='top', ha='left',
        xycoords='axes fraction', 
        fontsize='x-large', weight='bold',)
    ### Annotate the 2020 value
    plots.annotate(
        ax[coords[tech]], basecase,
        startyear, offsetstart.get(tech,(10,10)), color='C7',
        arrowprops={'arrowstyle':'-|>', 'color':'C7'})
if len(aggtechsplot) % 2:
    ax[-1,-1].axis('off')
handles, labels = ax[-1,0].get_legend_handles_labels()
leg = ax[-1,0].legend(
    handles, labels,
    fontsize='large', frameon=False, 
    loc='upper left', bbox_to_anchor=(0,0.95),
    handletextpad=0.3, handlelength=0.7,
    ncol=1,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')
ax[coords[list(aggtechsplot.keys())[0]]].set_xlim(startyear,lastyear)
ax[coords[list(aggtechsplot.keys())[0]]].set_ylim(0)
ax[coords[list(aggtechsplot.keys())[0]]].set_ylabel('Capacity [GW]', y=-0.075)
# for row in range(nrows):
#     ax[row,0].set_ylabel('Capacity [GW]')
ax[coords[list(aggtechsplot.keys())[0]]].set_ylim(0)

plots.despine(ax)
plt.draw()
plots.shorten_years(ax[1,0])
### Save it
slide = add_to_pptx('Capacity')
add_textbox(printstring, slide)
if interactive:
    print(printstring)
    plt.show()


#%%### Capacity and generation bars
toplot = {
    'Capacity': {'data': dictin_cap, 'values':'Capacity (GW)', 'label':'Capacity [GW]'},
    'Generation': {'data': dictin_gen, 'values':'Generation (TWh)', 'label':'Generation [TWh]'},
}
dfbase = {}
for slidetitle, data in toplot.items():
    plt.close()
    f,ax = plt.subplots(
        2, len(cases), figsize=(13.33, 6.8),
        sharex=True, sharey=sharey, dpi=None,
    )
    ax[0,0].set_ylabel(data['label'], y=-0.075)
    ax[0,0].set_xlim(2017.5, lastyear+2.5)
    ax[1,0].annotate(
        f'Diff\nfrom\n{basecase}', (0.03,0.03), xycoords='axes fraction',
        fontsize='x-large', weight='bold')
    ###### Absolute
    alltechs = set()
    for col, case in enumerate(cases):
        if case not in data['data']:
            continue
        dfplot = data['data'][case].pivot(index='year', columns='tech', values=data['values'])
        dfplot = (
            dfplot[[c for c in bokehcolors.index if c in dfplot]]
            .round(3).replace(0,np.nan)
            .dropna(axis=1, how='all')
        )
        if case == basecase:
            dfbase[slidetitle] = dfplot.copy()
        alltechs.update(dfplot.columns)
        plots.stackbar(df=dfplot, ax=ax[0,col], colors=techcolors, width=yearstep, net=False)
        ax[0,col].set_title(
            case.replace('__','\n'),
            fontsize=14, weight='bold', x=0, ha='left', pad=8,)
        ax[0,col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
        ax[0,col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))


    ### Legend
    handles = [
        mpl.patches.Patch(facecolor=techcolors[i], edgecolor='none', label=i.replace('Canada','imports'))
        for i in techcolors if i in alltechs
    ]
    leg = ax[0,-1].legend(
        handles=handles[::-1], loc='upper left', bbox_to_anchor=(1.0,1.0), 
        fontsize='medium', ncol=1,  frameon=False,
        handletextpad=0.3, handlelength=0.7, columnspacing=0.5, 
    )

    ###### Difference
    for col, case in enumerate(cases):
        ax[1,col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
        ax[1,col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))
        ax[1,col].axhline(0,c='k',ls='--',lw=0.75)

        if (case not in data['data']) or (case == basecase):
            continue
        dfplot = data['data'][case].pivot(index='year', columns='tech', values=data['values'])
        dfplot = (
            dfplot
            .round(3).replace(0,np.nan)
            .dropna(axis=1, how='all')
        )
        dfplot = dfplot.subtract(dfbase[slidetitle], fill_value=0)
        dfplot = dfplot[[c for c in bokehcolors.index if c in dfplot]].copy()
        alltechs.update(dfplot.columns)
        plots.stackbar(df=dfplot, ax=ax[1,col], colors=techcolors, width=yearstep, net=True)

    plots.despine(ax)
    plt.draw()
    plots.shorten_years(ax[1,0])
    ### Save it
    slide = add_to_pptx(slidetitle+' stack')
    if interactive:
        plt.show()


#%%### Hodgepodge 1: NPV, final capacity, final generation, runtime
plt.close()
f,ax = plt.subplots(
    2, 4, figsize=(13.33, 6.88), sharex=True,
    sharey=('col' if (sharey is True) else False),
)
handles = {}
### NPV of system cost
ax[0,0].set_ylabel('NPV of system cost [$B]', y=-0.075)
ax[0,0].axhline(0, c='k', ls='--', lw=0.75)
dfplot = pd.concat(dictin_npv, axis=1).T
dfplot = dfplot[[c for c in bokehcostcolors.index if c in dfplot]].copy()
## Absolute and difference
for (row, df) in enumerate([dfplot, dfplot - dfplot.loc[basecase]]):
    plots.stackbar(df=df, ax=ax[row,0], colors=bokehcostcolors, net=True, width=0.8)
    ymin, ymax = ax[row,0].get_ylim()
    ypad = (ymax - ymin) * 0.02
    ## label net value
    for x, case in enumerate(df.index):
        val = df.loc[case].sum()
        ax[row,0].annotate(
            f'{val:.0f}', (x, val - ypad), ha='center', va='top', color='k', size=9,
            path_effects=[pe.withStroke(linewidth=2.0, foreground='w', alpha=0.5)],
        )
## Legend info
handles['NPV'] = [
    mpl.patches.Patch(facecolor=bokehcostcolors[i], edgecolor='none', label=i)
    for i in bokehcostcolors.index if i in dfplot
]

### Final capacity and generation
toplot = {
    'Capacity': {'data': dictin_cap, 'values':'Capacity (GW)', 'label':f' {lastyear} Capacity [GW]'},
    'Generation': {'data': dictin_gen, 'values':'Generation (TWh)', 'label':f' {lastyear} Generation [TWh]'},
}
ax[0,2].axhline(0, c='k', ls='--', lw=0.75)
for col, (datum, data) in enumerate(toplot.items()):
    ax[0,col+1].set_ylabel(data['label'], y=-0.075)
    dfplot = pd.concat(
        {case:
         data['data'][case].loc[data['data'][case].year==lastyear]
         .set_index('tech')[data['values']]
         for case in cases},
        axis=1,
    ).T
    dfplot = (
        dfplot[[c for c in bokehcolors.index if c in dfplot]]
        .round(3).replace(0,np.nan).dropna(axis=1, how='all')
    )
    ## Absolute and difference
    for (row, df) in enumerate([dfplot, dfplot - dfplot.loc[basecase]]):
        plots.stackbar(df=df, ax=ax[row,col+1], colors=techcolors, width=0.8, net=row)
    ## Legend info
    handles[datum] = [
        mpl.patches.Patch(facecolor=techcolors[i], edgecolor='none', label=i)
        for i in techcolors if i in dfplot
    ]

### Runtime
ax[0,3].set_ylabel('Runtime [hours]', y=-0.075)
dfplot = pd.concat(dictin_runtime, axis=1).T
dfplot = dfplot[[c for c in colors_time.index if c in dfplot]].copy()
for (row, df) in enumerate([dfplot, dfplot - dfplot.loc[basecase]]):
    plots.stackbar(df=df, ax=ax[row,3], colors=colors_time, width=0.8, net=row)
## Legend info
handles['Runtime'] = [
    mpl.patches.Patch(facecolor=colors_time[i], edgecolor='none', label=i)
    for i in colors_time.index if i in dfplot
]

### Formatting
for col in range(4):
    ax[1,col].set_xticks(range(len(cases)))
    ax[1,col].set_xticklabels(cases.keys(), rotation=90)
    ax[1,col].annotate('Diff', (0.03,0.03), xycoords='axes fraction', fontsize='large')
    ax[1,col].axhline(0, c='k', ls='--', lw=0.75)
plt.tight_layout()
plots.despine(ax)
plt.draw()
### Save it
slide = add_to_pptx('Cost, capacity, generation, runtime')
if interactive:
    plt.show()

### Add legends as separate figure below the slide
plt.close()
f,ax = plt.subplots(1, 4, figsize=(13.33, 0.1))
for col, datum in enumerate(handles):
    leg = ax[col].legend(
        handles=handles[datum][::-1], loc='upper center', bbox_to_anchor=(0.5,1.0), 
        fontsize='medium', ncol=1, frameon=False,
        handletextpad=0.3, handlelength=0.7, columnspacing=0.5, 
    )
    ax[col].axis('off')
add_to_pptx(slide=slide, top=7.5)

#%% Diff

#%%### Hodgepodge 2: SCOE, CO2 emissions, Transmission TW-miles
plt.close()
f,ax = plt.subplots(1, 3, figsize=(10, 4.5))

### SCOE
for case in cases:
    df = dictin_scoe[case].groupby('year')['Average cost ($/MWh)'].sum().loc[years]
    ax[0].plot(df.index, df.values, label=case, color=colors[case])
    ## annotate the last value
    val = df.loc[max(years)]
    ax[0].annotate(
        f' {val:.0f}',
        (max(years), val), ha='left', va='center',
        color=colors[case], fontsize='medium',
    )
ax[0].set_ylim(0)
ax[0].set_ylabel('System cost of electricity [$/MWh]')
leg = ax[0].legend(
    loc='lower left', frameon=False, fontsize='large',
    handletextpad=0.3, handlelength=0.8,
)
for legobj in leg.legend_handles:
    legobj.set_linewidth(8)
    legobj.set_solid_capstyle('butt')

### CO2 emissions
for case in cases:
    df = dictin_emissions[case].reindex(years).fillna(0)
    ax[1].plot(df.index, df.values, label=case, color=colors[case])
    ## annotate the last value
    val = np.around(df.loc[max(years)], 0) + 0
    ax[1].annotate(
        f' {val:.0f}',
        (max(years), val), ha='left', va='center',
        color=colors[case], fontsize='medium',
    )
ax[1].set_ylim(0)
ax[1].axhline(phaseout_trigger, c='C7', ls='--', lw=0.75)
ax[1].set_ylabel(f'{pollutant} emissions [MMT/yr]')

### Transmission TW-miles
for case in cases:
    df = dictin_trans[case].groupby('year').sum()['Amount (GW-mi)'].reindex(years) / 1e3
    ax[2].plot(df.index, df.values, label=case, color=colors[case])
    ## annotate the last value
    val = np.around(df.loc[max(years)], 0) + 0
    ax[2].annotate(
        f' {val:.0f}',
        (max(years), val), ha='left', va='center',
        color=colors[case], fontsize='medium',
    )
ax[2].set_ylim(0)
ax[2].set_ylabel('Transmission capacity [TW-mi]')

### Formatting
plt.tight_layout()
plots.despine(ax)
plt.draw()
for col in range(3):
    ax[col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
    ax[col].xaxis.set_minor_locator(mpl.ticker.AutoMinorLocator(2))
    plots.shorten_years(ax[col])
### Save it
slide = add_to_pptx('Cost, emissions, and transmission')
if interactive:
    plt.show()


#%%### All-in-one maps
for case in cases:
    try:
        slide = add_to_pptx(
            case,
            file=os.path.join(
                cases[case],'outputs','maps',
                f'map_gencap_transcap-{lastyear}-sideplots.png'),
            width=None, height=6.88,
        )
    except FileNotFoundError:
        print(f'No all-in-one map for {os.path.basename(cases[case])}')


#%%### Generation capacity maps
### Capacity maps
nrows, ncols, coords = plots.get_coordinates(cases, aspect=2)
if (nrows == 1) or (ncols == 1):
    legendcoords = max(nrows, ncols) - 1
elif (nrows-1, ncols-1) in coords.values():
    legendcoords = (nrows-1, ncols-1)
else:
    legendcoords = (nrows-2, ncols-1)

### Shared data
base = cases[list(cases.keys())[0]]
val_r = dictin_cap_r[basecase].r.unique()
dfba = reedsplots.get_zonemap(base).loc[val_r]
dfstates = dfba.dissolve('st')

### Set up plot
for tech in maptechs:
    ### Get limits
    vmin = 0.
    vmax = float(pd.concat({
        case: dictin_cap_r[case].loc[
            (dictin_cap_r[case].i==tech)
            & (dictin_cap_r[case].t.astype(int)==lastyear)
        ].groupby('r').MW.sum()
        for case in cases
    }).max()) / 1e3
    if np.isnan(vmax):
        vmax = 0.
    if not vmax:
        print(f'{tech} has zero capacity in {lastyear}, so skipping maps')
        continue
    ### Set up plot
    plt.close()
    f,ax = plt.subplots(
        nrows, ncols, figsize=(13.33, 6.88),
        gridspec_kw={'wspace':0.0,'hspace':-0.1},
    )
    ### Plot it
    for case in cases:
        dfval = dictin_cap_r[case].loc[
            (dictin_cap_r[case].i==tech)
            & (dictin_cap_r[case].t.astype(int)==lastyear)
        ].groupby('r').MW.sum()
        dfplot = dfba.copy()
        dfplot['GW'] = (dfval / 1e3).fillna(0)

        ax[coords[case]].set_title(case)
        dfba.plot(
            ax=ax[coords[case]],
            facecolor='none', edgecolor='k', lw=0.1, zorder=10000)
        dfstates.plot(
            ax=ax[coords[case]],
            facecolor='none', edgecolor='k', lw=0.2, zorder=10001)
        dfplot.plot(
            ax=ax[coords[case]], column='GW', cmap=cmap, vmin=vmin, vmax=vmax,
            legend=False,
        )
        ## Legend
        if coords[case] == legendcoords:
            plots.addcolorbarhist(
                f=f, ax0=ax[coords[case]], data=dfplot.GW.values,
                title=f'{tech} {lastyear}\ncapacity [GW]', cmap=cmap, vmin=vmin, vmax=vmax,
                orientation='horizontal', labelpad=2.25, histratio=0.,
                cbarwidth=0.05, cbarheight=0.85,
                cbarbottom=-0.05, cbarhoffset=0.,
            )

    for row in range(nrows):
        for col in range(ncols):
            if nrows == 1:
                ax[col].axis('off')
            elif ncols == 1:
                ax[row].axis('off')
            else:
                ax[row,col].axis('off')
    ### Save it
    slide = add_to_pptx(f'{tech} capacity {lastyear} [GW]')
    if interactive:
        plt.show()

#%% Difference maps
### Set up plot
for tech in maptechs:
    ### Get limits
    dfval = pd.concat({
        case: dictin_cap_r[case].loc[
            (dictin_cap_r[case].i==tech)
            & (dictin_cap_r[case].t.astype(int)==lastyear)
        ].groupby('r').MW.sum()
        for case in cases
    }, axis=1).fillna(0) / 1e3
    dfdiff = dfval.subtract(dfval[basecase], axis=0)
    ### Get colorbar limits
    absmax = dfval.stack().max()
    diffmax = dfdiff.unstack().abs().max()

    if np.isnan(absmax):
        absmax = 0.
    if not absmax:
        print(f'{tech} has zero capacity in {lastyear}, so skipping maps')
        continue
    ### Set up plot
    plt.close()
    f,ax = plt.subplots(
        nrows, ncols, figsize=(13.33, 6.88),
        gridspec_kw={'wspace':0.0,'hspace':-0.1},
    )
    ### Plot it
    for case in cases:
        dfplot = dfba.copy()
        dfplot['GW'] = dfval[case] if case == basecase else dfdiff[case]

        ax[coords[case]].set_title(case)
        dfba.plot(
            ax=ax[coords[case]],
            facecolor='none', edgecolor='k', lw=0.1, zorder=10000)
        dfstates.plot(
            ax=ax[coords[case]],
            facecolor='none', edgecolor='k', lw=0.2, zorder=10001)
        dfplot.plot(
            ax=ax[coords[case]], column='GW',
            cmap=(cmap if case == basecase else cmap_diff),
            vmin=(0 if case == basecase else -diffmax),
            vmax=(absmax if case == basecase else diffmax),
            legend=False,
        )
        ## Difference legend
        if coords[case] == legendcoords:
            plots.addcolorbarhist(
                f=f, ax0=ax[coords[case]], data=dfplot.GW.values,
                title=f'{tech} {lastyear}\ncapacity, difference\nfrom {basecase} [GW]',
                cmap=(cmap if case == basecase else cmap_diff),
                vmin=(0 if case == basecase else -diffmax),
                vmax=(absmax if case == basecase else diffmax),
                orientation='horizontal', labelpad=2.25, histratio=0.,
                cbarwidth=0.05, cbarheight=0.85,
                cbarbottom=-0.05, cbarhoffset=0.,
            )
    ## Absolute legend
    plots.addcolorbarhist(
        f=f, ax0=ax[coords[basecase]], data=dfval[basecase].values,
        title=f'{tech} {lastyear}\ncapacity [GW]',
        cmap=cmap, vmin=0, vmax=absmax,
        orientation='horizontal', labelpad=2.25, histratio=0.,
        cbarwidth=0.05, cbarheight=0.85,
        cbarbottom=-0.05, cbarhoffset=0.,
    )

    for row in range(nrows):
        for col in range(ncols):
            if nrows == 1:
                ax[col].axis('off')
            elif ncols == 1:
                ax[row].axis('off')
            else:
                ax[row,col].axis('off')
    ### Save it
    slide = add_to_pptx(f'Difference: {tech} capacity {lastyear} [GW]')
    if interactive:
        plt.show()


#%%### Transmission maps
wscale = 0.0003
alpha = 0.8
for subtract_baseyear in [None, 2020]:
    plt.close()
    f,ax = plt.subplots(
        nrows, ncols, figsize=(13.33, 6.88),
        gridspec_kw={'wspace':0.0,'hspace':-0.1},
    )
    for case in cases:
        ### Plot it
        reedsplots.plot_trans_onecase(
            case=cases[case], pcalabel=False, wscale=wscale,
            yearlabel=False, year=lastyear, simpletypes=None,
            alpha=alpha, scalesize=8,
            f=f, ax=ax[coords[case]], title=False,
            subtract_baseyear=subtract_baseyear,
            thickborders='transreg', drawstates=False, drawzones=False, 
            label_line_capacity=10,
            scale=(True if case == basecase else False),
        )
        ax[coords[case]].set_title(case)
    ### Formatting
    title = (
        f'New interzonal transmission since {subtract_baseyear}' if subtract_baseyear
        else 'All interzonal transmission')
    for row in range(nrows):
        for col in range(ncols):
            if nrows == 1:
                ax[col].axis('off')
            elif ncols == 1:
                ax[row].axis('off')
            else:
                ax[row,col].axis('off')
    ### Save it
    slide = add_to_pptx(title)
    if interactive:
        plt.show()

#%% Transmission difference
plt.close()
f,ax = plt.subplots(
    nrows, ncols, figsize=(13.33, 6.88),
    gridspec_kw={'wspace':0.0,'hspace':-0.1},
)
for case in cases:
    ax[coords[case]].set_title(case)
    if case == basecase:
        ### Plot absolute
        reedsplots.plot_trans_onecase(
            case=cases[case], pcalabel=False, wscale=wscale,
            yearlabel=False, year=lastyear, simpletypes=None,
            alpha=alpha, scalesize=8,
            f=f, ax=ax[coords[case]], title=False,
            subtract_baseyear=subtract_baseyear,
            thickborders='transreg', drawstates=False, drawzones=False, 
            label_line_capacity=10,
            scale=(True if case == basecase else False),
        )
    else:
        ### Plot the difference
        reedsplots.plot_trans_diff(
            casebase=cases[basecase], casecomp=cases[case],
            pcalabel=False, wscale=wscale,
            yearlabel=False, year=lastyear, simpletypes=None,
            alpha=alpha,
            f=f, ax=ax[coords[case]],
            subtract_baseyear=subtract_baseyear,
            thickborders='transreg', drawstates=False, drawzones=False, 
            label_line_capacity=10,
            scale=False,
        )
### Formatting
title = 'Interzonal transmission difference'
for row in range(nrows):
    for col in range(ncols):
        if nrows == 1:
            ax[col].axis('off')
        elif ncols == 1:
            ax[row].axis('off')
        else:
            ax[row,col].axis('off')
### Save it
slide = add_to_pptx(title)
if interactive:
    plt.show()


#%%### Interregional transfer capability to peak demand ratio
try:
    f, ax, dfplot = reedsplots.plot_interreg_transfer_cap_ratio(
        case=list(cases.values()),
        colors={v: colors[k] for k,v in cases.items()},
        casenames={v:k for k,v in cases.items()},
        level=level, tstart=startyear,
        ymax=None,
    )
    ### Save it
    slide = add_to_pptx('Interregional transmission / peak demand')
    if interactive:
        plt.show()

except Exception:
    print(traceback.format_exc())


#%% Save the powerpoint file
prs.save(savename)
print(f'\ncompare_casegroup.py results saved to:\n{savename}')

### Open it
if sys.platform == 'darwin':
    sp.run(f"open '{savename}'", shell=True)
elif platform.system() == 'Windows':
    sp.run(f'"{savename}"', shell=True)
