"""
This module contains functions for plotting the uncertainty in the inputs and outputs of ReEDS.
It was originally developed to serve the MCS and MGA functionality in the model, but 
it can be used for other purposes, like mapping a large number of scenarios to help
understand the differences between specific inputs and outputs.
"""


#%% ===========================================================================
### --- IMPORTS ---
### ===========================================================================
import os
import sys
import numpy as np
import pandas as pd
import copy
import re
from tqdm import tqdm
import argparse
import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D
from matplotlib.ticker import FuncFormatter
import pptx
from pptx.util import Inches
import io
import seaborn as sns

# Local Imports
reeds_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
sys.path.append(reeds_path)
import reeds

# Define global matplotlib plot configurations
reeds.plots.plotparams()


#%% ===========================================================================
### --- General Helper Functions ---
### ===========================================================================
class Conventions:
    """
    Defines standard naming and styling conventions used for ReEDS post-processing and plots.

    Includes:
        - Technology renaming maps for various simplification levels.
        - Color palettes for technologies and transmission types.
        - Cost and transmission category mappings.
    """
    # --- Technology name mappings ---
    # tech_map1: Most simplified groupings
    tech_map1 = {
        **{f'upv_{i}':'PV' for i in range(20)},
        **{f'dupv_{i}':'PV' for i in range(20)},
        **{f'wind-ons_{i}':'Wind' for i in range(20)},
        **{f'wind-ofs_{i}':'Wind' for i in range(20)},
        **dict(zip(['nuclear','nuclear-smr', 'smr'], ['Nuclear']*20)),
        **dict(zip(
            ['gas-cc_re-cc','gas-ct_re-ct','re-cc','re-ct',
            'gas-cc_h2-ct','gas-ct_h2-ct','h2-cc','h2-ct',],
            ['H2 turbine']*20)),
        **{f'battery_{i}':'Storage' for i in range(20)},
        **{'battery_li':'Storage', 'pumped-hydro':'Storage'},
        **dict(zip(
            ['coal-igcc', 'coaloldscr', 'coalolduns', 'gas-cc', 'gas-ct', 'coal-new', 'o-g-s',],
            ['Fossil']*20)),
        **dict(zip(
            ['gas-cc_gas-cc-ccs_mod','gas-cc_gas-cc-ccs_max','gas-cc-ccs_mod','gas-cc-ccs_max',
            'gas-cc_gas-cc-ccs_mod','coal-igcc_coal-ccs_mod','coal-new_coal-ccs_mod',
            'coaloldscr_coal-ccs_mod','coalolduns_coal-ccs_mod','cofirenew_coal-ccs_mod',
            'cofireold_coal-ccs_mod','gas-cc_gas-cc-ccs_max','coal-igcc_coal-ccs_max',
            'coal-new_coal-ccs_max','coaloldscr_coal-ccs_max','coalolduns_coal-ccs_max',
            'cofirenew_coal-ccs_max','cofireold_coal-ccs_max',
            'coal-ccs_mod',
            ],
            ['Fossil+CCS']*50)),
        **dict(zip(['dac','beccs_mod','beccs_max'],['CO2 removal']*20)),
        **{f'egs_nearfield_{i}':'Geothermal' for i in range(20)},
        **{f'geohydro_allkm_{i}':'Geothermal' for i in range(20)},
        **{f'csp{i+1}_{j}':'CSP' for i in range(4) for j in range(20)},
        **{'csp-ns':'CSP'},
        **dict(zip(['hyded','hydend','hydnd','hydnpnd','hydud','hydund'], ['Hydro']*20)),
        **dict(zip(['biopower','can-imports','distpv','lfill-gas', "electrolyzer"], ['Others']*20)),
    }

    # tech_map2: Intermediate simplification
    tech_map2 = {
        **{f'upv_{i}': 'PV' for i in range(20)},
        **{f'dupv_{i}': 'PV' for i in range(20)},
        **{f'wind-ons_{i}': 'Wind' for i in range(20)},
        **{f'wind-ofs_{i}': 'Offshore Wind' for i in range(20)},
        'nuclear': 'Nuclear',
        **dict(zip(['nuclear-smr', 'smr'], ['Nuclear-SMR']*20)),
        **dict(zip(
            ['gas-cc_re-cc','gas-ct_re-ct','re-cc','re-ct',
            'gas-cc_h2-ct','gas-ct_h2-ct','h2-cc','h2-ct'],
            ['H2 Turbine'] * 8
        )),
        **{f'battery_{i}': 'Storage' for i in range(20)},
        **{'battery_li':'Storage', 'pumped-hydro':'Storage'},
        **dict(zip(
            ['coal-igcc', 'coaloldscr', 'coalolduns', 'coal-new'],
            ['Coal'] * 4
        )),
        **dict(zip(
            ['coal-ccs_mod',
            'coal-igcc_coal-ccs_mod','coal-igcc_coal-ccs_max',
            'coal-new_coal-ccs_mod','coal-new_coal-ccs_max',
            'coaloldscr_coal-ccs_mod','coaloldscr_coal-ccs_max',
            'coalolduns_coal-ccs_mod','coalolduns_coal-ccs_max',
            'cofirenew_coal-ccs_mod','cofirenew_coal-ccs_max',
            'cofireold_coal-ccs_mod','cofireold_coal-ccs_max'],
            ['Coal + CCS'] * 13
        )),
        'gas-cc': 'Gas CC',
        'gas-ct': 'Gas CT',
        **dict(zip(
            ['gas-cc-ccs_mod','gas-cc-ccs_max',
            'gas-cc_gas-cc-ccs_mod','gas-cc_gas-cc-ccs_max',
            'gas-ct_gas-ct-ccs_mod','gas-ct_gas-ct-ccs_max'],
            ['Gas + CCS'] * 6
        )),
        'o-g-s': 'o-g-s',
        **dict(zip(['dac','beccs_mod','beccs_max'], ['CO2 Removal'] * 3)),
        **{f'egs_nearfield_{i}': 'Geothermal' for i in range(20)},
        **{f'geohydro_allkm_{i}': 'Geothermal' for i in range(20)},
        **{f'csp{i+1}_{j}': 'CSP' for i in range(4) for j in range(20)},
        'csp-ns': 'CSP',
        **dict(zip(['hyded','hydend','hydnd','hydnpnd','hydud','hydund'], ['Hydro'] * 6)),
        **dict(zip(['biopower','can-imports','distpv','lfill-gas', "electrolyzer"], ['Others']*20)),
    }

    # tech_map3: Bokehpivot style
    tech_map3 = pd.read_csv(
        os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_map.csv'),
        index_col='raw')['display'].to_dict()

    # --- Technology colors ---
    tech_colors1 = {
        'PV': '#FFC903', 'CSP': '#FC761A', 'Wind': '#00B6EF', 'Offshore wind': '#106BA7',
        'Hydro': '#187F94', 'Geothermal': '#A96235',
        'Nuclear': '#820000', 'Nuclear-SMR': '#D40000',
        'Storage': '#FF4A88',
        'Fossil': '#5E1688', 'Fossil+CCS': '#9467BD', 'Coal': '#000000', 'Coal + CCS': '#666666',
        'Gas CC': '#5E1688', 'Gas CT': '#C2A1DB', 'Gas + CCS': '#9467BD', 'o-g-s': '#3D3376',
        'H2 turbine': '#708238', 'CO2 Removal': '#66C238', 'Others': '#A0A0A0',
    }
    tech_colors2 = tech_colors1

    tech_colors3 = pd.read_csv(
        os.path.join(reeds_path,'postprocessing','bokehpivot','in','reeds2','tech_style.csv'),
        index_col='order').squeeze(1)

    # --- Transmission type mapping and colors ---
    trtype_map = {
        'ac': 'AC',
        'b2b': 'B2B',
        'lcc': 'DC/LCC',
        'vsc': 'DC/VSC',
        'reinforcement': 'Reinf',
        'spur': 'Spur'
    }

    trtype_colors = {
        'AC (initial)': "#98DF8A",
        'DC, LCC (initial)': "#FFBB78",
        'AC': "#2CA02C",
        'B2B': "#9467BD",
        'DC/LCC': "#FF7F0E",
        'DC/VSC': "#D62728",
        'Reinf': "#339fff",
        'Spur': "#969696",
    }

    # --- Cost category mapping and colors ---
    costcat_map = pd.read_csv(
        os.path.join(
            reeds_path, 'postprocessing', 'bokehpivot', 'in', 'reeds2', 'cost_cat_map.csv'
        ),
    ).set_index('raw')["display"].to_dict()

    cost_colors = pd.read_csv(
        os.path.join(
            reeds_path,'postprocessing','bokehpivot','in','reeds2','cost_cat_style.csv'),
        index_col='order').squeeze(1)
    cost_colors['Total Cost'] = '#3e9c35'  


    # --- Map of tech/cost/trtype maps to color maps ---
    items_color_map = {
        "tech_map1": (tech_map1, tech_colors1),
        "tech_map2": (tech_map2, tech_colors2),
        "tech_map3": (tech_map3, tech_colors3),
        "trtype_map": (trtype_map, trtype_colors),
        "costcat_map": (costcat_map, cost_colors),
    }

    ### ===========================================================================
    ### --- Tech input manipulation (col rename and mapping) ---
    ### ===========================================================================
    tech_col_rename = {
        't': 'year', 'i': 'tech',
        'Year': 'year', 'Turbine': 'tech',
        'CF_mult': 'cf_mult', 'cf_improvement': 'cf_mult',
        'type': 'tech', 'Tech': 'tech',
        'capcost': 'capcost ($/kW)',
        'Cap cost 1000$/MW': 'capcost ($/kW)',
        'Overnight Cap Cost $/kW': 'capcost ($/kW)',
        'fom': 'fom ($/kW-yr)',
        'Fixed O&M $/(kW-yr)': 'fom ($/kW-yr)',
        'Fixed O&M 1000$/MW-yr': 'fom ($/kW-yr)',
        'vom': 'vom ($/MWh)',
        'Var O&M $/MWh': 'vom ($/MWh)',
        'heatrate': 'heatrate (MMBtu/MWh)',
    }


def get_default_color_map(items: list) -> dict:
    """
    Generate a default color map for a list of items.
    Cycles through colors if more items than available colors.

    Args:
        items (list): List of items to generate colors for.

    Returns:
        dict: Mapping of item names to colors.
    """
    color_map = {}
    base_colormap = plt.get_cmap('tab10') 

    for i, item in enumerate(items):
        color = base_colormap(i % base_colormap.N)  # Repeat colors as needed
        color_map[item] = mpl.colors.to_hex(color)

    return color_map


def map_cases(
    cases_location: str,
    prefix_old_new: list|str = None,
    highlight_runs: list|str = None,
    filter_unfinished: bool = True,
) -> dict:
    """
    Map ReEDS runs to a dictionary of case paths.

    Args:
        cases_location (str): 
            Path to the folder containing multiple ReEDS runs.

        prefix_old_new (list of tuple or str, optional):
            - List of (old_prefix, new_prefix) tuples used to rename and group cases.

            - If a string, it is interpreted as a CSV or Excel file path. Where
              the first column is the old prefix and the second column is the new prefix.

            - If None, the prefix defaults to the part of the case name before the
              last underscore.
              Example: 'i0_too_much_data_MCS1' -> prefix becomes 'i0_too_much_data'

            - Prefixes are used to group related cases and trigger comparative plots in the reports.

            - Prefixes like 'base', 'ref', or 'reference' receive special treatment
              and are treated as reference scenarios during comparison.

            - If prefix_old_new is not None, any case without a prefix match is skipped

        highlight_runs (list of tuple, optional):
            - List of runs you want to highlight in the report.
              Each tuple should contain (case_name, print_name, color)

            - If a string, it is interpreted as a CSV or Excel file path.
              The first column is the case_name, the second is the print_name.

        filter_unfinished (bool, optional):
            If True, only include runs that are considered finished—
            i.e., those containing a report.xlsx file under 
            outputs/reeds-report. Default is True.

    Returns:
        cases_dict (dict): A mapping of {prefix|case_name: full_case_path}.
        highlight_runs_f (list): Filtered list of runs to highlight with only those
            that are present in the cases_dict.
    """
    def is_reeds_case(folder_path: str) -> bool:
        """Check if a folder contains a ReEDS run (identified by gamslog.txt)."""
        return os.path.isfile(os.path.join(folder_path, 'gamslog.txt'))

    def is_finished_case(folder_path: str) -> bool:
        """Check if a ReEDS run has finished (identified by a report.xlsx file)."""
        return os.path.isfile(os.path.join(folder_path, 'outputs', 'reeds-report', 'report.xlsx'))

    print(f"Mapping cases in '{cases_location}'...")

    if prefix_old_new and isinstance(prefix_old_new, str):
        prefix_old_new_list = get_rename_cases_list(prefix_old_new)

    elif prefix_old_new and isinstance(prefix_old_new, list):
        prefix_old_new_list = prefix_old_new

    else:
        if prefix_old_new is not None:
            raise ValueError(
                "prefix_old_new must be a list of tuples or a string representing a file path."
            )

    if prefix_old_new is not None and len(prefix_old_new) == 0:
        raise ValueError("prefix_old_new must be a non-empty list of tuples.")

    # Check if the first element of each tuple is unique and not part of any other tuple
    prefixes = [old_prefix for old_prefix, _ in prefix_old_new_list]
    if len(prefixes) != len(set(prefixes)):
        raise ValueError(
            "The first element of each tuple in prefix_old_new must be unique."
        )

    for old_prefix, new_prefix in prefix_old_new_list:
        validation = [(old_prefix not in p) for p in prefixes if old_prefix != p]
        if not all(validation):
            raise ValueError(
                f"\nThe prefix '{old_prefix}' is part of another prefix in the list.\n"
                "This will cause conflicts in the mapping. \n"
                "Try using _MC or _MG to separate the prefixes."
            )

    # Get tuple list of highlight runs
    if highlight_runs and isinstance(highlight_runs, str):
        highlight_runs = get_rename_cases_list(highlight_runs)

    elif highlight_runs and isinstance(highlight_runs, list):
        highlight_runs = highlight_runs

    else:
        if highlight_runs is not None:
            raise ValueError(
                "highlight_runs must be a list of tuples or a string representing a file path."
            )

    cases_with_order = []

    for folder_name in os.listdir(cases_location):
        folder_path = os.path.join(cases_location, folder_name)

        if not os.path.isdir(folder_path):
            continue  # skip files

        if not is_reeds_case(folder_path):
            print(f"[Warning] Skipping '{folder_name}': missing gamslog.txt")
            continue

        prefix_rank = float('inf')  # Goes to the end
        old_case_name = os.path.basename(folder_path)

        # Default: prefix is everything before the last underscore
        case_key = None
        if prefix_old_new_list:

            for i, (old_prefix, new_prefix) in enumerate(prefix_old_new_list):
                new_prefix = new_prefix or old_prefix
                if folder_name.startswith(old_prefix):
                    case_key = f"{new_prefix}|{folder_name}"
                    prefix_rank = i
                    break

            if case_key is None:
                print(f"[Warning] Skipping '{folder_name}': no prefix match")
                continue # Skip case

        else:
            prefix = '_'.join(old_case_name.split('_')[:-1])
            case_key = f"{prefix}|{folder_name}"

            # If the default prefix is too big (+15 chars)
            # raise an error and ask the user to provide a prefix_old_new_list
            # Plots get messy with long prefixes
            if len(prefix) > 15:
                raise ValueError(
                    f"The default prefix '{prefix}' for '{folder_name}' is too long.\n"
                    "Please provide a prefix_old_new_list to rename the case.\n"
                    f"Example: python uncertainty_plots.py runs_folders -p {prefix}:small_prefix "
                    "other_prefix:small_prefix2"
                )
        cases_with_order.append((prefix_rank, case_key, folder_path))

    cases_with_order.sort()
    cases_dict = dict((case_key, folder_path) for _, case_key, folder_path in cases_with_order)

    # Check if we need to filter out unfinished cases
    if filter_unfinished:
        cases_dict = {
            key: value
            for key, value in cases_dict.items()
            if is_finished_case(value)
        }

    # Validate highlight list: Its items must be in the cases_dict
    highlight_runs_f = []
    if highlight_runs:
        valid_case_keys = list(cases_dict.keys())
        case_lookup = {case_key.split('|')[1]: case_key for case_key in valid_case_keys}

        for case_name, print_name, color in highlight_runs:
            case_key = case_lookup.get(case_name)
            if case_key:
                highlight_runs_f.append((case_key, print_name, color))
            else:
                print(
                    f"[Warning] Highlight run '{case_name}' not found in the mapped cases. "
                    "This run will be skipped."
                )

    return cases_dict, highlight_runs_f


def get_rename_cases_list(path: str) -> list:
    """
    Load case renaming or highlight configuration from a CSV or Excel file.

    The file must have:
      - Two columns: (old_prefix, new_prefix), for renaming cases.
      - Three columns: (case_name, print_name, color), for highlighting runs.

    Args:
        path (str): Path to the .csv or .xlsx file.

    Returns:
        list: A list of tuples:
              - (old_prefix, new_prefix) if 2 columns
              - (case_name, print_name, color) if 3 columns

    Raises:
        ValueError: If the file format is not supported or has an invalid number of columns.
    """
    if path.endswith('.csv'):
        df = pd.read_csv(path, header=None)
    elif path.endswith('.xlsx'):
        df = pd.read_excel(path, header=None)
    else:
        raise ValueError("Unsupported file format. Please provide a .csv or .xlsx file.")

    if df.shape[1] == 2:
        return list(zip(df.iloc[:, 0], df.iloc[:, 1]))

    if df.shape[1] == 3:
        return list(zip(df.iloc[:, 0], df.iloc[:, 1], df.iloc[:, 2]))

    raise ValueError(
        f"Invalid file format: expected 2 or 3 columns, got {df.shape[1]}.\n"
        "Accepted formats:\n"
        " - (old_prefix, new_prefix)\n"
        " - (case_name, print_name, color)"
    )


#%% ===========================================================================
### --- Read ReEDS Input/Output Class ---
### ===========================================================================
class DataFetcher:
    """
    This class manages the fetching of ReEDS input and output data.

    Args:
        cases_dict (dict): Dictionary mapping (prefix|case_name) to full case path.

    Attributes:
        data (dict): Stores fetched data dynamically.
        hierarchy (dict): Hierarchy mapping from the first case.
            We assume the same hierarchy for all cases.
        dfmap (dict): Dict of maps from the first case.
        FETCH_METHODS (dict): Combined mapping of keys to fetch functions.
    """
    def __init__(self, cases_dict: dict):
        self.cases_dict = cases_dict
        self.data = {} # Stores fetched data dynamically

        # Use first case as reference
        first_case_path = next(iter(cases_dict.values()))
        self.hierarchy = reeds.io.get_hierarchy(first_case_path)
        self.dfmap = reeds.io.get_dfmap(first_case_path)
        self.sw = reeds.io.get_switches(first_case_path)

        self.prefix_order = list(
            dict.fromkeys(key.split("|")[0] for key in cases_dict)
        )

        # --- Fetch methods for output data ---
        self.OUT_FETCH_METHODS = {
            # Naming as in e_report_params.csv
            "cap_out": self._fetch_cap_out,
            "gen_ann": self._fetch_gen_ann,
            "tran_mi_out_detail": self._fetch_tran_mi_out_detail,
            "tran_cap_energy": self._fetch_tran_cap_energy,
            "emit_r": self._fetch_emit_r,
            # System NPV for all regions
            "npv_r": self._fetch_npv_r,
        }

        # --- Fetch Methods for input data ---
        self.INPUT_FETCH_METHODS = {
            "demand": self._fetch_demand,
            "tech_inputs": self._fetch_tech_inputs,
            "supply_curves_capacity": self._fetch_supply_capacity,
            "ng_fuel_prices": self._fetch_fuel_prices,
        }

        self.FETCH_METHODS = {
            **self.OUT_FETCH_METHODS,
            **self.INPUT_FETCH_METHODS,
        }


    def fetch_all(self, data_keys: str|list = None, force_reload: bool = False) -> dict:
        """
        Fetch all data for the given data_keys.

        Args:
            data_keys (str or list, optional): Data types to fetch from self.FETCH_METHODS. 
                If None, fetch all data types.
            force_reload (bool, optional): Force reload data even if it has been fetched before.

        Returns:
            dict: Mapping of data_key to {case: DataFrame}
        """
        print("(DataFetcher) Fetching data...")

        if data_keys is None:
            data_keys = list(self.FETCH_METHODS.keys())
        elif isinstance(data_keys, str):
            data_keys = [data_keys]

        for data_key in data_keys:
            for case in tqdm(self.cases_dict, desc=f"(DataFetcher) Fetching {data_key}"):
                self.fetch_case(data_key, case)

        return {data_key: self.data[data_key] for data_key in data_keys}


    def fetch_case(
        self,
        data_key: str,
        case: str,
        force_reload: bool = False
    ) -> pd.DataFrame:
        """
        Fetch a specific data_key for a given case.

        Args:
            data_key (str): Data type to fetch from self.FETCH_METHODS
            case (str): (prefix|case_name) from cases_dict
            force_reload (bool, optional): Force reload data even if it has been fetched before.

        Returns:
            pd.DataFrame: Data for the given case and data_key
        """
        if data_key not in self.FETCH_METHODS:
            raise ValueError(
                f"Invalid fetch request '{data_key}' for case {case}. "
                f"Available keys: {list(self.FETCH_METHODS.keys())}"
            )

        if data_key not in self.data:
            # Initialize storage for this data type
            self.data[data_key] = {}

        if case not in self.data[data_key] or force_reload:
            # Only call if data has not been fetched yet
            self.data[data_key][case] = self.FETCH_METHODS[data_key](case)

        return self.data[data_key][case]


    ### ===========================================================================
    ### - Output Data Fetchers
    ### ===========================================================================
    def _fetch_cap_out(self, case: str) -> pd.DataFrame:
        """
        Fetch and process capacity data for a case (key) of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['tech', 'year', 'Capacity (GW)'].
        """
        df = reeds.io.read_output(case=self.cases_dict[case], filename='cap')
        df['Value'] /= 1000  # Convert MW to GW
        df.rename(
            columns={'Value': 'Capacity (GW)', 't': 'year', 'i': 'tech'},
            inplace=True,
        )
        return df


    def _fetch_gen_ann(self, case: str) -> pd.DataFrame:
        """
        Fetch and process annual generation data for a case (key) 
        of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['tech', 'year', 'Generation (TWh)'].
        """
        df = reeds.io.read_output(case=self.cases_dict[case], filename='gen_ann')
        df['Value'] /= 1e6  # Convert MWh to TWh
        df.rename(
            columns={'Value': 'Generation (TWh)', 't': 'year', 'i': 'tech'},
            inplace=True,
        )
        return df


    def _fetch_tran_mi_out_detail(self, case: str) -> pd.DataFrame:
        """
        Fetch transmission and reinforcement data for a case (key) 
        of the cases_dict.

        Returns:
            pd.DataFrame: Combined transmission and reinforcement data.
                Columns ['r', 'rr', 'trtype', 'year', 'Trans (TW-mi)']
        """
        df_trans = reeds.io.read_output(
            case=self.cases_dict[case],
            filename='tran_mi_out_detail',
        )
        df_trans['Value'] *= 1e-6  # Convert MW-mi to TW-mi
        df_trans.rename(
            columns={'Value': 'Trans (TW-mi)', 't': 'year'},
            inplace=True,
        )

        df_tech_trans = (
            reeds.output_calc.calc_reinforcement_spur_capacity_miles(self.cases_dict[case])
        )
        df_tech_trans['rr'] = df_tech_trans['r']  # Duplicate region data
        # Convert "Trans (GW-mi)" to "Trans (TW-mi)"
        df_tech_trans["Trans (TW-mi)"] = df_tech_trans["Trans (GW-mi)"] / 1000
        df_tech_trans.drop(columns="Trans (GW-mi)", inplace=True)

        return pd.concat([df_trans, df_tech_trans], ignore_index=True)


    def _fetch_tran_cap_energy(self, case: str) -> pd.DataFrame:
        """
        Fetch transmission capacity r->rr data for a case (key)
        of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['r', 'rr', 'trtype', 'year', 'Capacity (GW)']
        """
        df = reeds.io.read_output(case=self.cases_dict[case], filename='tran_cap_energy')
        df['Value'] /= 1000
        df.rename(
            columns={'Value': 'Capacity (GW)', 't': 'year'},
            inplace=True,
        )

        return df

    def _fetch_emit_r(self, case: str) -> pd.DataFrame:
        """
        Fetch emissions data for a case (key) of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['r', 'e_type', 'year', ' Em. (MMT)']
        """
        df = reeds.io.read_output(case=self.cases_dict[case], filename='emit_r')
        df['Value'] /= 1e6  # Convert from MT to MMT
        df.rename(
            columns = {'Value': 'Em. (MMT)', 't': 'year'},
            inplace=True,
        )

        # Remove (precombustion/combustion) type 
        df = df.groupby(['e', 'r', 'year'], as_index=False)['Em. (MMT)'].sum()
        df.rename(columns={'e': 'e_type'}, inplace=True)

        return df

    def _fetch_npv_r(self, case: str) -> pd.DataFrame:
        """
        Compute net present value of system cost for a case (key) 
        of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['year', 'r', 'cost_cat', 'Cost (Bil $)', 'Discounted Cost (Bil $)']
        """
        return reeds.output_calc.calc_systemcost(self.cases_dict[case])


    ### ===========================================================================
    ### - Input Data Fetchers
    ### ===========================================================================
    def _fetch_demand(self, case: str,  hourly = False, representative = True) -> pd.DataFrame:
        """
        Fetch demand data for a case (key) of the cases_dict.

        Args:
            case (str): (prefix|case_name) from cases_dict
            hourly (bool): If True, fetch hourly demand data. Else fetch annual data.
            representative (bool): If True, fetch only representative periods.

        Return:
            pd.DataFrame: 
                Columns ['r', 'year', 'Demand (TWh)']
                OR Columns ['r', 'year', 'h', 'Demand (TWh)']
        """

        if representative:
            # This files only get create if you run all the input processing
            numhours_path = os.path.join(
                self.cases_dict[case], 'inputs_case', 'rep', 'numhours.csv'
            )

            load_allyear_path = os.path.join(
                self.cases_dict[case], 'inputs_case', 'rep', 'load_allyear.csv'
            )

            numhours = pd.read_csv(numhours_path).rename(columns={'*h':'h'}).set_index('h').numhours
            demand = pd.read_csv(load_allyear_path).rename(
                columns={'*r':'r', 't':'year', 'MW':'Load (TWh)'}
            )

            # Convert from MW to TWh
            demand['Load (TWh)'] *= demand.h.map(numhours) / 1e6

            if not hourly:
                # Group by year and region
                demand = demand.groupby(['r','year'])['Load (TWh)'].sum().reset_index()

        else:
            # Read from the .h5 (take sone time)
            load_h5_path = os.path.join(
                self.cases_dict[case], 'inputs_case', 'load_hourly.h5'
            )
            demand = reeds.io.read_file(load_h5_path).reset_index()
            n_weather_years = demand.shape[0]/len(demand.year.unique())/(365*24)

            demand.rename(columns={'datetime': 'h'}, inplace=True)
            demand = demand.melt(id_vars=['year', 'h'], var_name='r', value_name='Load (TWh)')
            demand['Load (TWh)'] /= 1e6  # Convert from MW to TWh

            if not hourly:
                # Group by year and region
                demand = demand.groupby(['r','year'])['Load (TWh)'].sum()/n_weather_years
                demand.reset_index(inplace=True)

        # Create a item column
        demand.insert(len(demand.columns) - 1, "Load Type", "Total Load")
        return demand

    def _fetch_tech_inputs(self, case: str) -> pd.DataFrame:
        """Load and prepare all plant characteristic inputs for a given case."""
        input_dir = os.path.join(self.cases_dict[case], 'inputs_case')
        output_cols = [
            'year', 'tech', 'capcost ($/kW)', 'fom ($/kW-yr)',
            'vom ($/MWh)', 'cf_mult', 'heatrate (MMBtu/MWh)'
        ]

        configs = {
            'plantchar_upv.csv': {'insert': {'tech': 'UPV'}},
            'plantchar_onswind.csv': {'insert': {'tech': 'On. wind'}},
            'plantchar_ofswind.csv': {
                'drop': ['rsc_mult'],
                'post': lambda df: df.assign(tech='Off. wind ' + df['tech'].astype(str))
            },
            'plantchar_battery.csv': {'drop': ['rte']},
            'plantchar_csp.csv': {},
            'plantchar_geo.csv': {
                'drop': ['Geo class', 'Depth'],
                'pre': lambda df: df.assign(
                    tech=(
                        df['tech'].astype(str)
                        + '_' + df['Geo class'].astype(str)
                        + '_' + df['Depth'].astype(str)
                    )
                )
            },
            **{f'plantchar_{kind}.csv': {}
            for kind in [
                'gas', 'gas_ccs',
                'coal', 'coal_ccs',
                'nuclear', 'nuclear_smr'
            ]
            }
        }

        def load_and_prepare(fname: str, cfg: dict) -> pd.DataFrame:
            path = os.path.join(input_dir, fname)
            df = pd.read_csv(path)
            df = df.rename(
                columns={k: v for k, v in Conventions.tech_col_rename.items()
                        if k in df.columns}
            )
            df = df.pipe(cfg.get('pre', lambda x: x))

            if 'insert' in cfg:
                df = df.assign(**cfg['insert'])
            df = df.pipe(cfg.get('post', lambda x: x))
            return df.loc[:, [c for c in output_cols if c in df.columns]]

        dfs = [
            load_and_prepare(fname, cfg)
            for fname, cfg in configs.items()
        ]
        return pd.concat(dfs, ignore_index=True)

    def _fetch_supply_capacity(self, case: str) -> pd.DataFrame:
        """Get supply curve capacity data"""
        input_dir = os.path.join(self.cases_dict[case], 'inputs_case')

        files = {
            'UPV': 'upv_supply_curve.csv',
            'Wind On': 'wind-ons_supply_curve.csv',
            'Wind Off': 'wind-ofs_supply_curve.csv',
        }

        supply_curves = []
        for tech, fname in files.items():
            path = os.path.join(input_dir, fname)
            df = pd.read_csv(path)
            df.rename(
                columns={
                    'capacity': 'Capacity (GW)',
                    'region': 'r',
                },
            inplace=True)
            df['tech'] = tech
            df['Capacity (GW)'] /= 1000  # Convert from MW to GW
            df = df[['tech', 'r', 'Capacity (GW)']]

            supply_curves.append(df)

        # Concatenate all supply curves
        supply_curves = pd.concat(supply_curves, ignore_index=True)
        supply_curves = supply_curves.groupby(['tech', 'r'], as_index=False).sum()

        return supply_curves

    def _fetch_fuel_prices(self, case: str) -> pd.DataFrame:
        """
        Fetch fuel prices for a case (key) of the cases_dict.

        Returns:
            pd.DataFrame: Columns ['r', 'year', type, 'Price ($/MMBtu)']
        """
        df = pd.read_csv(
            os.path.join(self.cases_dict[case], 'inputs_case', 'natgas_price_cendiv.csv')
        )

        df = df.melt(
            id_vars=['year'],
            var_name='r',
            value_name='Price ($/MMBtu)',
        )

        df.insert(2, "type", "NG Price")

        return df


    ### ===========================================================================
    ### - Input/Output Transformations
    ### ===========================================================================
    def apply_transformations(
        self,
        data_key: str,
        items_map_name: str = None,
        r_filter: list = None,
        hierarchy_group: str = None,
        normalize_by_r: bool = False,
    ) -> dict:
        """
        Apply a series of transformations to the fetched data, optionally including:
        - filtering by region,
        - simplifying item names (e.g. techs, costs),
        - grouping by hierarchy (e.g. state, region, transreg, ...),
        - normalization (e.g. penetration shares per region).

        These steps can be composed to derive new metrics. For example, using "gen_ann"
        with normalization can produce technology penetration (% of generation by tech).

        Args:
            data_key (str): Key from self.data to apply transformations to.
            items_map_name (str): Name of the items map to use for simplifying item names.
                - Options: "tech_map1", "tech_map2", "tech_map3", "trtype_map", "costcat_map".
                  Based on Conventions.items_color_map.
                - If None, no mapping is applied (the original item names are used), and a 
                    default color map is used.
            r_filter (list, optional): List of regions to keep. If None, no filtering is applied.
            hierarchy_group (str, optional): Name of hierarchy level to group by (e.g., "state").
            normalize_by_r (bool, optional): Normalize values by region/year to compute shares (%).

        Returns:
            pd.DataFrame: Transformed data across all cases from cases_dict.
            str: The identified items column (e.g., 'tech', 'trtype').
            list: List of value columns (e.g., ['Capacity (GW)']).
            pd.DataFrame: Filtered dfmap for the hierarchy group (if provided).
            dict: Color map for the items (from Conventions, if available).
        """
        print(f"(DataFetcher) Transforming data for {data_key}...")

        # Fetch the data if data_key not already fetched
        if data_key not in self.data:
            self.fetch_all(data_keys=data_key)

        # Grab any available case to infer items/value columns
        sample_case = next(iter(self.data[data_key]))
        df_sample = self.data[data_key][sample_case]

        # Infer columns
        items_column = [c for c in df_sample.columns if c not in ['h', 'year', 'r', 'rr']][0]
        value_columns = [
            c for c in df_sample.columns if c not in ['h', 'year', 'r', 'rr', items_column]
        ]

        if items_map_name is not None:
            # If items_column is tech valid options are tech_map1, tech_map2, tech_map3
            if items_column == 'tech' and items_map_name not in ['tech_map1', 'tech_map2', 'tech_map3']:
                raise ValueError(
                    f"Invalid items_map_name '{items_map_name}' for items_column '{items_column}'. "
                    f"Available options: ['tech_map1', 'tech_map2', 'tech_map3']"
                )
            # If items_column is trtype valid options are trtype_map
            elif items_column == 'trtype' and items_map_name != 'trtype_map':
                raise ValueError(
                    f"Invalid items_map_name '{items_map_name}' for items_column '{items_column}'. "
                    f"Available options: ['trtype_map']"
                )
            # If items_column is cost_cat valid options are costcat_map
            elif items_column == 'cost_cat' and items_map_name != 'costcat_map':
                raise ValueError(
                    f"Invalid items_map_name '{items_map_name}' for items_column '{items_column}'. "
                    f"Available options: ['costcat_map']"
                )

            # Get the color map for the items
            color_map = Conventions.items_color_map[items_map_name][1]

        else:
            # Default color map for items
            print(
                f"No items_map_name provided. Using default color map for '{items_column}'."
            )
            color_map = get_default_color_map(df_sample[items_column].unique())

        # Get the dfmap for the hierarchy group
        if hierarchy_group is None:
            dfmap = copy.deepcopy(self.dfmap)
        else:
            dfmap = self.dfmap[hierarchy_group].copy()

        # Transform each case
        transformed_data = {}
        for case in self.data[data_key].keys():
            df = self.data[data_key][case].copy()

            # --- Filter regions ---
            if r_filter is not None:
                # Only have a r column no rr column
                if 'r' in df.columns and 'rr' not in df.columns:
                    df = df[df.r.isin(r_filter)].reset_index(drop=True)

                # Have both r and rr columns. Filter for cases
                # where either r or rr is in the list of regions
                elif 'r' in df.columns and 'rr' in df.columns:
                    df = df[df.r.isin(r_filter) | df.rr.isin(r_filter)].reset_index(drop=True)

                else:
                    # Just skip filtering
                    pass

            # --- Simplify item names ---
            if items_map_name:
                items_map = Conventions.items_color_map[items_map_name][0]
                if items_column == 'tech':
                    df['tech'] = df['tech'].str.lower().map(lambda x: items_map.get(x, x))
                    df = df.groupby(['tech', 'r', 'year']).sum().reset_index()

                    # Fill missing techs with 0
                    idx = pd.MultiIndex.from_product(
                        [sorted(set(items_map.values())), df['r'].unique(), df['year'].unique()],
                        names=['tech', 'r', 'year']
                    )
                    df = df.set_index(['tech', 'r', 'year']).reindex(idx, fill_value=0).reset_index()

                elif items_column == 'trtype':
                    df['trtype'] = df['trtype'].str.lower().map(lambda x: items_map.get(x, x))
                    df = df.groupby(['r', 'rr', 'trtype', 'year']).sum().reset_index()

                    # Fill missing transmission types with 0
                    idx = pd.MultiIndex.from_product(
                        [
                            df['r'].unique(), df['rr'].unique(),
                            sorted(set(items_map.values())), df['year'].unique()
                        ],
                        names=['r', 'rr', 'trtype', 'year']
                    )
                    df = (
                        df.set_index(['r', 'rr', 'trtype', 'year'])
                        .reindex(idx, fill_value=0)
                        .reset_index()
                    )

                elif items_column == 'cost_cat':
                    df['cost_cat'] = df['cost_cat'].str.lower().map(lambda x: items_map.get(x, x))
                    df = df.groupby(['year', 'r', 'cost_cat']).sum().reset_index()
                else:
                    raise ValueError(
                        f"Unknown items_column '{items_column}' for data_key '{data_key}'."
                    )

             # --- Group by hierarchy ---
            if hierarchy_group is not None:
                df = self._group_by_hierarchy(df, items_column, value_columns, hierarchy_group)

            # --- Normalize the data by region ---
            # This is done after hierarchy grouping
            if normalize_by_r:
                df, new_value_columns = self._normalize_by_region(df, items_column, value_columns)

            transformed_data[case] = df

        # Make a df with all cases
        transformed_data = pd.concat(transformed_data, names=["case"])
        transformed_data = transformed_data.reset_index(level=1, drop=True).reset_index()

        # If new_value_columns were created, update the value_columns
        if normalize_by_r:
            value_columns = new_value_columns

        # Oder the cases based on self.prefix_order
        transformed_data['prefix'] = transformed_data["case"].str.split("|").str[0]
        transformed_data["prefix"] = pd.Categorical(
            transformed_data["prefix"],
            categories=self.prefix_order,
            ordered=True
        )
        transformed_data.sort_values(["prefix","case"], inplace=True)
        transformed_data.drop(columns=["prefix"], inplace=True)
        transformed_data.reset_index(drop=True, inplace=True)

        # Replace small negative values (e.g., -1e-6) with 0.0, then round
        threshold = 1e-6
        df = transformed_data[value_columns] 
        df = df.mask((df < 0) & (df.abs() < threshold), 0.0)
        transformed_data[value_columns] = df.round(2)

        return transformed_data, items_column, value_columns, dfmap, color_map

    def _group_by_hierarchy(
        self,
        df: pd.DataFrame,
        items_column: str,
        value_columns: list,
        hierarchy_group: str,
    ) -> pd.DataFrame:
        """
        Group data by a specific hierarchy level.

        Args:
            df (pd.DataFrame): Data to be grouped.
            items_column (str): Column name for the type of data. (tech, trtype, etc.)
            value_columns (list): Columns with values will be summed in the grouping.
            hierarchy_group (str): Hierarchy level for grouping.

        Returns:
            pd.DataFrame: Grouped data.
        """

        # df must contain the columns 'r' and 'year' otherwise do nothing
        if 'r' not in df.columns:
            return df

        # Rename r and rr columns to have the hierarchy group name
        if hierarchy_group != 'cendiv':
            map_r_hierarchy = self.hierarchy[hierarchy_group]
        else:
            # NG case 
            map_r_hierarchy = pd.Series(
                index=self.hierarchy[hierarchy_group].unique(),
                data=self.hierarchy[hierarchy_group].unique(),
                name='cendiv'
            ).rename_axis('r')

        if "year" in df.columns:
            if 'rr' in df.columns:
                df['r'] = df['r'].map(map_r_hierarchy)
                df['rr'] = df['rr'].map(map_r_hierarchy)
                df = df.groupby(['r', 'rr', 'year', 'trtype'])[value_columns].sum()

            else:
                df['r'] = df['r'].map(map_r_hierarchy)
                df = df.groupby(['r', 'year', items_column])[value_columns].sum()

            # Sort region from west to east
            regions = self.dfmap[hierarchy_group].loc[
                self.hierarchy[hierarchy_group].unique()].bounds.minx.sort_values().index
            years = sorted(df.index.get_level_values('year').unique())
            items_type = df.index.get_level_values(items_column).unique()

            # Create a new DataFrame with all combinations of regions and years
            # and fill in the missing values with 0
            if 'rr' in df.index.names:
                index = pd.MultiIndex.from_product(
                    [regions, regions, years, items_type],
                    names=['r', 'rr', 'year', items_column]
                )
            else:
                index = pd.MultiIndex.from_product(
                    [regions, years, items_type],
                    names=['r', 'year', items_column]
                )

            df = df.reindex(index, fill_value=0).reset_index()

        else:
            # If no year column, just map the regions
            df['r'] = df['r'].map(map_r_hierarchy)
            df = df.groupby(['r', items_column])[value_columns].sum()

            regions = self.dfmap[hierarchy_group].loc[
                self.hierarchy[hierarchy_group].unique()].bounds.minx.sort_values().index
            items_type = df.index.get_level_values(items_column).unique()

            # Create a new DataFrame with all combinations of regions and items
            index = pd.MultiIndex.from_product(
                [regions, items_type],
                names=['r', items_column]
            )
            df = df.reindex(index, fill_value=0).reset_index()

        return df

    def _normalize_by_region(
        self,
        df: pd.DataFrame,
        items_column: str,
        value_columns: list,
    ) -> pd.DataFrame:
        """
        Normalize data by region. 
        Create a new column that sums to 100% for each year and region.

        Args:
            df (pd.DataFrame): Data to be normalized.
            items_column (str): Column name for the type of data. (tech, trtype, etc.)
            value_columns (list): Columns for the values to be normalized.

        Returns:
            pd.DataFrame: Normalized data.
        """ 

        # Only allow for capacity and generation
        if 'Capacity (GW)' not in df.columns and 'Generation (TWh)' not in df.columns:
            raise ValueError(
                "Data must contain either 'Capacity' or 'Generation' columns for normalization."
            )

        if 'rr' in df.columns:
            # Can't normalize by region if rr is in the columns
            raise ValueError(
                "Data can't be normalized by region if 'rr' is in the columns."
            )

        # Retain the original region order
        region_order = df['r'].unique()

        # Remove all rows for tech where the sum for at least one year is negative
        mask = (df[value_columns]< -0.01).any(axis=1)
        tech_negative = df[mask][items_column].unique()
        df = df[~df[items_column].isin(tech_negative)]

        # Normalize by region
        df = df.groupby(['year', 'r', items_column])[value_columns].sum()
        df = df / df.groupby(['year','r'])[value_columns].sum() * 100

        df = df.reset_index()

        # Rename value_columns to include percentage notation
        new_cols = [f"{col.split('(')[0].strip()} (%)" for col in value_columns]
        df.rename(columns=dict(zip(value_columns, new_cols)), inplace=True)

        # Sort rows by region order (preserved from original 'r' values)
        df['r'] = pd.Categorical(df['r'], categories=region_order, ordered=True)
        df = df.sort_values(['year', 'r', items_column]).reset_index(drop=True)

        return df, new_cols


#%% ===========================================================================
### --- Plot ReEDS In/Output Functions ---
### ===========================================================================
def hist_last(
    ax: plt.Axes,
    df: pd.DataFrame,
    bins: np.ndarray = None,
    color: str = 'k',
    xpad: int = 1,
    xscale: int = 4,
    alpha: float = 0.5,
) -> None:
    """
    Plot a horizontal histogram using the last row of the DataFrame.

    Args:
        ax (plt.Axes): Axis on which to plot.
        df (pd.DataFrame): Data with year as index.
        bins (np.ndarray, optional): Bin edges for the histogram.
        color (str): Color for the histogram.
        xpad (float): Horizontal offset.
        xscale (float): Factor to scale the bar widths.
        alpha (float): Transparency.

    Returns:
        - fig (plt.Figure): Figure object.
        - data (pd.DataFrame): Data used for plotting.
    """

    if df.empty:
        return  # Nothing to plot

    # x position of the histogram after the last year
    x = df.index[-1] + xpad
    values = df.iloc[-1].values

    if len(values) == 0:
        return  # No valid data

    # Set default bins if not provided or if they are too few
    if bins is None or len(bins) < 2:
        bins = np.linspace(values.min() - 1, values.max() + 1, num=51)
    else:
        # Warning and use default bins if bins not covering the data
        if np.any(values < bins[0]) or np.any(values > bins[-1]):
            print(
                "Warning: Histogram bins do not cover the data. "
                "Using default bins."
            )
            bins = np.linspace(values.min() - 1, values.max() + 1, num=51)

    hist_val, bin_edges = np.histogram(values, bins=bins)

    # Vertical spacing of the bars
    bin_height = bins[1] - bins[0]

    # Width of the bars (frequency normalized)
    freq_width = hist_val / hist_val.max() * xscale

    ax.barh(
        y=bin_edges[:-1],
        width=freq_width,
        height=bin_height,
        align='edge',
        left=x,
        color=color,
        alpha=alpha,
        edgecolor='none'
    )

def plot_time_series_item_r(
    data: pd.DataFrame,
    dfmap: pd.DataFrame,
    items_column: str = 'tech',
    plot_row_keys: list = ['Fossil', 'PV', 'Wind'],
    value_col: str = 'Capacity (GW)',
    highlight_runs: list = None,
    color_map: dict = Conventions.items_color_map['tech_map1'][1],
    interface_plot: bool = False,
    share_ymax: bool = True,
) -> tuple[plt.Figure, np.ndarray, pd.DataFrame]:
    """
    Plot time series for items (e.g., tech) across regions or interfaces.

    This function creates a series of subplots showing sample and base case time series data
    for each specified item (technology) and region or interface. It supports two modes:
    a) interface_plot mode that considers differing region pairs (r vs. rr) and 
    b) non-interface mode where each region is plotted individually.

    Args:
      data (pd.DataFrame): Input DataFrame with columns including 'year', items_column, 'case', and region columns ('r' and possibly 'rr').
      dfmap (pd.DataFrame): DataFrame used to plot region/interface boundary maps.
      items_column (str): Column name containing item labels (default: 'tech').
      plot_row_keys (list): List of items (e.g., technology names) to plot on rows.
      value_col (str): Column name containing the values to plot (default: 'Capacity (GW)').
      highlight_runs (list): List of runs to highlight in the plots: [(case_name, print_name, color),...]
      color_map (dict): Mapping from item names to colors.
      interface_plot (bool): If True, plot as interfaces (using 'r' and 'rr'); otherwise, use individual regions.
      share_ymax (bool): If True, share the y-axis maximum across items.

    Returns:
      tuple: (fig, ax, data) where fig is the matplotlib Figure object,
      ax is the array of Axes objects, and data is the DataFrame used for plotting.
    """

    print("(DataFetcher) Plot using: plot_time_series_item_r...")

    # Filter by year and plot_row_keys
    data = data[(data.year >= 2020) & (data[items_column].isin(plot_row_keys))].copy()

    regions = data.r.unique()
    region_index = {region: i for i, region in enumerate(regions)}

    if not interface_plot:
        if 'rr' in data.columns:
            # Not an interface plot but rr column exists: keep only r == rr
            data = data[data['r'] == data['rr']].drop(columns=['rr'])

        if 'r' in data.columns:
            data.rename(columns={'r': 'region'}, inplace=True)
        else:
            raise ValueError("Data must contain 'r' column for this plot.")

    if interface_plot:
        if 'rr' not in data.columns:
            raise ValueError("Interface plot requires 'rr' column in the data.")

        # Remove all cases where r == rr
        mask = data['r'] != data['rr']
        data = data[mask]

        # Original order of regions
        regions = data.r.unique()
        region_index = {region: i for i, region in enumerate(regions)}

        def _reorder_row(row):
            if region_index[row['r']] > region_index[row['rr']]:
                row['r'], row['rr'] = row['rr'], row['r']
            return row

        data = data.apply(_reorder_row, axis=1)
        data = data.groupby(['case', 'r', 'rr', 'year', items_column])[value_col].sum()
        data = data.reset_index()

        # Define a sorting key based on region_index
        data['r_order'] = data['r'].map(region_index)
        data['rr_order'] = data['rr'].map(region_index)

        # Sort using those keys
        data = data.sort_values(['case', 'r_order', 'rr_order']).reset_index(drop=True)
        data = data.drop(columns=['r_order', 'rr_order'])

        # Unique interfaces (r|rr)
        data['interface'] = data['r'] + '-' + data['rr']

        #If all the data from an interface is zero, remove it
        data = data.groupby('interface').filter(lambda x: (x[value_col] != 0).any())
        data.drop(columns=['r', 'rr'], inplace=True)
        data.rename(columns={'interface': 'region'}, inplace=True)

    # unique regions or interfaces
    unique_col_keys = data['region'].unique()

    # Separate the data on base and sample
    base_identifiers = ["base", "ref", "reference"]
    is_base_case = data["case"].str.split("|").str[0].str.lower().isin(base_identifiers)
    is_highlight = data["case"].isin([c[0] for c in highlight_runs])

    data_base = data[is_base_case]
    data_sample = data[~is_base_case]
    data_highlight = data[is_highlight]

    nrows = len(plot_row_keys) + 1
    ncols = len(unique_col_keys)

    ymax_by_tech = data.groupby(items_column)[value_col].max()
    ymin_by_tech = data.groupby(items_column)[value_col].min()

    ymin = {tech: np.floor(min(ymin_by_tech.get(tech, 0), 0) / 10) * 10 
            for tech in plot_row_keys}
    ymax = {tech: np.ceil(ymax_by_tech.get(tech, 0) / 10) * 10 
            for tech in plot_row_keys}

    if share_ymax:
        # Techs share max
        ymax_shared = np.ceil(ymax_by_tech.max() / 10) * 10
        ystep_tech = {}
        for tech in plot_row_keys:
            ystep_tech[tech] = {
                'major': np.ceil(ymax_shared / 5) if ymin[tech] >= 0 else np.ceil(np.abs(ymin[tech]) / 5),
                'minor': np.ceil(ymax_shared / 10) if ymin[tech] >= 0 else np.ceil(np.abs(ymin[tech]) / 10),
            }
    else:
        # Each technology has its own y-axis
        ystep_tech = {}
        for tech in plot_row_keys:
            ystep_tech[tech] = {
                'major': np.ceil(ymax[tech] / 5) if ymin[tech] >= 0 else np.ceil(np.abs(ymin[tech]) / 5),
                'minor': np.ceil(ymax[tech] / 10) if ymin[tech] >= 0 else np.ceil(np.abs(ymin[tech]) / 10)
            }

    fig, ax = plt.subplots(
        nrows, ncols, sharex='row', sharey='row', figsize=(1.5 * max(2, ncols), (nrows * 1.7 + 0.5)),
        gridspec_kw={'hspace': 0.2, 'height_ratios': [0.2] + [1 / (nrows - 1)] * (nrows - 1)},
        dpi=300,
    )

    # Ensure ax is always a 2D array
    ax = np.atleast_2d(ax)
    if ax.shape != (nrows, ncols):
        ax = ax.reshape((nrows, ncols))

    # Loop through each region or interface column
    for col, col_key in enumerate(unique_col_keys):
        # Plot top row map (e.g., state boundaries)
        dfmap.plot(ax=ax[0,col], facecolor='0.99', edgecolor='0.75', lw=0.2)

        if interface_plot:
            region1, region2 = col_key.split('-')
            dfmap.loc[[region1]].plot(
                ax=ax[0,col], facecolor='#2c3e50' , edgecolor='none')

            dfmap.loc[[region2]].plot(
                ax=ax[0,col], facecolor='#f1c40f', edgecolor='none')
        else:
            dfmap.loc[[col_key]].plot(
                ax=ax[0,col], facecolor='k', edgecolor='none')

        ax[0,col].axis('off')
        ax[0,col].patch.set_facecolor('none')

        # Filter data for the current region/interface.
        df_r_sample = data_sample[data_sample['region'] == col_key]
        df_r_base = data_base[data_base['region'] == col_key]
        df_r_highlight = data_highlight[data_highlight['region'] == col_key]

        # Plot sample cases for each technology.
        for row_index, tech in enumerate(plot_row_keys, start=1):
            df_tech_sample = df_r_sample[df_r_sample[items_column] == tech]
            df_tech_sample = df_tech_sample.pivot(index='year', columns='case', values=value_col)

            if not df_tech_sample.empty:
                # Fill between min and max to create a shaded envelope
                ax[row_index, col].fill_between(
                    df_tech_sample.index, df_tech_sample.max(axis=1), df_tech_sample.min(axis=1),
                    color=color_map[tech], alpha=0.5, lw=0,
                )
                # Plot individual sample curves.
                n_curves = df_tech_sample.shape[1]
                df_tech_sample.plot(
                    ax=ax[row_index, col], legend=False,
                    lw=0.5, color='k', alpha = [.15, .05, .01][int(n_curves > 20) + int(n_curves > 200)]
                )
                hist_last(ax[row_index, col], df_tech_sample, color=color_map[tech])


        # Plot base cases (dashed lines).
        if not data_base.empty:
            for row_index, tech in enumerate(plot_row_keys, start=1):
                df_tech_base = df_r_base[df_r_base[items_column] == tech]
                df_tech_base = df_tech_base.pivot(index='year', columns='case', values=value_col)
                if not df_tech_base.empty:
                    df_tech_base.plot(
                        ax=ax[row_index, col], legend=False,
                        lw=1.5, ls='--', color="k", alpha=0.8,
                    )

        # Highlight specific runs
        if not data_highlight.empty:
            for row_index, tech in enumerate(plot_row_keys, start=1):
                df_tech_highlight = df_r_highlight[df_r_highlight[items_column] == tech]
                df_tech_highlight = df_tech_highlight.pivot(index='year', columns='case', values=value_col)
                if not df_tech_highlight.empty:
                    for case_name, print_name, color in highlight_runs:
                        if case_name in df_tech_highlight.columns:
                            ax[row_index, col].plot(
                                df_tech_highlight.index,
                                df_tech_highlight[case_name],
                                lw=1.5, color=color, label=print_name
                            )

    unit_match = re.search(r'\((.*?)\)', value_col)
    Units = unit_match.group(1) if unit_match else ''
    for row_index, tech in enumerate(plot_row_keys, start=1):
        ax[row_index, 0].set_ylabel(f"{tech} [{Units}]")
        ax[row_index, 0].yaxis.set_minor_locator(mpl.ticker.MultipleLocator(ystep_tech[tech]['minor']))
        ax[row_index, 0].yaxis.set_major_locator(mpl.ticker.MultipleLocator(ystep_tech[tech]['major']))

        if share_ymax:
            ax[row_index, 0].set_ylim(ymin[tech], (0 if ymin[tech] < -0.01 else ymax_shared))

    # Format the x-axis ticks as abbreviated years.
    year_formatter = FuncFormatter(lambda x, pos: f"'{int(x) % 100:02d}")
    for col in range(ncols):
        for row in range(1, nrows):
            ax[row, col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
            ax[row, col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))
            ax[row, col].set_xlabel(None)
            ax[row, col].set_xticklabels([])

        ax[-1, col].xaxis.set_major_formatter(year_formatter)
        ax[-1, col].tick_params(axis='x', rotation=90)

    reeds.plots.despine(ax)

    if interface_plot:
        for col, col_key in enumerate(unique_col_keys):
            ax[1,col].annotate(
                col_key.replace('-', '\n'),
                (0.05, 0.98), xycoords='axes fraction',
                va='top', fontsize=10,
            )

    # Create custom legend handles
    custom_legend = []
    if not data_base.empty:
        custom_legend.append(Line2D([0], [0], color='k', lw=1.5, ls='--', label='Base Case'))

    if not data_highlight.empty:
        for case_name, print_name, color in highlight_runs or []:
            if case_name in data_highlight['case'].unique():
                custom_legend.append(Line2D([0], [0], color=color, lw=1.5, label=print_name))

    if custom_legend:
        fig.subplots_adjust(left=0.08,
                            right=0.95,
                            top=0.95,
                            bottom=0.08)

        fig.legend(
            handles=custom_legend,
            loc='center left',
            bbox_to_anchor=(0.98, 0.5),
            fontsize=10,
            frameon=False,
            ncol=1,
            borderaxespad=0.0,
            borderpad=0.0,
            handletextpad=0.3,
            labelspacing=0.2
        )

    return fig, ax, data


def plot_multi_time_series_matrix(
    data: pd.DataFrame,
    color_map: dict,
    items_column: str = 'tech',
    plot_tech_names: list = ["UPV", "Nuclear"],
    value_cols: list = ['capcost ($/kW)', 'fom ($/kW-yr)'],
    highlight_runs: list = None,
    share_ymax: bool = False,
) -> tuple[plt.Figure, np.ndarray, pd.DataFrame]:
    """
    Plot time series where each column is an item (e.g. technology) 
    and each row an attribute of that item (e.g. capex, opex, etc.).

    Args:
      data (pd.DataFrame): Input DataFrame with columns including 'year', items_column, 'case', and region columns ('r' and possibly 'rr').
      items_column (str): Column name containing item labels (default: 'tech').
      plot_tech_names (list): List of items (e.g., technology names) to plot on rows.
      value_col (list): Column names containing the values to plot
      highlight_runs (list): List of runs to highlight in the plots: [(case_name, print_name, color),...]
      color_map (dict): Mapping from item names to colors.
      share_ymax (bool): If True, share the y-axis maximum across items.

    Returns:
      tuple: (fig, ax, data) where fig is the matplotlib Figure object,
      ax is the array of Axes objects, and data is the DataFrame used for plotting.
    """

    print("(DataFetcher) Plot using: plot_multi_time_series_matrix...")

    # Filter by year, plot_tech_names
    data = data[(data.year >= 2020) & (data[items_column].isin(plot_tech_names))].copy()
    data = data[data[items_column].isin(plot_tech_names)].reset_index(drop=True)

    # Separate the data on base and sample
    base_identifiers = ["base", "ref", "reference"]
    is_base_case = data["case"].str.split("|").str[0].str.lower().isin(base_identifiers)
    is_highlight = data["case"].isin([c[0] for c in highlight_runs])
    data_base = data[is_base_case]
    data_sample = data[~is_base_case]
    data_highlight = data[is_highlight]

    nrows = len(value_cols)
    ncols = len(plot_tech_names)

    fig, ax = plt.subplots(
        nrows,
        ncols,
        sharex='row',
        sharey='row' if share_ymax else False,
        figsize=(1.5 * max(2, ncols), nrows * 1.7),
        gridspec_kw={
            'hspace': 0.2,
            'wspace': 0.6 if not share_ymax else 0.2, 
            'height_ratios': [1 / (nrows - 1)] * (nrows),
        },
        dpi=300,
    )

    # Ensure ax is always a 2D array
    ax = np.atleast_2d(ax)
    if ax.shape != (nrows, ncols):
        ax = ax.reshape((nrows, ncols))

    # Plot sample cases for each technology
    for col_i, tech in enumerate(plot_tech_names):
        # Loop through each value_col
        for row_i, row_key in enumerate(value_cols):

            # Plot sample
            df_tech_sample = data_sample[data_sample[items_column] == tech]
            df_tech_sample = df_tech_sample.pivot(index='year', columns='case', values=row_key)
            if not df_tech_sample.empty:
                # Fill between min and max to create a shaded envelope
                ax[row_i, col_i].fill_between(
                    df_tech_sample.index, df_tech_sample.max(axis=1), df_tech_sample.min(axis=1),
                    color=color_map[tech], alpha=0.5, lw=0,
                )
                # Plot individual sample curves.
                n_curves = df_tech_sample.shape[1]
                df_tech_sample.plot(
                    ax=ax[row_i, col_i], legend=False,
                    lw=0.5, color='k', alpha = [.15, .05, .01][int(n_curves > 20) + int(n_curves > 200)]
                )
                hist_last(ax[row_i, col_i], df_tech_sample, color=color_map[tech])

            # Plot base cases (dashed lines).
            df_tech_base = data_base[data_base[items_column] == tech]
            df_tech_base = df_tech_base.pivot(index='year', columns='case', values=row_key)
            if not df_tech_base.empty:
                df_tech_base.plot(
                    ax=ax[row_i, col_i], legend=False,
                    lw=1.5, ls='--', color="k", alpha=0.8,
                )

            # Highlight specific runs
            if highlight_runs:
                df_tech_highlight = data_highlight[data_highlight[items_column] == tech]
                df_tech_highlight = df_tech_highlight.pivot(index='year', columns='case', values=row_key)
                if not df_tech_highlight.empty:
                    for case_name, print_name, color in highlight_runs:
                        if case_name in df_tech_highlight.columns:
                            ax[row_i, col_i].plot(
                                df_tech_highlight.index,
                                df_tech_highlight[case_name],
                                lw=1.5, color=color, label=print_name
                            )

    for row_i, value_name in enumerate(value_cols):
        row_label = value_name.split('(')[0].strip()

        if len(value_name.split('(')) > 1:
            row_units = value_name.split('(')[1].split(')')[0].strip()
            ax[row_i, 0].set_ylabel(f"{row_label}\n[{row_units}]")
        else:
            ax[row_i, 0].set_ylabel(row_label)

    # Format the x-axis ticks as abbreviated years.
    year_formatter = FuncFormatter(lambda x, pos: f"'{int(x) % 100:02d}")
    for col in range(ncols):
        for row in range(nrows):
            ax[row, col].xaxis.set_major_locator(mpl.ticker.MultipleLocator(10))
            ax[row, col].xaxis.set_minor_locator(mpl.ticker.MultipleLocator(5))
            ax[row, col].set_xlabel(None)
            ax[row, col].set_xticklabels([])

        ax[0, col].set_title(plot_tech_names[col], fontsize=12, fontweight='bold')
        ax[-1, col].xaxis.set_major_formatter(year_formatter)
        ax[-1, col].tick_params(axis='x', rotation=90)

    reeds.plots.despine(ax)

    # Create custom legend handles
    custom_legend = []
    if not data_base.empty:
        custom_legend.append(Line2D([0], [0], color='k', lw=1.5, ls='--', label='Base Case'))

    if not data_highlight.empty:
        for case_name, print_name, color in highlight_runs or []:
            if case_name in data_highlight['case'].unique():
                custom_legend.append(Line2D([0], [0], color=color, lw=1.5, label=print_name))


    if custom_legend:
        fig.subplots_adjust(left=0.08,
                            right=0.95,
                            top=0.95,
                            bottom=0.08)

        fig.legend(
            handles=custom_legend,
            loc='center left',
            bbox_to_anchor=(0.98, 0.5),
            fontsize=10,
            frameon=False,
            ncol=1,
            borderaxespad=0.0,
            borderpad=0.0,
            handletextpad=0.3,
            labelspacing=0.2
        )

    return fig, ax, data


def add_raincloud_to_ax(
    ax: plt.Axes,
    data: pd.DataFrame,
    x_column: str,
    y_column: str,
    color_map: dict = None,
    mark_base_case_ranges: bool = False,
    highlight_runs: list = None,
) -> None:
    """
    Add RainCloud-style plots to an axis, combining violin, jitter, and box plots.
    Each unique element in the x_column will have its own violin plot,
    with jittered scatter points and box plots overlaid.

    Args:
        ax (plt.Axes): Axis to plot on.
        data (pd.DataFrame): Data to plot.
        x_column (str):  Column from data for x-axis.
        y_column (str): Column from data for y-axis.  
        color_map (dict, optional): Color mapping for the x_column values.
        mark_base_case_ranges (bool, optional): If True, mark base case ranges 
        highlight_runs (list, optional): List of runs to highlight in the plots: [(case_name, print_name, color),...]
    """
    if color_map is None:
        color_map = get_default_color_map(data[x_column].unique())

    # Make a plot to mimic the data distribution (viloin plot)
    sns.violinplot(
        data=data,
        x=x_column,
        y=y_column,
        hue=x_column,
        legend=False, 
        inner=None,
        density_norm="width",
        palette=color_map,
        linewidth=1.5,
        cut=0,
        ax=ax,
        alpha=0.7
    )

    # Remove the left-most violin plot
    for artist in ax.collections:
        path = artist.get_paths()[0]
        vertices = path.vertices
        vertices[:, 0] = np.clip(vertices[:, 0], vertices[:, 0].mean(), None)

    # Get tick locations and labels from Seaborn
    tick_locs = ax.get_xticks()
    tick_labels = [tick.get_text() for tick in ax.get_xticklabels()]

    jitter_strength = 0.3
    box_offset = 0.2

    # plot one group at a time
    for i, label in enumerate(tick_labels):
        group_mask = data[x_column].astype(str) == label
        group_vals = data.loc[group_mask, y_column].dropna()

        # single jitter array for the whole group
        jitter = np.random.uniform(0, jitter_strength, size=len(group_vals))
        x_vals = np.full_like(group_vals, tick_locs[i], dtype=float) + jitter
        n_curves = group_vals.count()

        # base scatter (all runs, faint)
        ax.scatter(
            x_vals,
            group_vals,
            s=10,
            facecolors="none",
            edgecolors="black",
            linewidths=1,
            alpha = [.3, .05, .01][int(n_curves > 20) + int(n_curves > 200)]
        )

        # recolor points that belong to any highlight run
        if highlight_runs:
            for case_name, print_name, color in highlight_runs:
                case_mask = (data["case"] == case_name) & group_mask
                hl_vals = data.loc[case_mask, y_column].dropna()
                if hl_vals.empty:
                    continue

                # positions of these rows inside group_vals
                idx_in_group = group_vals.index.get_indexer(hl_vals.index)
                ax.scatter(
                    x_vals[idx_in_group],
                    hl_vals,
                    s=20,
                    color=color,
                    alpha=0.8,
                    label=print_name,
                )

        # box layer
        ax.boxplot(
            group_vals,
            positions=[tick_locs[i] - box_offset],
            widths=0.2,
            patch_artist=True,
            boxprops=dict(facecolor="none", color="black", linewidth=1.5),
            medianprops=dict(color="black", linewidth=1.5),
            whiskerprops=dict(color="black", linewidth=1.5),
            capprops=dict(color="black", linewidth=1.5),
            flierprops=dict(markeredgecolor="black", markersize=4),
        )

    ax.set_xticks(tick_locs)
    ax.set_xticklabels(tick_labels, rotation=90, ha="center", fontsize=12)
    ax.grid(True, linestyle="--", linewidth=0.5)
    ax.set_xlabel("")
    ax.set_ylabel("")

    # option to draw base case ranges
    if mark_base_case_ranges:
        data_base = data[data["case"].str.split("|").str[0].str.lower() == "base"]
        for i, label in enumerate(tick_labels):
            y_vals = data_base[data_base[x_column].astype(str) == label][y_column].dropna()
            for y in y_vals:
                ax.plot(
                    [tick_locs[i] - 0.4, tick_locs[i] + 0.4],
                    [y, y],
                    color="black",
                    linewidth=0.7,
                )

def out_plot_raincloud(
    data: pd.DataFrame,
    x_axis: str,
    y_axis: str,
    items_column: str,
    plot_items: list,
    mark_base_case_ranges: bool = False,
    highlight_runs: list = None,
    color_map: dict = None,
    share_ymax: bool = True,
    figsize: tuple = None,
) -> tuple[plt.Figure, np.ndarray, pd.DataFrame]:
    """
    Create a raincloud plot layout from the given data.

    Depending on x_axis, plots are arranged as follows:
      - If x_axis is "year": each prefix is plotted in a separate column and each 
        item (from plot_items) is plotted in a separate row.
      - If x_axis is "prefix": a single row is used, with each item as a separate column.
      - If x_axis it not "year" or "prefix", a single column is used, with each item as a separate row.


    Args:
        data (pd.DataFrame): Data frame that must contain a 'case' column and depending
            on x_axis, a 'year' column.
        x_axis (str): 'year', 'prefix', or other column name used for x-axis values.
        y_axis (str): Column name used for the y-axis values.
        items_column (str): Column used to select specific items that will be plotted.
        plot_items (list): List of items to plot.
        mark_base_case_ranges (bool, optional): If True, mark base case ranges.
        highlight_runs (list, optional): List of runs to highlight in the plots: [(case_name, print_name, color),...]
        color_map (dict, optional): A mapping for colors. For x_axis 'prefix', if none
            is provided, a default will be created.
        share_ymax (bool, optional): Whether to share y-axis scales across rows.
        figsize (tuple, optional): Size of the figure. If None, it will be calculated

    Returns:
        tuple: (fig, ax, data) 
    """
    # Validate data structure based on x_axis
    if x_axis == "year":
        if "year" not in data.columns:
            raise ValueError("Data must contain a 'year' column for x_axis='year'.")

        # Check if we have unique combinations of case, items_col, and year
        if data.groupby(["case", items_column, "year"]).size().max() > 1:
            raise ValueError("Data must have unique combinations of case, items_col, and year.")

    # Get the prefix of each case
    data["prefix"] = data.case.str.split("|").str[0]
    unique_prefixes = data["prefix"].unique()

    if x_axis == "prefix" and color_map is None:
        # If x_axis is prefix and no color_map is provided, create a default one
        color_map = get_default_color_map(data['prefix'].unique())

    # Determine layout dimensions
    if x_axis == "year":
        ncols = len(unique_prefixes)
        nrows = len(plot_items)
    elif x_axis == "prefix":
        ncols = len(plot_items)
        nrows = 1
    else:
        ncols = 1
        nrows = len(plot_items)

    if figsize is None:
        figsize=(ncols*5, nrows*(3 if x_axis == "year" else 5))

    fig, ax = plt.subplots(
        nrows=nrows, ncols=ncols, figsize=figsize,
        sharex="row", sharey="row" if share_ymax else False,
        gridspec_kw={'hspace': 0.2, 'wspace': 0.2},
        dpi=300,
    )

    # Ensure ax is always a 2D array
    ax = np.atleast_2d(ax)
    if ax.shape != (nrows, ncols):
        ax = ax.reshape((nrows, ncols))

    if x_axis == "year":
        for col, prefix in enumerate(unique_prefixes):
            # Filter data for the current prefix
            df_prefix = data[data["prefix"] == prefix]

            # Loop through each item
            for row, item in enumerate(plot_items):
                # Filter data for the current item
                df_item = df_prefix[df_prefix[items_column] == item]
                if df_item.empty:
                    continue

                # Create a temporary color map applying the same color to all years       
                cmap_tmp = {year: color_map.get(item) 
                    for year in df_item.year.unique()
                } 

                add_raincloud_to_ax(
                    ax=ax[row, col],
                    data=df_item,
                    x_column=x_axis,
                    y_column=y_axis,
                    color_map=cmap_tmp,
                    mark_base_case_ranges=mark_base_case_ranges,
                    highlight_runs=highlight_runs,
                )
    elif x_axis == "prefix":
        # If x_axis is prefix, we have a single row
        # and each item in a different column
        for col, item in enumerate(plot_items):
            # Filter data for the current item
            df_item = data[data[items_column] == item]
            if df_item.empty:
                continue

            add_raincloud_to_ax(
                ax=ax[0, col],
                data=df_item,
                x_column=x_axis,
                y_column=y_axis,
                color_map=color_map,
                mark_base_case_ranges=False,
                highlight_runs=highlight_runs,
            )
    else:
        for row, item in enumerate(plot_items):
            # Filter data for the current item
            df_item = data[data[items_column] == item]
            if df_item.empty:
                continue

            add_raincloud_to_ax(
                ax=ax[row, 0],
                data=df_item,
                x_column=x_axis,
                y_column=y_axis,
                color_map=color_map,
                mark_base_case_ranges=mark_base_case_ranges,
                highlight_runs=highlight_runs,
            )

    # Organize titles and labels
    if x_axis == "year":
        # Set titles at the top of each prefix column
        for col, prefix in enumerate(unique_prefixes):
            ax[0, col].set_title(prefix, fontsize=12, weight='bold')
            # Remove x-tick labels for all but the bottom row
            for row in range(nrows - 1):
                ax[row, col].set_xticklabels([])
            ax[-1, col].tick_params(axis='x', rotation=90) 

        # Set y-axis labels for each row representing an item
        for row, item in enumerate(plot_items):
            unit_match = re.search(r'\((.*?)\)', y_axis)
            units = unit_match.group(1) if unit_match else ''
            ax[row, 0].set_ylabel(f"{item} [{units}]")

    elif x_axis == "prefix":
        # For x_axis "prefix": use y-axis label only for the first subplot
        ax[0, 0].set_ylabel(y_axis.replace(" (", " [").replace(")", "]"))
        # Set a bold title for each item column
        for col, item in enumerate(plot_items):
            ax[0, col].set_title(item, fontsize=12, weight='bold')

    else:
        for row, item in enumerate(plot_items):
            unit_match = re.search(r'\((.*?)\)', y_axis)
            units = unit_match.group(1) if unit_match else ''
            ax[row, 0].set_ylabel(f"{item} [{units}]")

    # Remove top and right frames and set tick directions
    for a in ax.flat:
        a.tick_params(
            axis='both',
            direction='out',  # Ticks point outside the plot box
            bottom=True, top=False,
            left=True, right=False,
            labelbottom=True,
            labelleft=True
        )
        a.spines['top'].set_visible(False)
        a.spines['right'].set_visible(False)

    if share_ymax:
        # Collect all y-limits from every subplot
        all_ymins = []
        all_ymaxs = []

        for a in ax.flat:
            ymin, ymax = a.get_ylim()
            all_ymins.append(ymin)
            all_ymaxs.append(ymax)

        # Compute global min and max
        global_ymin = min(all_ymins)
        global_ymax = max(all_ymaxs)

        # Set the same limits on all axes
        for a in ax.flat:
            a.set_ylim(global_ymin, global_ymax)

    # Plot legend for base and highlight runs if applicable
    custom_legend = []
    if mark_base_case_ranges and data['case'].str.contains(r'base|ref|reference', case=False, regex=True).any():
        custom_legend.append(Line2D([0], [0], color='black', marker='o', lw=0, label='Base Case'))
    if highlight_runs:
        for case_name, print_name, color in highlight_runs:
            custom_legend.append(Line2D([0], [0], color=color, marker='o', lw=0, label=print_name))

    if custom_legend:
        fig.subplots_adjust(left=0.08,
                            right=0.95,
                            top=0.95,
                            bottom=0.08)

        fig.legend(
            handles=custom_legend,
            loc='center left',
            bbox_to_anchor=(0.98, 0.5),
            fontsize=12,
            frameon=False,
            ncol=1,
            borderaxespad=0.0,
            borderpad=0.0,
            handletextpad=0.3,
            labelspacing=0.2
        )

    return fig, ax, data


#%% ===========================================================================
### --- Plot ReEDS Input/Output Class ---
### ===========================================================================
class DataPlotter:
    """
        Creates plots for the input/output data fetched by 
        a DataFetcher instance and manages summary reports 
        which are saved to a PowerPoint presentation.

        Args:
            data_fetcher (DataFetcher): Instance of DataFetcher to fetch data.
            highlight_runs (list, optional): List of tuples with (run_name, print_name, color)
                to highlight specific runs in plots.
            save_path (str, optional): Path to save the PowerPoint presentation
                and a excel files with the figures data.
                If None, a folder at os.path.join(reeds_path, 'runs', 'comparisons')
                will be created.
        Attributes:
            figures_data (dict(dict)):
                Dictionary to store metadata for each figure, including title and axes.
            prs (pptx.Presentation): PowerPoint presentation object.
            plot_methods (dict): Dictionary mapping plot types to methods.
    """
    def __init__(
        self,
        data_fetcher: DataFetcher,
        highlight_runs: list = [],
        save_path: str = None,
    ):
        self.data_fetcher = data_fetcher

        # Dictionary to store metadata for each figure, including title and axes
        self.figures_data = {}
        self.highlight_runs = highlight_runs
        self.save_path = save_path

        # Load PowerPoint template
        self.prs = reeds.report_utils.init_pptx()  

        self.input_plot_methods ={
            "demand_year_r": self.plot_demand_year_r,
            "tech_inputs_year": self.plot_tech_inputs_year,
            "supply_curves_cap": self.plot_supply_capacity,
            "ng_fuel_prices": self.plot_ng_fuel_prices,
        }

        self.output_plot_methods = {
            "capacity_year_tech_r": self.plot_capacity_year_tech_r,
            "generation_year_tech_r": self.plot_generation_year_tech_r,
            "gen_penetration_tech_year_r": self.plot_gen_penetration_tech_year_r,
            "transmission_year_tw_mi_r": self.plot_transmission_year_tw_mi_r,
            "transmission_year_tw_mi_interface": self.plot_transmission_year_tw_mi_interface,
            "transmission_year_cap_interface": self.plot_transmission_year_cap_interface,
            "emissions_year_r": self.plot_emissions_r,
            "system_npv_prefix_comparison": self.plot_system_npv_prefix_comparison,
            "tech_cap_prefix_comp_year": self.plot_tech_cap_prefix_comp_year,
            "tech_gen_prefix_comp_year": self.plot_tech_gen_prefix_comp_year,
            "tech_cap_lastyear_prefix_comp": self.plot_tech_cap_lastyear_prefix_comp,
            "tech_gen_lastyear_prefix_comp": self.plot_tech_gen_lastyear_prefix_comp,
        }

        self.plot_methods = {
            **self.input_plot_methods,
            **self.output_plot_methods
        }

    ### =========================================================================
    ### --- Helpers 
    ### =========================================================================
    def create_all_plots(self, methods_list: str | list = None):
        """
        Create plots for all cases in methods_list, store the figures in the class,
        and also save them to a PowerPoint presentation.

        Args:
            methods_list (str or list): 
                The type of plots to create (keys from self.plot_methods).
                if None, all plots will be created.
        """
        if methods_list is None:
            methods_list = list(self.plot_methods.keys())
        elif isinstance(methods_list, str):
            methods_list = [methods_list]

        for plot_type in methods_list:
            plot_func = self.plot_methods.get(plot_type)
            if not plot_func:
                print(f"[WARNING] Unknown plot type '{plot_type}'. Skipping.")
                continue
            try:
                print(f"\n>> (DataPlotter) Get Plot Using: {plot_type}")
                plot_func()
            except Exception as e:
                print(f"[ERROR] Failed to generate '{plot_type}': {e}")

        self.create_summary_pptx()


    def create_summary_pptx(self):
        """Export all stored figures to a PowerPoint presentation."""

        self.save_path = self._create_ppt_folder()
        print(f'\n>> Saving pptx to {self.save_path}')

        for fig_key, fig_info in self.figures_data.items():
            title = fig_info.get("title", fig_key)
            fig = fig_info.get("fig")

            if fig is None:
                print(f"[WARNING] No figure found for '{fig_key}'. Skipping.")
                continue

            print(f"Adding '{title}' to the PowerPoint")
            self._add_to_pptx(fig, title=title)

        pptx_path = os.path.join(self.save_path, 'summary_report.pptx')
        self.prs.save(pptx_path)

        print(f">> Done: {pptx_path}")

        # Create an Excel file with all the data from all figs
        self._save_to_excel()


    def _save_to_excel(self):
        """
        Create an Excel file with all the data from all figs.
        """
        if self.save_path is None:
            self.save_path = self._create_ppt_folder()

        print(f'\n>> Saving fig data to {self.save_path}')
        excel_path = os.path.join(self.save_path, 'figures_data.xlsx')
        with pd.ExcelWriter(excel_path, engine='xlsxwriter') as writer:
            for fig_key, fig_info in self.figures_data.items():
                title = fig_info.get("title", fig_key)
                df = fig_info.get("data")

                if df is None:
                    print(f"[WARNING] No data found for '{fig_key}'. Skipping.")
                    continue

                # Write the DataFrame to the specified sheet.
                df.to_excel(writer, sheet_name=title[:31], index=False)

        print(f">> Done: {excel_path}")


    def _add_to_pptx(self, fig: plt.Figure, title: str,
        left: float = None, top: float = None,
        width: float = None, height: float = None,
        center: bool = True,
    ) -> None:
        """
        Add a figure to a new slide in the PowerPoint.

        Args:
            fig (plt.Figure): Matplotlib figure.
            title (str, optional): Slide title.
            left, top, width, height (float): Dimensions in inches.
            center (bool): If True, center the figure on the slide.
        """
        # Set default dimensions
        max_width = reeds.report_utils.SLIDE_WIDTH
        max_height = reeds.report_utils.SLIDE_HEIGHT

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[3])

        if title:
            slide.shapes.title.text = title

        # Use figure size unless dimensions are passed
        fig_width, fig_height = fig.get_size_inches()

        # Determine max scaling ratio that fits both width and height
        scale_w = max_width / fig_width
        scale_h = max_height / fig_height
        scale = min(scale_w, scale_h, 1.0)  # Never upscale beyond 100%

        # Apply user overrides if passed
        final_width = width if width is not None else fig_width * scale
        final_height = height if height is not None else fig_height * scale

        width = Inches(final_width)
        height = Inches(final_height)

        if center:
            left = (max_width - final_width) / 2
            top = 0.62 + (max_height - final_height) / 2


        # Save figure to in-memory PNG
        img_stream = io.BytesIO()
        fig.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
        img_stream.seek(0)

        slide.shapes.add_picture(
            img_stream,
            left=Inches(left),
            top=Inches(top),
            width=width,
            height=height,
        )

    def _create_ppt_folder(self):
        """
        Create an output folder for the summary PowerPoint report and 
        return any changes to the path.

        Args:
            self.save_path: Base path to save results. 
            If None, uses a default based on {reeds_path}/runs/comparisons.

        Returns:
            str: Final path to the output folder.
        """

        datetime_str = pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')

        if self.save_path is None:
            outpath = os.path.join(reeds_path, 'runs', 'comparisons')
        else:
            outpath = self.save_path

        if os.path.exists(outpath):
            print(f"\nFolder {outpath} already exists. Renaming it with a timestamp.")
            outpath = os.path.join(outpath, datetime_str)

        os.makedirs(outpath, exist_ok=True)
        return outpath


    ### =========================================================================
    ### --- Input Plotting Functions
    ### =========================================================================
    def plot_demand_year_r(self):
        """
        Plot demand data by region.
        """
        data_key = "demand"
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['Total Load']

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map,
        )

        # Store figure in the class
        self.figures_data['demand_year_r'] = {
            "title": "Inputs - Demand by Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_tech_inputs_year(self):
        """
        Plot technology inputs by region.
        """
        data_key = "tech_inputs"

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
        )
        data = data.fillna(0)

        # ---- First plot
        plot_tech_names = ['UPV', 'On. wind', 'Off. wind fixed', 'battery_li']
        fig, ax, data_save = plot_multi_time_series_matrix(
            data=data,
            color_map=color_map,
            items_column = items_column,
            plot_tech_names = plot_tech_names,
            value_cols = ['capcost ($/kW)', 'fom ($/kW-yr)'],
            highlight_runs=self.highlight_runs,
        )

        # Store figure in the class
        self.figures_data['tech_inputs_P1'] = {
            "title": "Inputs - Technology Inputs - P1",
            "fig": fig,
            "ax": ax,
            "data": data_save,
            "adjusts": {}, # Adjustments for the pptx
        }

        # ---- Second plot 
        plot_tech_names = ['Nuclear', 'Nuclear-SMR', 'Gas-CC', 'Gas-CT', 'Coal-new']
        fig, ax, data_save = plot_multi_time_series_matrix(
            data=data,
            color_map=color_map,
            items_column = items_column,
            plot_tech_names = plot_tech_names,
            value_cols = ['capcost ($/kW)', 'fom ($/kW-yr)','vom ($/MWh)', 'heatrate (MMBtu/MWh)'],
            highlight_runs=self.highlight_runs,
        )

        # Store figure in the class
        self.figures_data['tech_inputs_P2'] = {
            "title": "Inputs - Technology Inputs - P2",
            "fig": fig,
            "ax": ax,
            "data": data_save,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_supply_capacity(self):
        """
        Plot technology supply curve capacity by region,
        """
        data_key = 'supply_curves_capacity'
        hierarchy_group = 'st'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            hierarchy_group=hierarchy_group,
        )
        # Sort data using the 'r' column
        data = data.sort_values(by=['r'])

        fig, ax, data_save = out_plot_raincloud(
            data, 
            x_axis = 'r',
            y_axis = value_columns[0],
            items_column = items_column,
            plot_items = ['UPV',  'Wind On', 'Wind Off'],
            mark_base_case_ranges = True,
            highlight_runs = self.highlight_runs,
            color_map = None,
            share_ymax = False,
            figsize = (
                min(20, max(2, len(data["r"].unique()))*1.5),
                8 if len(data["r"].unique()) < 5 else 10
            )
        )

        # Store figure in the class
        self.figures_data['supply_curves'] = {
            "title": "Inputs - Avail. Capacity by r",
            "fig": fig,
            "ax": ax,
            "data": data_save,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_ng_fuel_prices(self):
        """
        Plot natural gas fuel prices by region.
        """
        data_key = "ng_fuel_prices"
        hierarchy_group = 'cendiv'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['NG Price']

        fig, ax, data_save = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['ng_fuel_prices'] = {
            "title": "Inputs - Natural Gas Fuel Prices",
            "fig": fig,
            "ax": ax,
            "data": data_save,
            "adjusts": {}, # Adjustments for the pptx
        }

    ### =========================================================================
    ### --- Output Plotting Functions
    ### =========================================================================
    def plot_capacity_year_tech_r(self):
        """
        Plot capacity data by technology and region.
        """
        data_key = "cap_out"
        items_map_name = 'tech_map1'
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['capacity_year_tech_r'] = {
            "title": "Capacity by Technology and Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_generation_year_tech_r(self):
        """
        Plot generation data by technology and region.
        """
        data_key = "gen_ann"
        items_map_name = 'tech_map1'
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['generation_year_tech_r'] = {
            "title": "Generation by Technology and Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_gen_penetration_tech_year_r(self):
        """
        Plot generation penetration by technology and region.
        """
        data_key = "gen_ann"
        items_map_name = 'tech_map1'
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
            normalize_by_r=True,
        )

        plot_row_keys = ['PV', 'Wind','Fossil', 'Nuclear']

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['gen_penetration_tech_year_r'] = {
            "title": "Generation Penetration by Technology and Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_transmission_year_tw_mi_r(self):
        """
        Plot transmission data by technology and region.
        """
        data_key = "tran_mi_out_detail"
        items_map_name = 'trtype_map'

        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['AC', 'Reinf', 'Spur']
        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=False,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['transmission_year_tw_mi_r'] = {
            "title": "Transmission TW-Mile by Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_transmission_year_tw_mi_interface(self):
        """
        Plot transmission data by technology and interface.
        """
        data_key = "tran_mi_out_detail"
        items_map_name = 'trtype_map'
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        plot_row_keys = (
            ['AC', 'B2B', 'DC/LCC'] if region_prefix in ['country', 'interconnect']
            else  ['AC']
        )

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=True,
            share_ymax=False,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['transmission_year_tw_mi_interface'] = {
            "title": "Transmission TW-Mile by Inteface",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_transmission_year_cap_interface(self):
        """
        Plot transmission capacity by technology and interface.
        """
        data_key = "tran_cap_energy"
        items_map_name = 'trtype_map'
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        plot_row_keys = (
            ['AC', 'B2B', 'DC/LCC'] if region_prefix in ['country', 'interconnect']
            else  ['AC']
        )

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=True,
            share_ymax=False,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['transmission_year_cap_interface'] = {
            "title": "Transmission Capacity by Inteface",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_emissions_r(self):
        """
        Plot emissions data by region.
        """
        data_key = "emit_r"
        region_prefix = self.data_fetcher.sw.GSw_Region.split("/")[0]
        hierarchy_group = 'st' if region_prefix == 'st' else 'transreg'

        # Get the raw data
        _ = self.data_fetcher.fetch_all(data_keys=data_key)

        # Apply transformations
        data, items_column, value_columns, dfmap, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            hierarchy_group=hierarchy_group,
        )

        plot_row_keys = ['CO2e']
        fig, ax, data = plot_time_series_item_r(
            data=data,
            dfmap=dfmap,
            items_column=items_column,
            value_col=value_columns[0],
            highlight_runs=self.highlight_runs,
            plot_row_keys=plot_row_keys,
            interface_plot=False,
            share_ymax=True,
            color_map=color_map
        )

        # Store figure in the class
        self.figures_data['emissions_year_r'] = {
            "title": "Emissions by Region",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_system_npv_prefix_comparison(self):
        """
        Plot system NPV comparison by prefix.
        """
        data_key = "npv_r"
        hierarchy_group = 'country'
        plot_items = ["Total Cost", "Gen & Stor Capex", "Gen & Stor O&M", "Fuel"]
        x_axis="prefix"
        y_axis="Discounted Cost (Bil $)"
        share_ymax=False
        color_map = None

        # Apply transformations
        data, items_column, _, _, _ = self.data_fetcher.apply_transformations(
            data_key=data_key,
            hierarchy_group=hierarchy_group,
            items_map_name='costcat_map',
        )

        # Create a Total Cost column
        data = data.groupby(["case", items_column]).sum()[y_axis].reset_index()
        df_tot_cost = data.groupby(["case"]).sum()[y_axis].reset_index()
        df_tot_cost[items_column] = "Total Cost"
        data = pd.concat([data, df_tot_cost], axis=0) 

        # Sort the data by prefix
        data["prefix"] = data.case.str.split("|").str[0]
        data["prefix"] = pd.Categorical(
            data["prefix"],
            categories=self.data_fetcher.prefix_order,
            ordered=True,
        )
        data.sort_values(["prefix", items_column], inplace=True)
        data.drop(columns=["prefix"], inplace=True)

        fig, ax, data = out_plot_raincloud(
            data=data,
            x_axis=x_axis,
            y_axis=y_axis,
            items_column=items_column,
            plot_items=plot_items,
            highlight_runs = self.highlight_runs,
            color_map=color_map,
            share_ymax=share_ymax,
        )

        # Store figure in the class
        self.figures_data['system_npv_prefix_comparison'] = {
            "title": "System NPV Comparison by Prefix",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_tech_cap_lastyear_prefix_comp(self):
        """
        Plot technology capacity in the last year of the simulation.
        """
        data_key = "cap_out"
        items_map_name = 'tech_map1'
        hierarchy_group = 'country'
        plot_items = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']
        x_axis="prefix"
        y_axis="Capacity (GW)"
        color_map = None

        # Apply transformations
        data, items_column, _, _, _ = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        year = data.year.max()
        # Filter data for the last year
        data = data[data.year == year].drop(columns=["year"])

        fig, ax, data = out_plot_raincloud(
            data=data,
            x_axis=x_axis,
            y_axis=y_axis,
            items_column=items_column,
            plot_items=plot_items,
            highlight_runs = self.highlight_runs,
            color_map=color_map,
            share_ymax=True,
        )

        # Store figure in the class
        self.figures_data['tech_cap_lastyear_prefix_comp'] = {
            "title": f"Technology Capacity in {year}",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_tech_gen_lastyear_prefix_comp(self):
        """
        Plot technology generation in the last year of the simulation.
        """
        data_key = "gen_ann"
        items_map_name = 'tech_map1'
        hierarchy_group = 'country'
        plot_items = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']
        x_axis="prefix"
        y_axis="Generation (TWh)"
        color_map = None

        # Apply transformations
        data, items_column, _, _, _ = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        year = data.year.max()
        # Filter data for the last year
        data = data[data.year == year].drop(columns=["year"])

        fig, ax, data = out_plot_raincloud(
            data=data,
            x_axis=x_axis,
            y_axis=y_axis,
            items_column=items_column,
            plot_items=plot_items,
            highlight_runs = self.highlight_runs,
            color_map=color_map,
            share_ymax=True,
        )

        # Store figure in the class
        self.figures_data['tech_gen_lastyear_prefix_comp'] = {
            "title": f"Technology Generation in {year}",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_tech_cap_prefix_comp_year(self):
        """
        Plot technology capacity for all years (prefix comparison)
        """
        data_key = "cap_out"
        items_map_name = 'tech_map1'
        hierarchy_group = 'country'
        plot_items = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']
        x_axis="year"
        y_axis="Capacity (GW)"

        # Apply transformations
        data, items_column, _, _, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        # Plot for only after 2020
        data = data[data.year >= 2020]

        fig, ax, data = out_plot_raincloud(
            data=data,
            x_axis=x_axis,
            y_axis=y_axis,
            items_column=items_column,
            plot_items=plot_items,
            highlight_runs = self.highlight_runs,
            color_map=color_map,
            share_ymax=True,
        )

        # Store figure in the class
        self.figures_data['tech_cap_prefix_comp_year'] = {
            "title": "Technology Capacity by Year",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }

    def plot_tech_gen_prefix_comp_year(self):
        """
        Plot technology generation for all years (prefix comparison)
        """
        data_key = "gen_ann"
        items_map_name = 'tech_map1'
        hierarchy_group = 'country'
        plot_items = ['PV', 'Wind','Fossil', 'Nuclear', 'Storage']
        x_axis="year"
        y_axis="Generation (TWh)"

        # Apply transformations
        data, items_column, _, _, color_map = self.data_fetcher.apply_transformations(
            data_key=data_key,
            items_map_name=items_map_name,
            hierarchy_group=hierarchy_group,
        )

        # Plot for only after 2020
        data = data[data.year >= 2020]

        fig, ax, data = out_plot_raincloud(
            data=data,
            x_axis=x_axis,
            y_axis=y_axis,
            items_column=items_column,
            plot_items=plot_items,
            highlight_runs = self.highlight_runs,
            color_map=color_map,
            share_ymax=True
        )

        # Store figure in the class
        self.figures_data['tech_gen_prefix_comp_year'] = {
            "title": "Technology Generation by Year",
            "fig": fig,
            "ax": ax,
            "data": data,
            "adjusts": {}, # Adjustments for the pptx
        }


#%% ===========================================================================
### --- MAIN Procedure ---
### ===========================================================================
def main(
    cases_location: str,
    plot_type: str = 'InOut',
    prefix_old_new: list|str = None,
    highlight_runs: list|str = None,
    save_path: str = None,
    filter_unfinished: bool = True,
):
    """
    Main function to automate the plotting of ReEDS input and output data.
    Args:
        cases_location (str): 
            - Path to a folder containing multiple ReEDS runs.
        plot_type (str): 
            - Type of plot to create. Options: 'In' (input only), 'Out' (output only), 'InOut (both)'.
        prefix_old_new (list|str, optional):
            - List of tuples with old and new prefixes for ReEDS cases.
            - Example: [('old_prefix1', 'new_prefix1'), ('old_prefix2', 'new_prefix2')]
            - If a single string is provided, it should be a path to a CSV or Excel file
            - If None, the default prefixes will be used.
        highlight_runs (list|str, optional):
            - List of tuples with run_name, print_name, and color to highlight in plots.
            - Example: [('run1', 'Run 1', 'red'), ('run2', 'Run 2', 'blue')]
            - If a single string is provided, it should be a path to a CSV or Excel file.
        save_path (str, optional):
            - Path to save the PowerPoint presentation and Excel files with the figures data.
            - If None, a folder at os.path.join(reeds_path, 'runs', 'comparisons') will be created.
        filter_unfinished (bool, optional):
            If True, only finished ReEDS runs will be included in the dictionary.
            A run is considered finished if it contains a 'report.xlsx' file in the 'outputs/reeds-report' folder.
            Default is True.
    """
    # Get dictionary of cases
    cases_dict, highlight_runs = map_cases(
        cases_location,
        prefix_old_new=prefix_old_new,
        highlight_runs=highlight_runs,
        filter_unfinished=filter_unfinished,
    )

    # Initialize DataFetcher
    data_fetcher = DataFetcher(cases_dict)

    # Initialize DataPlotter
    data_plotter = DataPlotter(
        data_fetcher,
        highlight_runs,
        save_path=save_path,
    )

    # Create all plots and save them to PowerPoint
    if plot_type.lower() in ['in', 'input', 'inputs']: 
        plot_methods = data_plotter.input_plot_methods
    elif plot_type.lower() in ['out', 'output', 'outputs']:
        plot_methods = data_plotter.output_plot_methods
    elif plot_type.lower() in ['inout', 'both', 'all']:
        plot_methods = data_plotter.plot_methods

    data_plotter.create_all_plots(plot_methods)


if __name__ == '__main__' and not hasattr(sys, 'ps1'):
    def _parse_prefix_mapping(arg_list):
        if not arg_list:
            return None
        if len(arg_list) == 1 and os.path.isfile(arg_list[0]):
            return arg_list[0]
        else:
            return [tuple(item.split(':')) for item in arg_list]

    def _parse_highlight_runs(arg_list):
        if not arg_list:
            return None
        if len(arg_list) == 1 and os.path.isfile(arg_list[0]):
            return arg_list[0]
        else:
            return [tuple(item.split(':')) for item in arg_list]

    # ---- Argument parser setup ----
    parser = argparse.ArgumentParser(
        description="Generate uncertainty plots for ReEDS runs and export to PowerPoint."
    )

    parser.add_argument(
        'cases_location',
        type=str,
        help="Path(s) to one or more ReEDS case folders or parent directories."
    )

    parser.add_argument(
        '--plot_type', '-t',
        type=str,
        choices=['in', 'out', 'inout'],
        default='InOut',
        help=(
            "Type of plot to create. Options: 'In' (input only), "
            "'Out' (output only), 'InOut' (both). Default is 'InOut'."
        )
    )

    parser.add_argument(
        '--prefix_old_new', '-p',
        type=str,
        nargs='+',
        default=None,
        help=(
            "Optional list of old:new prefix mappings or a single file path (CSV or Excel). "
            "Examples:\n"
            "  -p old1:new1 old2:new2\n"
            "  -p prefix_mapping.csv"
        )
    )

    parser.add_argument(
        '--highlight_runs', '-hl',
        type=str,
        nargs='+',
        default=None,
        help=(
            "Optional list of run_name:print_name:color to highlight in plots or a single file path (CSV or Excel). "
            "Example:\n"
            "   -hl July11_MC001:MC001:red July11_MC002:MC002:blue\n"
            "   -hl highlight_runs.csv"
        )
    )

    parser.add_argument(
        '--save_path', '-s',
        type=str,
        default=None,
        help="Directory where the PowerPoint file will be saved. If omitted, a timestamped folder will be created."
    )

    # ---- Parse and unpack arguments ----
    args = parser.parse_args()
    prefix_old_new  = _parse_prefix_mapping(args.prefix_old_new)
    highlight_runs = _parse_highlight_runs(args.highlight_runs)
    save_path = args.save_path

    # If the user provided a path for prefix_old_new,
    # and the save_path is None, create a default save path
    if isinstance(prefix_old_new, str) and args.save_path is None:
        save_folder = os.path.splitext(os.path.basename(prefix_old_new))[0]
        save_path = os.path.join(reeds_path, 'runs', 'comparisons', save_folder)

    # Only filter unfinished runs when plot type is not 'in'
    if args.plot_type == 'in':
        filter_unfinished = False
    else:
        filter_unfinished = True

    print(">> Running: Uncertainty Plots for ReEDS")
    main(
        cases_location=args.cases_location,
        plot_type=args.plot_type,
        prefix_old_new=prefix_old_new,
        highlight_runs=highlight_runs,
        save_path=save_path,
        filter_unfinished=filter_unfinished,
    )

    print("\n>> Done: Plots (pptx) and data (xlsx) saved.")
